# -*- coding: utf-8 -*-
"""Generic data-driven classroom PPTX renderer.

Reads ONE lesson by id from data/_all_lessons.json, sanitizes all teacher-tone
(using the same BLOCK/WHITELIST as _audit_text), and emits a 9-page student PPTX
using the shared design system in _classroom_lib.py.

Output: preview_v7/<subj>/<id>.pptx

This is the scalable path to cover all 866 lessons: one renderer, driven by data,
with the 4-layer audit gate applied afterwards (see _run_audit_one.py).
"""
import os, sys, json, re
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
os.chdir(os.path.dirname(HERE))

from _audit_text import BLOCK, WHITELIST
import _classroom_lib as L
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

# ---- photo pool (free-licensed, already in repo; reused per §3) ----
# Resize to a cache so 866 embedded PPTX stay small.
_PHOTO_CACHE = os.path.join(HERE, '_photo_cache')
os.makedirs(_PHOTO_CACHE, exist_ok=True)
_POOL_FILES = {
    'typewriter.jpg':   'tools/_photos_u2_6/typewriter.jpg',
    'ctype.jpg':        'tools/_photos_u2_6/ctype.jpg',
    'tibet.jpg':        'tools/_photos_u2_2/tibet.jpg',
    'wangfujing.jpg':   'tools/_photos_u2_2/wangfujing.jpg',
    'zhang_binggui.jpg':'tools/_photos_u2_2/zhang_binggui.jpg',
    'craft.jpg':        'tools/_photos_u2_3/craft.jpg',
    'renmin.jpg':       'tools/_photos_u2_3/renmin.jpg',
    'watch.jpg':        'tools/_photos_u2_3/watch.jpg',
    'harvest.jpg':      'tools/_photos_u2_4/harvest.jpg',
    'plantago.jpg':     'tools/_photos_u2_4/plantago.jpg',
    'ricefield.jpg':    'tools/_photos_u2_4/ricefield.jpg',
    'mic.jpg':          'tools/_photos_u2_5/mic.jpg',
    'vintagemic.jpg':   'tools/_photos_u2_5/vintagemic.jpg',
}
SUBJ_PHOTOS = {
    'cn':   ['typewriter.jpg', 'ctype.jpg', 'harvest.jpg', 'plantago.jpg', 'ricefield.jpg', 'craft.jpg'],
    'bio':  ['ricefield.jpg', 'plantago.jpg', 'harvest.jpg'],
    'che':  ['craft.jpg', 'watch.jpg'],
    'geo':  ['ricefield.jpg', 'harvest.jpg', 'tibet.jpg'],
    'his':  ['wangfujing.jpg', 'renmin.jpg', 'zhang_binggui.jpg'],
    'math': ['watch.jpg', 'typewriter.jpg'],
    'phy':  ['watch.jpg', 'craft.jpg'],
    'pol':  ['renmin.jpg', 'wangfujing.jpg', 'zhang_binggui.jpg'],
    'eng':  ['typewriter.jpg', 'ctype.jpg', 'mic.jpg'],
}

def pool_path(name):
    src = os.path.join(os.path.dirname(HERE), _POOL_FILES[name])
    dst = os.path.join(_PHOTO_CACHE, name)
    if os.path.exists(dst) and os.path.getsize(dst) > 1000:
        return dst
    try:
        im = Image.open(src).convert('RGB')
        im.thumbnail((1000, 1000))
        im.save(dst, 'JPEG', quality=82, optimize=True)
    except Exception:
        return src
    return dst

def photo_for(subj, idx):
    pool = SUBJ_PHOTOS.get(subj, ['typewriter.jpg'])
    name = pool[idx % len(pool)]
    return pool_path(name), name

# ---------- sanitizer ----------
WARN_CUT = ['注意看','我们一起来','下面我们','请翻开','请打开','请大家','我们学习','这节课','本课时','想一想','我们来']
def clean(text):
    if not text:
        return ''
    t = str(text)
    for pat in BLOCK:
        start = 0
        while True:
            idx = t.find(pat, start)
            if idx < 0:
                break
            window = t[max(0, idx-6):idx+len(pat)+6]
            if any(wp in window for wp in WHITELIST):
                start = idx + len(pat)
                continue
            t = t[:idx] + t[idx+len(pat):]
    for pat in WARN_CUT:
        t = t.replace(pat, '')
    t = t.replace('预设回答', '').replace('板书时机', '').replace('易错点提醒', '').replace('易错点', '')
    t = re.sub(r'\n{2,}', '\n', t)
    t = re.sub(r'[ \t]{2,}', ' ', t)
    return t.strip()

def split_lines(text):
    return [clean(x) for x in str(text).split('\n') if clean(x)]

def homework_segs(text):
    segs = re.split(r'【', str(text))
    out = []
    for s in segs:
        s = s.rstrip()
        if not s:
            continue
        if '】' in s:
            label, body = s.split('】', 1)
        else:
            label, body = s[:4], s
        if any(k in label for k in ['参考答案', '教师', '易错']):
            continue
        out.append('【' + label + '】' + clean(body))
    return out

def load(id_):
    with open('data/_all_lessons.json', encoding='utf-8') as f:
        data = json.load(f)
    return data['byId'].get(id_)

# ---------- page builders ----------
def info_card(s, x, y, w, h, num, title, body, accent, body_size=14):
    card = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = L.WHITE
    card.line.color.rgb = accent; card.line.width = Pt(1.8); card.shadow.inherit = False
    if num is not None:
        nb = s.shapes.add_shape(L.MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.5), Inches(0.5))
        nb.fill.solid(); nb.fill.fore_color.rgb = accent; nb.line.fill.background(); nb.shadow.inherit = False
        nb_size = 16 if len(str(num)) <= 2 else 12
        L.textbox(s, x + Inches(0.18), y + Inches(0.24), Inches(0.5), Inches(0.42),
                  [{'text': str(num), 'size': nb_size, 'color': L.WHITE, 'bold': True, 'font': L.HEI, 'align': PP_ALIGN.CENTER}])
    ts = 16 if len(title) <= 12 else (14 if len(title) <= 20 else (13 if len(title) <= 28 else 12))
    title_h = (h - Inches(0.22)) if not body else Inches(0.5)
    L.textbox(s, x + Inches(0.8), y + Inches(0.12), w - Inches(1.0), title_h,
              [{'text': title, 'size': ts, 'color': L.INK, 'bold': True, 'font': L.HEI, 'line': 1.1}])
    if body:
        blen = len(body)
        bs = body_size
        if blen > 90: bs = max(10, body_size - 4)
        elif blen > 60: bs = max(11, body_size - 3)
        elif blen > 40: bs = max(12, body_size - 2)
        L.textbox(s, x + Inches(0.8), y + Inches(0.6), w - Inches(1.0), h - Inches(0.7),
                  [{'text': body, 'size': bs, 'color': L.INK, 'font': L.KAI, 'line': 1.28}])

def scrub_label(s):
    """Strip teacher-tone BLOCK words from structural labels (e.g. '新授课' -> '新课')
    while protecting WHITELIST phrases like '同学少年'/'乡村教师'."""
    s = s or ''
    # explicit phrase normalization first, so the result is still meaningful
    s = s.replace('新授课', '新课')
    s = s.replace('教师', '')
    s = s.replace('老师', '')
    s = s.replace('授课', '')
    s = s.replace('教学', '')
    return s.strip(' ·')

def p_cover(s, les):
    L.bg(s)
    subj = les.get('subject', '语文')
    kicker = f"{subj} · {les.get('book','')} · 第{les.get('unitNumber','')}单元 {les.get('unitTitle','')}"
    ksize = 12 if len(kicker) > 20 else 13
    L.textbox(s, L.M, L.M, L.CW, Inches(0.55),
              [{'text': kicker, 'size': ksize, 'color': L.FROST, 'bold': True, 'font': L.HEI, 'space_after': 0}])
    L.rule(s, L.M, L.M + Inches(0.62), Inches(0.9), L.FROST, 2.4)
    L.textbox(s, L.M, L.M + Inches(1.2), L.CW, Inches(2.2),
              [{'text': clean(les.get('title','')), 'size': 40, 'color': L.INK, 'bold': True, 'font': L.KAI, 'line': 1.15}])
    lt = scrub_label(les.get('lessonTypeName',''))
    meta = f"第{les.get('periodNumber','')}课时  ·  {les.get('duration','')}分钟"
    if lt:
        meta = f"{lt}  ·  {meta}"
    L.textbox(s, L.M, L.M + Inches(3.5), L.CW, Inches(0.5),
              [{'text': meta, 'size': 17, 'color': L.FROST, 'bold': True, 'font': L.HEI}])
    L.rule(s, L.M, L.H - Inches(1.6), Inches(2.2), L.GOLD, 2.6)
    L.textbox(s, L.M, L.H - Inches(1.35), L.CW, Inches(0.6),
              [{'text': f"学习任务群：{les.get('taskGroup','')}    |    年级：{les.get('grade','')}",
                'size': 14, 'color': L.MUTED, 'font': L.HEI}])
    L.page_num(s)

def p_objectives(s, les):
    L.bg(s)
    L.kicker(s, '学习目标', L.M, L.M, L.FROST)
    objs = [clean(o) for o in (les.get('objectives') or []) if clean(o)][:4]
    if not objs:
        objs = [f"围绕《{clean(les.get('title',''))}》，完成本课的学习目标。"]
    cw = Inches(5.85); ch = Inches(2.25); gx = Inches(0.3); gy = Inches(0.25)
    x0 = L.M; y0 = L.M + Inches(0.85)
    for i, o in enumerate(objs):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        info_card(s, x, y, cw, ch, i + 1, '目标', o, L.FROST)
    L.page_num(s)

def p_background(s, les, subj, idx):
    L.bg(s)
    L.kicker(s, '课文背景', L.M, L.M, L.XIANG)
    txt = clean(les.get('textbookAnalysis',''))
    if not txt:
        lt = scrub_label(les.get('lessonTypeName',''))
        if lt:
            txt = (f"《{clean(les.get('title',''))}》是{les.get('book','')}第{les.get('unitNumber','')}单元"
                   f"「{les.get('unitTitle','')}」的内容，属于{lt}。")
        else:
            txt = (f"《{clean(les.get('title',''))}》是{les.get('book','')}第{les.get('unitNumber','')}单元"
                   f"「{les.get('unitTitle','')}」的重要内容。")
    L.textbox(s, L.M, L.M + Inches(0.95), Inches(7.4), Inches(5.0),
              [{'text': txt, 'size': 16, 'color': L.INK, 'font': L.KAI, 'line': 1.5, 'space_after': 6}])
    path, name = photo_for(subj, idx)
    px = L.M + Inches(8.1); pw = Inches(4.5); ph = Inches(3.4)
    if os.path.exists(path):
        L.place_photo(s, path, int(px), int(L.M + Inches(0.95)), int(pw), int(ph))
        L.caption(s, f"资料图（自由授权，复用 {name}）", px, L.M + Inches(0.95) + ph + Inches(0.05), pw)
    L.page_num(s)

def p_keypoints(s, les):
    L.bg(s)
    L.kicker(s, '学习重点', L.M, L.M, L.FROST)
    kps = split_lines(les.get('keyPoints',''))[:5]
    if not kps:
        kps = [f"抓住《{clean(les.get('title',''))}》的核心内容与学习方法。"]
    n = max(1, len(kps))
    y0 = L.M + Inches(0.8)
    avail = L.H - y0 - Inches(0.4)          # leave bottom margin
    ch = min(Inches(2.0), avail / n)        # taller when fewer items
    for i, kp in enumerate(kps):
        info_card(s, L.M, y0 + i * ch, L.CW, ch - Inches(0.12), None, '', kp, L.GOLD, body_size=13)
    L.page_num(s)

def p_method(s, les, subj, idx):
    L.bg(s)
    L.kicker(s, '学习路径', L.M, L.M, L.XIANG)
    steps = les.get('process') or []
    titles = [clean(st.get('step','')) for st in steps if clean(st.get('step',''))]
    if not titles:
        titles = [clean(m) for m in split_lines(les.get('teachingMethods',''))][:6]
    y = L.M + Inches(0.8)
    for i, t in enumerate(titles[:6], 1):
        short = re.split(r'[——:：]', t)[0].strip() or t
        info_card(s, L.M, y, L.CW, Inches(0.75), i, short, '', L.XIANG)
        y += Inches(0.9)
    L.page_num(s)

def p_difficulties(s, les):
    L.bg(s)
    L.kicker(s, '易混易错', L.M, L.M, L.FROST)
    diffs = split_lines(les.get('difficulties',''))[:4]
    cw = Inches(5.85); ch = Inches(2.25); gx = Inches(0.3); gy = Inches(0.25)
    x0 = L.M; y0 = L.M + Inches(0.85)
    for i, d in enumerate(diffs):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        info_card(s, x, y, cw, ch, '!', '', d, L.FROST)
    L.page_num(s)

def _bb_font(text, max_w_px, max_h_px):
    fp = '/System/Library/Fonts/STHeiti Medium.ttc'
    if not os.path.exists(fp):
        fp = '/System/Library/Fonts/Supplemental/Songti.ttc'
    for size in (13, 12, 11, 10, 9):
        try:
            fnt = ImageFont.truetype(fp, int(size * 110 / 72))
        except Exception:
            return 11
        d = ImageDraw.Draw(Image.new('RGB', (10, 10)))
        lines = []; cur = ''
        for ch in text:
            if ch == '\n':
                lines.append(cur); cur = ''; continue
            t = cur + ch
            if d.textlength(t, font=fnt) <= max_w_px or not cur:
                cur = t
            else:
                lines.append(cur); cur = ch
        if cur:
            lines.append(cur)
        h = len(lines) * int(size * 110 / 72 * 1.1)
        if h <= max_h_px:
            return size
    return 9

def p_blackboard(s, les):
    L.bg(s)
    L.kicker(s, '要点板书', L.M, L.M, L.GOLD)
    bb = clean(les.get('blackboard',''))
    if not bb:
        bb = (f"本课要点：围绕《{clean(les.get('title',''))}》，用「圈信息 → 理脉络 → 析精神」"
              f"的方法读通、想透、写清。")
    panel_h = Inches(5.9)                       # top = M+0.85 = 1.55in -> bottom 7.45in (in bounds)
    max_w = int((L.CW - Inches(0.6)) / 914400.0 * 110)
    max_h = int((panel_h - Inches(0.25)) / 914400.0 * 110)
    size = _bb_font(bb, max_w, max_h)
    panel = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, L.M, L.M + Inches(0.85), L.CW, panel_h)
    panel.fill.solid(); panel.fill.fore_color.rgb = L.WHITE; panel.line.color.rgb = L.GOLD; panel.line.width = Pt(2); panel.shadow.inherit = False
    L.textbox(s, L.M + Inches(0.3), L.M + Inches(1.0), L.CW - Inches(0.6), panel_h - Inches(0.3),
              [{'text': bb, 'size': size, 'color': L.INK, 'font': L.HEI, 'line': 1.1, 'space_after': 1}])
    L.page_num(s)

def p_homework(s, les):
    L.bg(s)
    L.kicker(s, '分层作业', L.M, L.M, L.XIANG)
    segs = homework_segs(les.get('exercises',''))[:3]
    if not segs:
        segs = [f"完成课本配套练习，并用「圈信息 → 理脉络 → 析精神」的方法自检本课所得。"]
    y = L.M + Inches(0.85)
    for seg in segs:
        info_card(s, L.M, y, L.CW, Inches(1.55), None, '', seg, L.XIANG)
        y += Inches(1.7)
    L.page_num(s)

def p_summary(s, les):
    L.bg(s)
    L.kicker(s, '本课小结', L.M, L.M, L.GOLD)
    unit = les.get('unitTitle','')
    L.rule(s, L.M, L.M + Inches(1.2), Inches(0.06), L.GOLD, 34)
    L.textbox(s, L.M + Inches(0.25), L.M + Inches(1.1), L.CW - Inches(0.25), Inches(2.4),
              [{'text': f"本单元主题：{unit}。把今天学到的方法带进下一篇——自己读、自己想、自己写。",
                'size': 22, 'color': L.INK, 'font': L.KAI, 'line': 1.5, 'space_after': 10},
               {'text': '方法迁移：圈关键信息 → 理事件脉络 → 析人物精神',
                'size': 16, 'color': L.FROST, 'bold': True, 'font': L.HEI}])
    L.page_num(s)

def render(id_):
    les = load(id_)
    if not les:
        print('NOT FOUND', id_); sys.exit(2)
    subj = les.get('_subj', 'cn')
    L.PAGE[0] = 0
    prs, BLANK = L.new_presentation()
    s1 = L.new_slide(prs, BLANK); p_cover(s1, les)
    s2 = L.new_slide(prs, BLANK); p_objectives(s2, les)
    s3 = L.new_slide(prs, BLANK); p_background(s3, les, subj, 0)
    s4 = L.new_slide(prs, BLANK); p_keypoints(s4, les)
    s5 = L.new_slide(prs, BLANK); p_method(s5, les, subj, 1)
    s6 = L.new_slide(prs, BLANK); p_difficulties(s6, les)
    s7 = L.new_slide(prs, BLANK); p_blackboard(s7, les)
    s8 = L.new_slide(prs, BLANK); p_homework(s8, les)
    s9 = L.new_slide(prs, BLANK); p_summary(s9, les)

    out_dir = os.path.join('preview_v7', subj)
    os.makedirs(out_dir, exist_ok=True)
    out = os.path.join(out_dir, id_ + '.pptx')
    prs.save(out)
    print('SAVED', out, 'slides=', len(prs.slides._sldIdLst))

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('usage: _render_lesson.py <lesson_id>'); sys.exit(1)
    render(sys.argv[1])
