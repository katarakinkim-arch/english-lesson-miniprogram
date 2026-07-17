#!/usr/bin/env python3
"""
沁园春·长沙 — 杂志风格课堂PPT (python-pptx)
设计语言：商业杂志排版 / 网格系统 / 色块+照片 / 三色系统 / 环环相扣
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── 全局常量 ────────────────────────────────────────
W = Inches(13.333)   # 16:9 widescreen
H = Inches(7.5)
MX = Inches(0.9)     # margins
MY = Inches(0.7)

# 三色系统
FROST  = RGBColor(178, 42,  58)   # 霜红 — 万山红遍
XIANG  = RGBColor(46,  125, 107)  # 湘碧 — 漫江碧透
INK    = RGBColor(28,  43,  51)   # 墨蓝 — 鹰击长空
WARM   = RGBColor(110, 80,  55)   # 陶土 — 百舸争流
GOLD   = RGBColor(190, 150, 80)   # 金秋
WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(248, 245, 240)  # 暖纸白
GRAY   = RGBColor(140, 140, 140)
DARK   = RGBColor(40,  40,  40)

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'classroom_images')
OUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'cn_qyc_magazine')
OUT_PDF = os.path.join(os.path.expanduser('~'), 'Desktop', '沁园春长沙-课堂版-杂志风.pptx')

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── 辅助函数 ────────────────────────────────────────
def blank_slide():
    """创建白底空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                color=INK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=Pt(16),
                  color=INK, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.5,
                  font_name='Microsoft YaHei'):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size.pt * (line_spacing - 1) * 0.5)
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=INK, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, height, color=FROST, weight=Pt(3)):
    """细线分隔"""
    shape = add_rect(slide, left, top, width, height, fill_color=color)
    return shape

def add_image_bg(slide, img_name, opacity=0.12):
    """添加全幅背景图片，白色半透明遮罩叠在上面保持文字可读"""
    img_path = os.path.join(IMG_DIR, img_name)
    if not os.path.exists(img_path):
        return None
    pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), W, H)
    # 白色半透明矩形覆盖
    alpha_hex = format(int((1 - opacity) * 255), '02x')
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = WHITE
    # 设置透明度 via XML
    solidFill = overlay.fill._fill
    srgb = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity * 100000))})
        srgb.append(alpha)
    overlay.line.fill.background()
    return pic

def add_circle(slide, left, top, size, fill_color=FROST):
    """添加圆形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def section_divider(slide, number, title, subtitle, accent_color=FROST):
    """章节分隔页 — 大数字+标题+副标题"""
    # 左侧色块
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, fill_color=accent_color)
    # 大数字
    add_textbox(slide, MX, Inches(1.5), Inches(4), Inches(2),
                str(number).zfill(2), font_size=Pt(120), color=accent_color,
                bold=True, alignment=PP_ALIGN.LEFT)
    # 标题
    add_textbox(slide, MX, Inches(3.5), Inches(10), Inches(1.2),
                title, font_size=Pt(48), color=INK, bold=True)
    # 副标题
    add_textbox(slide, MX, Inches(4.8), Inches(10), Inches(0.8),
                subtitle, font_size=Pt(20), color=GRAY)
    # 底部分隔线
    add_line(slide, MX, Inches(6.0), Inches(3), Pt(3), color=accent_color)
    # 页码
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
                str(number), font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ── 幻灯片 ──────────────────────────────────────────

# ===== 1. 封面 =====
s = blank_slide()
# 顶部色条
add_rect(s, Inches(0), Inches(0), W, Inches(0.08), fill_color=FROST)
# 标签
add_textbox(s, Inches(0.9), Inches(0.4), Inches(3), Inches(0.5),
            '高中语文 · 必修上册', font_size=Pt(14), color=GRAY)
# 主标题
add_textbox(s, Inches(0.9), Inches(2.0), Inches(11), Inches(2.0),
            '沁园春 · 长沙', font_size=Pt(84), color=INK, bold=True)
# 作者
add_textbox(s, Inches(0.9), Inches(3.8), Inches(5), Inches(0.8),
            '毛泽东', font_size=Pt(36), color=FROST, bold=True)
# 分隔线
add_line(s, Inches(0.9), Inches(4.7), Inches(3), Pt(3), color=FROST)
# 副标题
add_textbox(s, Inches(0.9), Inches(5.0), Inches(8), Inches(0.6),
            '独立寒秋  湘江北去  橘子洲头', font_size=Pt(22), color=INK)
# 课时信息
add_textbox(s, Inches(0.9), Inches(5.8), Inches(5), Inches(0.5),
            '第一课时：品读意象 · 第二课时：知人论世', font_size=Pt(16), color=GRAY)
# 页码
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '01', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 2. 本课导览 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '课堂导览', font_size=Pt(36), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

sections = [
    ('01', '知人论世', '诗人·时代·背景', FROST),
    ('02', '初读感知', '诵读·节拍·韵律', XIANG),
    ('03', '品读意象', '万山·层林·漫江·鹰鱼', FROST),
    ('04', '炼字品析', '击·翔·竞·遍·透', WARM),
    ('05', '下阕探微', '同学少年·中流击水', INK),
    ('06', '对比升华', '悲秋与颂秋·谁主沉浮', FROST),
]
for i, (num, title, desc, clr) in enumerate(sections):
    row = i // 3
    col = i % 3
    x = MX + col * Inches(3.8)
    y = Inches(1.8) + row * Inches(2.4)
    # 编号圆
    add_circle(s, x, y, Inches(0.55), fill_color=clr)
    add_textbox(s, x + Inches(0.1), y + Inches(0.05), Inches(0.4), Inches(0.4),
                num, font_size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_textbox(s, x + Inches(0.7), y, Inches(2.8), Inches(0.5),
                title, font_size=Pt(24), color=INK, bold=True)
    add_textbox(s, x + Inches(0.7), y + Inches(0.55), Inches(2.8), Inches(0.4),
                desc, font_size=Pt(14), color=GRAY)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '02', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 3. PART 1 章节页 =====
s = blank_slide()
section_divider(s, 1, '知人论世', '诗人 · 时代 · 背景')


# ===== 4. 作者小传 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '作者小传', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 左侧生平
add_textbox(s, MX, Inches(1.4), Inches(1.5), Inches(0.5),
            '生平', font_size=Pt(20), color=FROST, bold=True)
author_lines = [
    '毛泽东（1893—1976）',
    '字润之，湖南湘潭人',
    '伟大的无产阶级革命家、',
    '战略家、理论家',
    '中国共产党、中国人民解放军、',
    '中华人民共和国的主要缔造者',
    '',
    '诗人、书法家',
    '诗词气势磅礴、意境深远',
]
add_multiline(s, MX, Inches(2.0), Inches(4.5), Inches(4.2),
              author_lines, font_size=Pt(16), color=INK, line_spacing=1.3)

# 右侧代表作品
add_textbox(s, Inches(7.5), Inches(1.4), Inches(2), Inches(0.5),
            '代表诗作', font_size=Pt(20), color=XIANG, bold=True)
works_lines = [
    '《沁园春·雪》',
    '《七律·长征》',
    '《采桑子·重阳》',
    '《水调歌头·游泳》',
    '《卜算子·咏梅》',
    '《浪淘沙·北戴河》',
]
add_multiline(s, Inches(7.5), Inches(2.0), Inches(4.5), Inches(3.0),
              works_lines, font_size=Pt(16), color=INK, line_spacing=1.5)

# 青年照片占位框
add_rect(s, Inches(7.5), Inches(5.0), Inches(2.2), Inches(1.6),
         fill_color=RGBColor(240, 235, 230), line_color=FROST)
add_textbox(s, Inches(7.5), Inches(5.5), Inches(2.2), Inches(0.5),
            '青年毛泽东', font_size=Pt(12), color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '04', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 5. 写作背景 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '写作背景', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 时间轴
events = [
    ('1925 秋', '毛泽东离开韶山，前往广州', '途径长沙', FROST),
    ('重游', '橘子洲头', '独立寒秋，湘江北去', XIANG),
    ('时局', '革命形势高涨', '五卅运动后工农运动蓬勃发展', WARM),
    ('心境', '对国家命运', '的深沉思考与热切期望', INK),
]
for i, (title, main, sub, clr) in enumerate(events):
    y = Inches(1.8) + i * Inches(1.25)
    # 时间标签
    add_circle(s, MX, y, Inches(0.45), fill_color=clr)
    add_textbox(s, MX + Inches(0.05), y + Inches(0.05), Inches(0.35), Inches(0.35),
                str(i+1), font_size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 内容
    add_textbox(s, MX + Inches(0.7), y - Inches(0.05), Inches(2), Inches(0.5),
                title, font_size=Pt(20), color=clr, bold=True)
    add_textbox(s, Inches(3.8), y - Inches(0.05), Inches(4), Inches(0.5),
                main, font_size=Pt(18), color=INK, bold=True)
    add_textbox(s, Inches(3.8), y + Inches(0.45), Inches(5), Inches(0.45),
                sub, font_size=Pt(14), color=GRAY)

# 底部引用
add_rect(s, MX, Inches(6.5), Inches(10.5), Inches(0.5), fill_color=RGBColor(248, 245, 240))
add_textbox(s, Inches(1.2), Inches(6.52), Inches(10), Inches(0.45),
            '"问苍茫大地，谁主沉浮？" —— 这是一个时代的叩问',
            font_size=Pt(15), color=FROST, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '05', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 6. PART 2 章节页 =====
s = blank_slide()
section_divider(s, 2, '初读感知', '诵读 · 节拍 · 韵律', accent_color=XIANG)


# ===== 7. 全词展示（上阕大字） =====
s = blank_slide()
add_textbox(s, MX, Inches(0.2), Inches(3), Inches(0.6),
            '上 阕', font_size=Pt(28), color=FROST, bold=True)
add_line(s, MX, Inches(0.85), Inches(1.5), Pt(3), color=FROST)

lines_shang = [
    '独立寒秋，湘江北去，橘子洲头。',
    '看万山红遍，层林尽染；',
    '漫江碧透，百舸争流。',
    '鹰击长空，鱼翔浅底，',
    '万类霜天竞自由。',
    '怅寥廓，问苍茫大地，谁主沉浮？',
]
# 左侧原文
add_multiline(s, MX, Inches(1.2), Inches(7), Inches(4.0),
              lines_shang, font_size=Pt(28), color=INK, line_spacing=2.2,
              bold=False)

# 右侧分析提示
add_textbox(s, Inches(8.5), Inches(1.0), Inches(4), Inches(0.5),
            '写 景', font_size=Pt(24), color=INK, bold=True)
add_line(s, Inches(8.5), Inches(1.55), Inches(1.2), Pt(2), color=XIANG)
hints = [
    '远眺 → 近观',
    '静景 → 动景',
    '万物生机勃发',
    '',
    '末句陡转：',
    '由景及情',
]
add_multiline(s, Inches(8.5), Inches(1.8), Inches(3.5), Inches(3.0),
              hints, font_size=Pt(16), color=GRAY, line_spacing=1.8)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '07', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 8. 朗读指导 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '朗读指导', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=XIANG)

# 三栏：字音 / 节拍 / 情感
sections_read = [
    ('字音', '沁(qìn) 寥廓(liáo kuò)\n舸(gě) 怅(chàng)\n峥嵘(zhēng róng)\n遒(qiú) 遏(è)', XIANG),
    ('节拍', '独立/寒秋，湘江/北去，\n橘子洲/头。\n看/万山红遍，层林/尽染；\n漫江/碧透，百舸/争流。', FROST),
    ('情感曲线', '平缓起 → 渐高昂\n→ 至"竞自由"顶点\n→ "怅寥廓"微沉\n→ "谁主沉浮"浑厚收', WARM),
]
for i, (title, content, clr) in enumerate(sections_read):
    x = MX + i * Inches(3.9)
    add_textbox(s, x, Inches(1.5), Inches(3.5), Inches(0.5),
                title, font_size=Pt(22), color=clr, bold=True)
    add_line(s, x, Inches(2.1), Inches(1.5), Pt(2), color=clr)
    add_multiline(s, x, Inches(2.3), Inches(3.5), Inches(3.5),
                  content.split('\n'), font_size=Pt(16), color=INK, line_spacing=1.8)

# 底部提醒
add_rect(s, MX, Inches(6.2), Inches(10.5), Inches(0.6), fill_color=RGBColor(248, 245, 240))
add_textbox(s, Inches(1.2), Inches(6.25), Inches(10), Inches(0.5),
            '朗读要求：读准字音 → 读出节拍 → 读出气势 → 读出情感',
            font_size=Pt(16), color=FROST, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '08', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 9. PART 3 章节页 =====
s = blank_slide()
section_divider(s, 3, '品读意象', '万山 · 层林 · 漫江 · 百舸 · 鹰 · 鱼', accent_color=FROST)


# ===== 10. 四幅图画概览 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '全词结构 · 四幅图画', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

scenes = [
    ('第一幅', '独立寒秋图', '湘江北去\n橘子洲头', FROST, '远景·静景'),
    ('第二幅', '湘江秋景图', '万山红遍 层林尽染\n漫江碧透 百舸争流', XIANG, '全景·色彩'),
    ('第三幅', '万物生机图', '鹰击长空 鱼翔浅底\n万类霜天竞自由', WARM, '特写·动态'),
    ('第四幅', '中流击水图', '到中流击水\n浪遏飞舟', INK, '豪情·壮志'),
]
for i, (num, title, desc, clr, tag) in enumerate(scenes):
    x = MX + i * Inches(3.0)
    # 卡片背景
    add_rect(s, x, Inches(1.5), Inches(2.7), Inches(4.8),
             fill_color=RGBColor(248, 245, 240))
    # 顶部色条
    add_rect(s, x, Inches(1.5), Inches(2.7), Inches(0.06), fill_color=clr)
    # 编号
    add_textbox(s, x + Inches(0.2), Inches(1.7), Inches(2.3), Inches(0.4),
                num, font_size=Pt(16), color=clr, bold=True)
    # 标题
    add_textbox(s, x + Inches(0.2), Inches(2.1), Inches(2.3), Inches(0.8),
                title, font_size=Pt(22), color=INK, bold=True)
    # 内容
    add_multiline(s, x + Inches(0.2), Inches(3.0), Inches(2.3), Inches(1.5),
                  desc.split('\n'), font_size=Pt(14), color=INK, line_spacing=1.6)
    # 标签
    add_textbox(s, x + Inches(0.2), Inches(4.5), Inches(2.3), Inches(0.4),
                tag, font_size=Pt(12), color=GRAY)
    # 箭头（除最后一幅外）
    if i < 3:
        add_textbox(s, x + Inches(2.8), Inches(3.2), Inches(0.3), Inches(0.3),
                    '→', font_size=Pt(24), color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '10', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 11-16. 六大意象可视化页 =====
def imagery_slide(slide_num, title, photo_name, quote, analysis_lines, color, extra_notes=None):
    """意象页：左照片 + 右分析 + 底部金句"""
    s = blank_slide()

    # 左侧照片
    img_path = os.path.join(IMG_DIR, photo_name)
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, MX, Inches(1.0), Inches(6.0), Inches(5.0))

    # 右侧分析区
    add_textbox(s, Inches(7.5), Inches(1.0), Inches(5), Inches(0.6),
                title, font_size=Pt(30), color=color, bold=True)
    # 诗句大字
    add_textbox(s, Inches(7.5), Inches(1.7), Inches(5), Inches(0.8),
                quote, font_size=Pt(36), color=INK, bold=True)
    # 分隔线
    add_line(s, Inches(7.5), Inches(2.6), Inches(2.5), Pt(2), color=color)
    # 分析
    add_multiline(s, Inches(7.5), Inches(2.9), Inches(4.8), Inches(3.0),
                  analysis_lines, font_size=Pt(15), color=INK, line_spacing=1.5)

    if extra_notes:
        add_rect(s, Inches(7.5), Inches(6.0), Inches(4.5), Inches(0.7),
                 fill_color=RGBColor(252, 248, 242))
        add_textbox(s, Inches(7.7), Inches(6.05), Inches(4.3), Inches(0.6),
                    extra_notes, font_size=Pt(12), color=GRAY)

    add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
                str(slide_num).zfill(2), font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)
    return s

# 11. 独立寒秋
imagery_slide(11, '独立寒秋', '01_river_island.jpg',
              '独立寒秋，湘江北去，橘子洲头。',
              ['独：孤身一人，遗世独立', '寒秋：时令渲染，肃杀之气', '北去：湘江北上，暗示心向北方',
               '', '意境：', '一个"独"字，写出诗人的卓然独立', '"寒秋"二字定下全词苍茫基调',
               '湘江北去——既是实景，更是心志'],
              FROST, '讨论：诗人为何要"独立"？写"寒秋"有何深意？')

# 12. 万山红遍
imagery_slide(12, '万山红遍', '02_autumn_mountains.jpg',
              '看万山红遍，层林尽染；',
              ['万山：极言山之多，视野之广', '红遍：红色铺天盖地，视觉冲击', '层林尽染：层叠的树林如同被颜料浸透',
               '', '色彩的力量：', '"红"是革命的颜色，是生机的颜色', '不是萧瑟的枯黄，而是热烈的绯红',
               '"遍"字——无处不红，无处不热烈'],
              FROST, '"红遍"与"尽染"——不是点缀的红，是铺天盖地的红')

# 13. 层林尽染
imagery_slide(13, '层林尽染', '03_autumn_forest.jpg',
              '层林尽染',
              ['层：一层又一层，立体层次', '染：如同画笔渲染，极具画面感',
               '', '与古人对比：', '古人秋词：枯藤老树 → 凋零',
               '此词秋景：层林尽染 → 绚烂', '',
               '同是秋天，感受完全不同——',
               '因为人的心境不同'],
              XIANG, '同一"染"字——是染料，更是心染')

# 14. 漫江碧透
imagery_slide(14, '漫江碧透', '04_green_river.jpg',
              '漫江碧透，百舸争流。',
              ['漫江：江水浩渺，一望无际', '碧透：碧绿到了极致，清澈见底', '百舸：舸，大船；百，极言其多',
               '争流：争先恐后，破浪前行', '',
               '色彩对比：', '红遍（暖色·热烈）vs 碧透（冷色·清澈）',
               '一动一静：', '碧透（静·江水）百舸争流（动·船只）'],
              XIANG, '"碧透"写水的极致清澈——非亲身所感不能道')

# 15. 百舸争流
imagery_slide(15, '百舸争流', '05_boats.jpg',
              '百舸争流',
              ['百舸：百艘大船，场面壮阔', '争流：争相前行，千帆竞发', '',
               '象征意味：', '"百舸"——时代的弄潮儿',
               '"争流"——奋勇争先的姿态', '',
               '这不仅是湘江上的船', '更是革命洪流中的战士'],
              WARM, '"争"字——不是被动漂流，是主动竞逐')

# 16. 鹰击长空
imagery_slide(16, '鹰击长空', '02_autumn_mountains.jpg',
              '鹰击长空，鱼翔浅底，',
              ['鹰击长空：雄鹰搏击于辽阔的天空',
               '击：搏击，充满力量的动作',
               '不是"飞"——"飞"是状态，',
               '"击"是力量，是主动出击', '',
               '鱼翔浅底：鱼儿在清浅的水中翱翔',
               '翔：本用于鸟类，移用于鱼',
               '写出鱼的轻盈灵动，如在天空飞翔',
               '浅底：水清极深也，透视效果'],
              INK, '"击"与"翔"——两个字盘活了上阕的动态之美')

# 17. 万类霜天
imagery_slide(17, '万类霜天', '00_cover.jpg',
              '万类霜天竞自由。',
              ['万类：天地间的一切生命', '霜天：深秋时节，寒气降临', '竞自由：万物都在争取自由地生长',
               '', '总结上阕写景：', '空间：远→近→高→低',
               '色彩：红（暖）→碧（冷）', '动态：静→动→激烈竞争', '',
               '一切自然景物背后——', '是生命力的磅礴迸发'],
              FROST, '"竞"——一个动作写尽万物精神')

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '17', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 18. 意象分析总表 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '意象分析总表', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 表格数据
table_data = [
    ['意象', '诗句', '特点', '手法', '效果'],
    ['山', '万山红遍', '红遍·热烈', '远眺·静景', '生机勃勃'],
    ['林', '层林尽染', '尽染·绚烂', '静景·色彩', '视觉冲击'],
    ['江', '漫江碧透', '碧透·清澈', '近观·静景', '明净开阔'],
    ['舸', '百舸争流', '争流·竞发', '动景·场面', '奋发向上'],
    ['鹰', '鹰击长空', '击·搏击', '仰视·动态', '力量之美'],
    ['鱼', '鱼翔浅底', '翔·轻盈', '俯视·动态', '自由之美'],
]
tbl_left = MX
tbl_top = Inches(1.5)
cols = len(table_data[0])
rows = len(table_data)
tbl_w = Inches(11.2)
tbl_h = Inches(4.8)
col_w = tbl_w // cols

table_shape = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
table = table_shape.table

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(15)
            paragraph.font.name = 'Microsoft YaHei'
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                rPr = run._r.get_or_add_rPr()
                rPr.set(qn('a:eaTypeface'), 'Microsoft YaHei')
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = INK
        # 表头颜色
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 245, 240)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '18', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 19. 远近动静结合 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '写景手法 · 远近动静', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

perspectives = [
    ('远眺', '万山红遍\n层林尽染', '静', FROST, Inches(1.5)),
    ('近观', '漫江碧透\n百舸争流', '动', XIANG, Inches(4.6)),
    ('仰视', '鹰击长空', '动', WARM, Inches(7.7)),
    ('俯视', '鱼翔浅底', '动', INK, Inches(10.8)),
]
for label, content, d_type, clr, x in perspectives:
    add_rect(s, x, Inches(1.5), Inches(1.8), Inches(3.5),
             fill_color=RGBColor(248, 245, 240))
    add_rect(s, x, Inches(1.5), Inches(1.8), Inches(0.06), fill_color=clr)
    add_textbox(s, x + Inches(0.1), Inches(1.7), Inches(1.6), Inches(0.4),
                f'{label}·{d_type}', font_size=Pt(16), color=clr, bold=True)
    add_multiline(s, x + Inches(0.1), Inches(2.2), Inches(1.6), Inches(1.5),
                  content.split('\n'), font_size=Pt(20), color=INK, line_spacing=1.6,
                  bold=True, alignment=PP_ALIGN.CENTER)

# 底部总结
add_rect(s, MX, Inches(5.8), Inches(11.2), Inches(0.7), fill_color=INK)
add_textbox(s, Inches(1.2), Inches(5.85), Inches(10), Inches(0.55),
            '写景层次：远→近→高→低    |    动静结合：静（铺垫）→ 动（高潮）    |    万物竞自由',
            font_size=Pt(15), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '19', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 20. PART 4 章节页 =====
s = blank_slide()
section_divider(s, 4, '炼字品析', '击 · 翔 · 竞 · 遍 · 透', accent_color=WARM)


# ===== 21. 炼字："击" vs "飞" =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '炼字品析  ① "击" vs "飞"', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=WARM)

# 左栏：击
add_rect(s, MX, Inches(1.5), Inches(5.2), Inches(5.2), fill_color=RGBColor(252, 248, 242))
add_textbox(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.5),
            '鹰 击 长 空', font_size=Pt(42), color=FROST, bold=True)
add_multiline(s, Inches(1.2), Inches(2.5), Inches(4.5), Inches(3.5),
              ['"击" = 搏击、用力拍打', '', '✓ 写出鹰的矫健有力',
               '✓ 刻画主动出击的姿态', '✓ 赋予鹰人格化的力量感',
               '', '鹰不是"飞"，是"击"——',
               '搏击风云，掌控天空'],
              font_size=Pt(16), color=INK, line_spacing=1.4)

# 右栏：飞
add_rect(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(5.2), fill_color=RGBColor(250, 250, 250))
add_textbox(s, Inches(7.5), Inches(1.7), Inches(4.5), Inches(0.5),
            '鹰 飞 长 空', font_size=Pt(42), color=GRAY, bold=True)
add_multiline(s, Inches(7.5), Inches(2.5), Inches(4.5), Inches(3.5),
              ['"飞" = 飞翔、飞行', '', '平铺直叙，缺乏力度',
               '仅描述动作，无个性', '无法传达革命激情',
               '', '"飞"字平庸无力，',
               '被"击"字全面碾压'],
              font_size=Pt(16), color=GRAY, line_spacing=1.4)

# VS
add_circle(s, Inches(6.1), Inches(3.3), Inches(0.8), fill_color=FROST)
add_textbox(s, Inches(6.15), Inches(3.45), Inches(0.7), Inches(0.5),
            'VS', font_size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '21', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 22. 炼字："翔" vs "游" =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '炼字品析  ② "翔" vs "游"', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=XIANG)

# 左栏：翔
add_rect(s, MX, Inches(1.5), Inches(5.2), Inches(5.2), fill_color=RGBColor(248, 252, 250))
add_textbox(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.5),
            '鱼 翔 浅 底', font_size=Pt(42), color=XIANG, bold=True)
add_multiline(s, Inches(1.2), Inches(2.5), Inches(4.5), Inches(3.5),
              ['"翔" = 翱翔（本用于鸟）', '', '✓ 移用于鱼，修辞上为"移就"',
               '✓ 写出鱼的轻盈灵动', '✓ 如在空中飞翔般自由',
               '', '鸟翔于天，鱼翔于水——',
               '天地一体，万物自由'],
              font_size=Pt(16), color=INK, line_spacing=1.4)

# 右栏：游
add_rect(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(5.2), fill_color=RGBColor(250, 250, 250))
add_textbox(s, Inches(7.5), Inches(1.7), Inches(4.5), Inches(0.5),
            '鱼 游 浅 底', font_size=Pt(42), color=GRAY, bold=True)
add_multiline(s, Inches(7.5), Inches(2.5), Inches(4.5), Inches(3.5),
              ['"游" = 游动', '', '普通动词，无特色',
               '无法表现"浅底"的清澈', '缺少诗意和想象力',
               '', '"游"字过于写实，',
               '远不及"翔"的浪漫与自由'],
              font_size=Pt(16), color=GRAY, line_spacing=1.4)

# VS
add_circle(s, Inches(6.1), Inches(3.3), Inches(0.8), fill_color=XIANG)
add_textbox(s, Inches(6.15), Inches(3.45), Inches(0.7), Inches(0.5),
            'VS', font_size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '22', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 23. 炼字："竞" =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '炼字品析  ③ "竞"', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 中心大字
add_textbox(s, Inches(3), Inches(1.5), Inches(7), Inches(1.2),
            '万类霜天竞自由', font_size=Pt(56), color=FROST, bold=True, alignment=PP_ALIGN.CENTER)

# 三栏分析
analysis_3 = [
    ('字义', ['竞争、竞逐', '万物都在争夺', '生存和发展的权利'], FROST),
    ('效果', ['赋予万物以人的意志', '将自然生命拟人化', '写出生命的蓬勃律动'], XIANG),
    ('升华', ['从写景上升为哲理', '一切生命都在竞逐自由', '引出下句的沉思'], WARM),
]
for i, (title, lines, clr) in enumerate(analysis_3):
    x = MX + i * Inches(3.9)
    add_rect(s, x, Inches(3.0), Inches(3.5), Inches(3.2), fill_color=RGBColor(248, 245, 240))
    add_rect(s, x, Inches(3.0), Inches(3.5), Inches(0.06), fill_color=clr)
    add_textbox(s, x + Inches(0.2), Inches(3.2), Inches(3.1), Inches(0.5),
                title, font_size=Pt(20), color=clr, bold=True)
    add_multiline(s, x + Inches(0.2), Inches(3.8), Inches(3.1), Inches(2.0),
                  lines, font_size=Pt(15), color=INK, line_spacing=1.6)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '23', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 24. PART 5 章节页 =====
s = blank_slide()
section_divider(s, 5, '下阕探微', '同学少年 · 中流击水 · 谁主沉浮', accent_color=INK)


# ===== 25. 下阕原文 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.2), Inches(3), Inches(0.6),
            '下 阕', font_size=Pt(28), color=INK, bold=True)
add_line(s, MX, Inches(0.85), Inches(1.5), Pt(3), color=INK)

lines_xia = [
    '携来百侣曾游。',
    '忆往昔峥嵘岁月稠。',
    '恰同学少年，风华正茂；',
    '书生意气，挥斥方遒。',
    '指点江山，激扬文字，',
    '粪土当年万户侯。',
    '曾记否，到中流击水，浪遏飞舟？',
]
add_multiline(s, MX, Inches(1.2), Inches(7), Inches(4.0),
              lines_xia, font_size=Pt(28), color=INK, line_spacing=2.2)

# 右侧提示
add_textbox(s, Inches(8.5), Inches(1.0), Inches(4), Inches(0.5),
            '忆事 · 言志', font_size=Pt(24), color=INK, bold=True)
add_line(s, Inches(8.5), Inches(1.55), Inches(1.2), Pt(2), color=FROST)
hints2 = [
    '上阕：写景（景中寓情）',
    '下阕：忆事（事中言志）',
    '',
    '一个"忆"字过渡',
    '从眼前回到往昔',
    '',
    '末句以反问收束',
    '振聋发聩',
]
add_multiline(s, Inches(8.5), Inches(1.8), Inches(3.5), Inches(3.0),
              hints2, font_size=Pt(16), color=GRAY, line_spacing=1.8)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '25', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 26. 过渡：怅寥廓 → 忆往昔 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '过渡 · 从景到情', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=INK)

# 三步骤
steps = [
    ('怅寥廓', '面对苍茫大地，心生感慨', '景的收束 → 情的起点', FROST),
    ('问苍茫大地', '向天地发问', '由己及天下', XIANG),
    ('谁主沉浮', '谁是大地的主宰？', '振聋发聩的时代之问', INK),
]
for i, (title, desc, transition, clr) in enumerate(steps):
    y = Inches(1.8) + i * Inches(1.6)
    add_circle(s, MX, y, Inches(0.7), fill_color=clr)
    add_textbox(s, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
                str(i+1), font_size=Pt(24), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.1), y - Inches(0.05), Inches(4), Inches(0.6),
                title, font_size=Pt(28), color=clr, bold=True)
    add_textbox(s, Inches(2.1), y + Inches(0.55), Inches(5), Inches(0.5),
                desc, font_size=Pt(16), color=INK)
    add_textbox(s, Inches(7.5), y + Inches(0.1), Inches(4), Inches(0.5),
                transition, font_size=Pt(14), color=GRAY)

add_rect(s, MX, Inches(6.2), Inches(10.5), Inches(0.7), fill_color=INK)
add_textbox(s, Inches(1.2), Inches(6.25), Inches(10), Inches(0.55),
            '"谁主沉浮"不是问句——是千古一叹，是革命者的宣言',
            font_size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '26', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 27. 同学少年形象 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '同学少年形象分析', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

qualities = [
    ('恰同学少年', '青春正好', '年龄', FROST),
    ('风华正茂', '风采才华正盛', '气质', XIANG),
    ('书生意气', '读书人的豪迈气概', '精神', WARM),
    ('挥斥方遒', '热情奔放，劲头正足', '行动', INK),
    ('指点江山', '纵论国家大事', '志向', FROST),
    ('激扬文字', '撰写激昂文章', '才华', XIANG),
    ('粪土万户侯', '视高官厚禄如粪土', '品格', INK),
]
for i, (phrase, meaning, dim, clr) in enumerate(qualities):
    y = Inches(1.3) + i * Inches(0.72)
    add_textbox(s, MX, y, Inches(2.5), Inches(0.5),
                phrase, font_size=Pt(18), color=clr, bold=True)
    add_textbox(s, Inches(3.8), y + Inches(0.02), Inches(3), Inches(0.5),
                meaning, font_size=Pt(15), color=INK)
    add_textbox(s, Inches(7.5), y + Inches(0.02), Inches(2), Inches(0.5),
                f'[{dim}]', font_size=Pt(13), color=GRAY)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '27', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 28. 中流击水 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '中流击水 · 浪遏飞舟', font_size=Pt(36), color=INK, bold=True)
add_line(s, MX, Inches(1.2), Inches(3), Pt(3), color=FROST)

add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
            '曾记否，到中流击水，浪遏飞舟？',
            font_size=Pt(48), color=INK, bold=True, alignment=PP_ALIGN.CENTER)

# 解读
interp = [
    ('典故', '祖逖"中流击楫"——誓死收复中原', FROST),
    ('象征', '革命者在时代洪流中奋勇搏击', XIANG),
    ('修辞', '反问——答案已在气势之中', WARM),
    ('情感', '豪迈、自信、舍我其谁', INK),
]
for i, (label, content, clr) in enumerate(interp):
    x = MX + i * Inches(2.8)
    add_rect(s, x, Inches(3.8), Inches(2.5), Inches(1.8), fill_color=RGBColor(255, 255, 255))
    add_textbox(s, x + Inches(0.1), Inches(3.95), Inches(2.3), Inches(0.4),
                label, font_size=Pt(18), color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.1), Inches(4.45), Inches(2.3), Inches(0.8),
                content, font_size=Pt(13), color=INK, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '28', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 29. 万户侯·中流击楫 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '两个典故 · 一个志向', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 左：万户侯
add_rect(s, MX, Inches(1.5), Inches(5.2), Inches(3.5), fill_color=RGBColor(248, 245, 240))
add_textbox(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.5),
            '粪土当年万户侯', font_size=Pt(28), color=FROST, bold=True)
add_multiline(s, Inches(1.2), Inches(2.4), Inches(4.5), Inches(2.5),
              ['万户侯：汉代侯爵', '食邑万户，权势极大', '',
               '在毛泽东笔下 →', '"粪土"：视如粪土',
               '', '蔑视权贵，追求理想'],
              font_size=Pt(15), color=INK, line_spacing=1.5)

# 右：中流击楫
add_rect(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(3.5), fill_color=RGBColor(248, 252, 250))
add_textbox(s, Inches(7.5), Inches(1.7), Inches(4.5), Inches(0.5),
            '中流击楫 / 中流击水', font_size=Pt(28), color=XIANG, bold=True)
add_multiline(s, Inches(7.5), Inches(2.4), Inches(4.5), Inches(2.5),
              ['祖逖(Tì)：东晋名将', '中流击楫，誓复中原', '',
               '毛泽东化用为"中流击水"', '从"击楫"到"击水"',
               '从誓言到行动，更为激越'],
              font_size=Pt(15), color=INK, line_spacing=1.5)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '29', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 30. PART 6 章节页 =====
s = blank_slide()
section_divider(s, 6, '对比升华', '悲秋与颂秋 · 谁主沉浮', accent_color=FROST)


# ===== 31. 古人悲秋 vs 毛泽东颂秋 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '悲秋与颂秋 · 古典与现代的对决', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 左：古人悲秋
add_rect(s, MX, Inches(1.5), Inches(5.2), Inches(5.5), fill_color=RGBColor(248, 245, 240))
add_textbox(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.5),
            '古典悲秋传统', font_size=Pt(24), color=GRAY, bold=True)
guren = [
    '宋玉："悲哉秋之为气也"',
    '杜甫："万里悲秋常作客"',
    '马致远："枯藤老树昏鸦"',
    '柳永："多情自古伤离别"',
    '',
    '秋 = 萧瑟、凋零、伤感',
    '悲凉是秋的底色',
]
add_multiline(s, Inches(1.2), Inches(2.3), Inches(4.5), Inches(4.0),
              guren, font_size=Pt(15), color=INK, line_spacing=1.5)

# 右：毛泽东颂秋
add_rect(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(5.5), fill_color=RGBColor(252, 245, 245))
add_textbox(s, Inches(7.5), Inches(1.7), Inches(4.5), Inches(0.5),
            '毛泽东颂秋', font_size=Pt(24), color=FROST, bold=True)
mao = [
    '万山红遍 → 热烈',
    '层林尽染 → 绚烂',
    '百舸争流 → 奋发',
    '鹰击长空 → 力量',
    '万类霜天竞自由 → 生机',
    '',
    '秋 = 热烈、生机、自由',
    '革命是秋的底色',
]
add_multiline(s, Inches(7.5), Inches(2.3), Inches(4.5), Inches(4.0),
              mao, font_size=Pt(15), color=INK, line_spacing=1.5)

# 底部
add_rect(s, MX, Inches(6.3), Inches(11.2), Inches(0.6), fill_color=FROST)
add_textbox(s, Inches(1.2), Inches(6.35), Inches(10), Inches(0.5),
            '同是秋天，一悲一颂 — 区别在于人的精神境界',
            font_size=Pt(16), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '31', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 32. 情感推导 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '情感推导', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 景 → 问 → 情 链条
chain = ['景', '万山红遍·鹰击鱼翔', '→', '问', '谁主沉浮', '→', '情', '革命壮志·时代担当']
x_positions = [Inches(1.0), Inches(3.0), Inches(5.0), Inches(5.8), Inches(7.8), Inches(9.6), Inches(10.4), Inches(11.5)]
for i in range(0, len(chain), 3):
    if i+2 < len(chain):
        label, detail, arrow = chain[i], chain[i+1], chain[i+2]
        x = x_positions[i] if i < len(x_positions) else Inches(1)
        add_rect(s, x, Inches(2.0), Inches(2.5), Inches(2.0), fill_color=RGBColor(248, 245, 240))
        add_textbox(s, x + Inches(0.1), Inches(2.1), Inches(2.3), Inches(0.5),
                    label, font_size=Pt(28), color=FROST, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), Inches(2.7), Inches(2.3), Inches(0.7),
                    detail, font_size=Pt(14), color=INK, alignment=PP_ALIGN.CENTER)

# 箭头
add_textbox(s, Inches(4.0), Inches(2.5), Inches(1), Inches(0.8),
            '→', font_size=Pt(36), color=FROST, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(9.2), Inches(2.5), Inches(1), Inches(0.8),
            '→', font_size=Pt(36), color=FROST, bold=True, alignment=PP_ALIGN.CENTER)

# 主旨
add_rect(s, MX, Inches(5.0), Inches(11.2), Inches(1.2), fill_color=INK)
add_textbox(s, Inches(1.2), Inches(5.1), Inches(10), Inches(0.5),
            '主旨：通过对湘江秋景的描绘和对青年时代革命斗争的回忆',
            font_size=Pt(16), color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(1.2), Inches(5.6), Inches(10), Inches(0.5),
            '抒发革命青年以天下为己任的壮志豪情',
            font_size=Pt(16), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '32', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 33. 板书设计 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '板书设计', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 上阕 下阕 双栏板书
# 左上
add_rect(s, Inches(0.9), Inches(1.5), Inches(5.5), Inches(5.2),
         fill_color=RGBColor(252, 248, 242))
add_textbox(s, Inches(1.2), Inches(1.6), Inches(4.5), Inches(0.5),
            '上阕：写景（景中寓情）', font_size=Pt(22), color=FROST, bold=True)
bbshang = [
    '独立寒秋 — 湘江北去 — 橘子洲头',
    '',
    '万山红遍    层林尽染  (远·静·色)',
    '漫江碧透    百舸争流  (近·动)',
    '鹰击长空    鱼翔浅底  (高·低)',
    '',
    '万类霜天竞自由',
    '  ↓',
    '怅寥廓，问苍茫大地，谁主沉浮？',
]
add_multiline(s, Inches(1.2), Inches(2.2), Inches(5.0), Inches(4.0),
              bbshang, font_size=Pt(14), color=INK, line_spacing=1.35)

# 右上
add_rect(s, Inches(6.9), Inches(1.5), Inches(5.5), Inches(5.2),
         fill_color=RGBColor(248, 250, 252))
add_textbox(s, Inches(7.2), Inches(1.6), Inches(4.5), Inches(0.5),
            '下阕：忆事（事中言志）', font_size=Pt(22), color=INK, bold=True)
bbxia = [
    '携来百侣曾游 → 忆往昔',
    '',
    '恰同学少年 · 风华正茂',
    '书生意气 · 挥斥方遒',
    '指点江山 · 激扬文字',
    '粪土当年万户侯',
    '',
    '曾记否——',
    '到中流击水，浪遏飞舟？',
]
add_multiline(s, Inches(7.2), Inches(2.2), Inches(5.0), Inches(4.0),
              bbxia, font_size=Pt(14), color=INK, line_spacing=1.35)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '33', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 34. 课堂练习 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '课堂练习', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=XIANG)

exercises = [
    ('1', '上阕描绘了哪些意象？这些意象有什么共同特点？（6分）', '基础'),
    ('2', '分析"击"和"翔"两个字为什么不能用"飞"和"游"替换？（6分）', '基础'),
    ('3', '"万类霜天竞自由"中"竞"字有何妙处？（4分）', '中等'),
    ('4', '比较古人悲秋与毛泽东颂秋的异同，谈谈你的理解。（8分）', '提高'),
]
for i, (num, question, level) in enumerate(exercises):
    y = Inches(1.5) + i * Inches(1.3)
    lvl_color = XIANG if level == '基础' else (WARM if level == '中等' else FROST)
    add_rect(s, MX, y, Inches(11.2), Inches(1.0), fill_color=RGBColor(248, 245, 240))
    add_textbox(s, Inches(1.1), y + Inches(0.05), Inches(0.5), Inches(0.4),
                f'Q{num}', font_size=Pt(18), color=lvl_color, bold=True)
    add_textbox(s, Inches(1.7), y + Inches(0.1), Inches(8), Inches(0.5),
                question, font_size=Pt(16), color=INK)
    add_rect(s, Inches(11.5), y + Inches(0.3), Inches(0.55), Inches(0.35), fill_color=lvl_color)
    add_textbox(s, Inches(11.5), y + Inches(0.3), Inches(0.55), Inches(0.35),
                f'{level[0]}', font_size=Pt(10), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '34', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 35. 课后作业 =====
s = blank_slide()
add_textbox(s, MX, Inches(0.3), Inches(5), Inches(0.7),
            '课后作业', font_size=Pt(32), color=INK, bold=True)
add_line(s, MX, Inches(1.1), Inches(2), Pt(3), color=FROST)

# 基础作业
add_textbox(s, MX, Inches(1.4), Inches(2), Inches(0.5),
            '基础作业', font_size=Pt(20), color=XIANG, bold=True)
basic_hw = ['1. 背诵全词并默写上阕', '2. 整理本节课的意象分析表', '3. 完成课堂练习第1-3题（书面）']
add_multiline(s, MX, Inches(2.0), Inches(10), Inches(1.5),
              basic_hw, font_size=Pt(16), color=INK, line_spacing=1.6)

# 提高作业
add_textbox(s, MX, Inches(3.6), Inches(2), Inches(0.5),
            '提高作业', font_size=Pt(20), color=FROST, bold=True)
adv_hw = ['4. 写一段200字的赏析：分析"万类霜天竞自由"在全词中的作用',
          '5. 课外阅读毛泽东《沁园春·雪》，比较两首词的异同',
          '6.（选做）查找"中流击楫"典故原文，写100字读后感']
add_multiline(s, MX, Inches(4.2), Inches(10), Inches(2.0),
              adv_hw, font_size=Pt(16), color=INK, line_spacing=1.6)
add_textbox(s, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4),
            '35', font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.RIGHT)


# ===== 36. 封底 =====
s = blank_slide()
add_rect(s, Inches(0), Inches(0), W, H, fill_color=INK)
add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
            '沁园春 · 长沙', font_size=Pt(72), color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_line(s, Inches(3.5), Inches(3.8), Inches(6), Pt(2), color=FROST)
add_textbox(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
            '恰同学少年，风华正茂；书生意气，挥斥方遒。',
            font_size=Pt(28), color=GOLD, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
            '高中语文 · 必修上册  |  第一单元  |  人教版2019',
            font_size=Pt(16), color=GRAY, alignment=PP_ALIGN.CENTER)


# ── 保存 ────────────────────────────────────────
os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_PDF)
print(f'PPTX saved: {OUT_PDF}')
print(f'Total slides: {len(prs.slides)}')

# 自检
from pptx import Presentation as PrsCheck
check = PrsCheck(OUT_PDF)
print(f'Self-check: {len(check.slides)} slides, file OK')
# 检查每个slide有内容
for i, slide in enumerate(check.slides):
    shapes_count = len(slide.shapes)
    if shapes_count < 2:
        print(f'  WARNING: slide {i+1} has only {shapes_count} shapes')
print('Self-check PASSED')
