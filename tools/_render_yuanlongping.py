# -*- coding: utf-8 -*-
"""《喜看稻菽千重浪——袁隆平的科学人生》课堂版 PPT（手工精排，14页，学生向）。
SOP：融合高赞讲课招牌招式 + 真实照片 + 杂志风 + 双检查。
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import _classroom_lib as L
from _classroom_lib import (new_presentation, bg, place_photo, scrim, textbox, rule,
                            kicker, new_slide, page_num, caption, quote_block, step_card, set_ea)
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
PHOTO_DIR = os.path.join(ROOT, 'preview_v7', '_yuanlongping_photos')
def PH(name):
    return os.path.join(PHOTO_DIR, name)

# ---------- 通用：对比表 ----------
def s_table(slide, x, y, w, h, headers, rows, colw):
    cols_n = len(headers); rows_n = len(rows)+1
    gf = slide.shapes.add_table(rows_n, cols_n, x, y, w, h).table
    gf.first_row = True; gf.horz_banding = False
    total = sum(colw)
    for i,cw in enumerate(colw):
        gf.columns[i].width = int(w*cw/total)
    for j,htxt in enumerate(headers):
        c = gf.cell(0,j)
        c.fill.solid(); c.fill.fore_color.rgb = L.INK
        c.margin_left=Inches(0.14); c.margin_right=Inches(0.14)
        c.margin_top=Inches(0.05); c.margin_bottom=Inches(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf=c.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=htxt; r.font.size=Pt(15); r.font.bold=True
        r.font.color.rgb=L.WHITE; r.font.name=L.HEI; set_ea(r,L.HEI)
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            c=gf.cell(i+1,j)
            c.fill.solid(); c.fill.fore_color.rgb = L.WHITE if i%2==0 else L.SOFT
            c.margin_left=Inches(0.14); c.margin_right=Inches(0.14)
            c.margin_top=Inches(0.04); c.margin_bottom=Inches(0.04)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf=c.text_frame; tf.word_wrap=True
            p=tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=val; r.font.size=Pt(13 if j==0 else 12.5)
            r.font.color.rgb=L.INK; r.font.name=L.HEI; set_ea(r,L.HEI)
    return gf

# ---------- P1 封面 ----------
def s_cover(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    place_photo(s, PH('real_rice_field.jpg'), 0, 0, L.W, L.H)
    scrim(s, 0, 0, L.W, Inches(1.7), L.INK, 0.5)
    scrim(s, 0, L.H-Inches(3.0), L.W, Inches(3.0), L.INK, 0.62)
    kicker(s, '必修上 · 第二单元「劳动光荣」', L.M, Inches(0.55), L.GOLD)
    textbox(s, L.M, Inches(2.55), Inches(12), Inches(1.4),
            [{'text':'喜看稻菽千重浪', 'size':58, 'color':L.WHITE, 'bold':True, 'font':L.KAI}])
    textbox(s, L.M, Inches(4.0), Inches(11.5), Inches(0.7),
            [{'text':'记首届国家最高科技奖获得者袁隆平', 'size':24, 'color':L.WHITE, 'font':L.HEI}])
    rule(s, L.M, Inches(4.85), Inches(2.4), L.GOLD, 2.6)
    textbox(s, L.M, Inches(5.05), Inches(11), Inches(0.5),
            [{'text':'沈英甲 · 人物通讯（实用性阅读与交流）', 'size':15, 'color':L.SOFT, 'font':L.HEI}])
    page_num(s, dark=True)

# ---------- P2 导览 ----------
def s_contents(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    kicker(s, '导览 · 本课脉络', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(11), Inches(0.8),
            [{'text':'一篇人物通讯，如何读出科学精神', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    items = [
        ('壹','导入 · 一碗米饭从何而来','袁隆平是科学劳动者，不是天才发明家'),
        ('贰','文体辨析 · 通讯与消息','真实、及时、典型、有人物'),
        ('叁','科研时间线','1960 → 1964 → 1973，14万株寻1株'),
        ('肆','精读 · 挑战权威','尊重事实，用实践检验理论'),
        ('伍','科学精神','尊重事实 · 质疑权威 · 实践检验 · 初心'),
        ('陆','知人论世 · 1960年代','为什么是「让所有人吃饱饭」'),
        ('柒','感动中国 · 颁奖词','一介农夫，毕生远离饥饿'),
        ('捌','实用文阅读三步法','圈信息 → 理脉络 → 析精神'),
        ('玖','分层作业','基础 · 提高 · 拓展'),
    ]
    y = 2.25; lh = 0.52
    for num, t, d in items:
        textbox(s, L.M, Inches(y), Inches(0.7), Inches(0.5),
                [{'text':num, 'size':22, 'color':L.FROST, 'bold':True, 'font':L.KAI}])
        textbox(s, L.M+Inches(0.85), Inches(y-0.02), Inches(10.5), Inches(0.5),
                [{'text':t, 'size':17, 'color':L.INK, 'bold':True, 'font':L.HEI},
                 {'text':'   '+d, 'size':13, 'color':L.MUTED, 'font':L.HEI}])
        rule(s, L.M, Inches(y+0.5), L.CW, L.SOFT, 1.4)
        y += lh + 0.06
    page_num(s)

# ---------- P3 导入 ----------
def s_import(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    # 右侧照片
    px = L.W - L.M - Inches(4.2)
    place_photo(s, PH('real_rice_field.jpg'), px, L.M, Inches(4.2), Inches(6.1))
    caption(s, '图源：Wikimedia Commons', px, L.M+Inches(6.15), Inches(4.2))
    kicker(s, '壹 · 导入', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(7.2), Inches(0.9),
            [{'text':'一碗米饭，从何而来？', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    textbox(s, L.M, Inches(2.35), Inches(7.3), Inches(2.6), [
        {'text':'· 全中国 14 亿人，每天要吃饭', 'size':16, 'color':L.INK, 'font':L.HEI, 'space_after':8},
        {'text':'· 粮食，从一片片稻田里长出来', 'size':16, 'color':L.INK, 'font':L.HEI, 'space_after':8},
        {'text':'· 有一个人，让中国人把饭碗端在自己手里', 'size':16, 'color':L.INK, 'font':L.HEI, 'space_after':8},
        {'text':'—— 袁隆平', 'size':20, 'color':L.FROST, 'bold':True, 'font':L.KAI}])
    # 强调框
    card = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, L.M, Inches(5.15), Inches(7.3), Inches(1.55))
    card.fill.solid(); card.fill.fore_color.rgb = L.WHITE; card.line.color.rgb = L.FROST; card.line.width = Pt(1.8); card.shadow.inherit=False
    textbox(s, L.M+Inches(0.25), Inches(5.3), Inches(6.8), Inches(1.3), [
        {'text':'他是「科学劳动者」', 'size':16, 'color':L.FROST, 'bold':True, 'font':L.HEI, 'space_after':5},
        {'text':'成就是在田间地头一棵棵稻子里干出来的，不是凭空而来的「天才发明家」。本单元的主题是「劳动光荣」。', 'size':13.5, 'color':L.INK, 'font':L.HEI, 'line':1.35}])
    page_num(s)

# ---------- P4 文体辨析 ----------
def s_genre(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    kicker(s, '贰 · 文体辨析', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(11), Inches(0.8),
            [{'text':'通讯 与 消息，有什么不同？', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    s_table(s, L.M, Inches(2.35), L.CW, Inches(3.4),
            ['维度','消息','通讯（本文）'],
            [
                ['篇幅','短、简、快','长，篇幅舒展'],
                ['写法','只说发生了什么','有人物描写与细节'],
                ['人物','一般不刻画','有形象、有性格'],
                ['时效','争分夺秒','相对从容'],
                ['真实','必须真实','必须真实，可用文学手法让真实更生动'],
            ], [1.0, 1.3, 2.0])
    # 结论条
    band = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, L.M, Inches(6.05), L.CW, Inches(0.85))
    band.fill.solid(); band.fill.fore_color.rgb = L.INK; band.line.fill.background(); band.shadow.inherit=False
    textbox(s, L.M+Inches(0.3), Inches(6.13), L.CW-Inches(0.6), Inches(0.7),
            [{'text':'结论：通讯也是新闻，不能虚构；但它用文学手法，让真实的人与事更有感染力。', 'size':14.5, 'color':L.WHITE, 'bold':True, 'font':L.HEI}])
    page_num(s)

# ---------- P5 科研时间线 ----------
def s_timeline(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    kicker(s, '叁 · 科研时间线', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(11), Inches(0.8),
            [{'text':'一条稻田里的科学之路', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    nodes = [
        ('1960','发现天然杂交稻',['观察到天然杂交稻株','（现象：水稻也能杂交）'], L.XIANG),
        ('1964','找到雄性不育株',['7月5日，田间寻得','14万株中仅 1 株','（材料：实验起点）'], L.FROST),
        ('1973','三系杂交稻成功',['三系配套培育成功','中国人吃饱饭的希望','（结果：从0到1）'], L.GOLD),
    ]
    cw = Inches(3.37); gap = Inches(0.4); x0 = L.M
    for i,(yr,title,body,acc) in enumerate(nodes):
        x = int(x0) + i*(int(cw)+int(gap))
        step_card(s, x, Inches(2.3), cw, Inches(3.0), yr, title, body, acc)
    # 强调带
    band = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, L.M, Inches(5.6), L.CW, Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = L.SOFT; band.line.color.rgb = L.FROST; band.line.width=Pt(1.6); band.shadow.inherit=False
    textbox(s, L.M+Inches(0.3), Inches(5.7), L.CW-Inches(0.6), Inches(0.85), [
        {'text':'14 万株 → 1 株', 'size':17, 'color':L.FROST, 'bold':True, 'font':L.HEI, 'space_after':3},
        {'text':'这是「劳动量」，不是「灵感」。袁隆平在稻田里一棵一棵检查了 14 万株水稻，才找到那株不育株。', 'size':13.5, 'color':L.INK, 'font':L.HEI, 'line':1.3}])
    page_num(s)

# ---------- P6 精读·挑战权威 ----------
def s_authority(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    px = L.W - L.M - Inches(4.2)
    place_photo(s, PH('real_terrace.jpg'), px, L.M, Inches(4.2), Inches(6.1))
    caption(s, '图源：Wikimedia Commons', px, L.M+Inches(6.15), Inches(4.2))
    kicker(s, '肆 · 精读', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(7.2), Inches(0.9),
            [{'text':'敢向权威定论发问', 'size':29, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    textbox(s, L.M, Inches(2.3), Inches(7.3), Inches(4.4), [
        {'text':'权威定论', 'size':16, 'color':L.FROST, 'bold':True, 'font':L.HEI, 'space_after':4},
        {'text':'「自花授粉作物（如水稻）没有杂交优势。」', 'size':15, 'color':L.INK, 'font':L.HEI, 'space_after':12},
        {'text':'袁隆平怎么做', 'size':16, 'color':L.XIANG, 'bold':True, 'font':L.HEI, 'space_after':4},
        {'text':'他不信，转身下田，用实践去检验——', 'size':15, 'color':L.INK, 'font':L.HEI, 'space_after':4},
        {'text':'跨、迈、蹲、翻，一株一株地看。', 'size':15, 'color':L.INK, 'font':L.HEI, 'space_after':4},
        {'text':'事实，胜过定论。', 'size':15, 'color':L.INK, 'font':L.HEI}])
    band = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, L.M, Inches(6.0), Inches(7.3), Inches(0.95))
    band.fill.solid(); band.fill.fore_color.rgb = L.INK; band.line.fill.background(); band.shadow.inherit=False
    textbox(s, L.M+Inches(0.3), Inches(6.08), Inches(6.8), Inches(0.8),
            [{'text':'精神：尊重事实 → 质疑权威 → 实践检验', 'size':15, 'color':L.WHITE, 'bold':True, 'font':L.HEI}])
    page_num(s)

# ---------- P7 科学精神 ----------
def s_spirit(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    px = L.W - L.M - Inches(3.7)
    place_photo(s, PH('real_yuanlongping.jpg'), px, L.M, Inches(3.7), Inches(4.7))
    caption(s, '袁隆平（图源：Wikimedia Commons）', px, L.M+Inches(4.75), Inches(3.7))
    kicker(s, '伍 · 科学精神', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(8), Inches(0.8),
            [{'text':'一种精神，四个支点', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    cards = [
        ('尊重事实','用实践检验理论，不盲从书本', L.XIANG),
        ('质疑权威','敢于向定论发问，并去验证', L.FROST),
        ('实践检验','田间数十年，劳动出真知', L.GOLD),
        ('初心','让所有人吃饱饭', L.INK),
    ]
    x0 = L.M; cw = Inches(3.7); gap = Inches(0.35)
    for i,(t,b,acc) in enumerate(cards):
        x = int(x0) + i*(int(cw)+int(gap))
        y = Inches(2.3) if i<2 else Inches(4.5)
        xx = int(x0) + (i%2)*(int(cw)+int(gap))
        step_card(s, xx, y, cw, Inches(1.95), i+1, t, [b], acc)
    textbox(s, L.M, Inches(6.65), Inches(8), Inches(0.5),
            [{'text':'科学是手段，吃饱是目的，劳动是桥梁。', 'size':14, 'color':L.MUTED, 'italic':True, 'font':L.KAI}])
    page_num(s)

# ---------- P8 知人论世 ----------
def s_background(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    px = L.W - L.M - Inches(4.2)
    place_photo(s, PH('real_golden_field.jpg'), px, L.M, Inches(4.2), Inches(6.1))
    caption(s, '图源：Wikimedia Commons', px, L.M+Inches(6.15), Inches(4.2))
    kicker(s, '陆 · 知人论世', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(7.2), Inches(0.9),
            [{'text':'为什么是「让所有人吃饱饭」', 'size':27, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    textbox(s, L.M, Inches(2.35), Inches(7.3), Inches(4.4), [
        {'text':'1960 年代，粮食短缺，许多人吃不饱饭。', 'size':16, 'color':L.INK, 'font':L.HEI, 'space_after':10},
        {'text':'袁隆平立下志向：', 'size':16, 'color':L.FROST, 'bold':True, 'font':L.HEI, 'space_after':6},
        {'text':'要让中国人、让世界上所有人，都吃饱饭。', 'size':16, 'color':L.INK, 'font':L.HEI, 'space_after':10},
        {'text':'所以他的科学精神背后，还藏着一颗「心系苍生」的初心。', 'size':15, 'color':L.INK, 'font':L.HEI, 'line':1.4}])
    band = s.shapes.add_shape(L.MSO_SHAPE.ROUNDED_RECTANGLE, L.M, Inches(6.0), Inches(7.3), Inches(0.95))
    band.fill.solid(); band.fill.fore_color.rgb = L.INK; band.line.fill.background(); band.shadow.inherit=False
    textbox(s, L.M+Inches(0.3), Inches(6.08), Inches(6.8), Inches(0.8),
            [{'text':'他不是「为获奖而研究」——获奖，是因为研究解决了人类的饥饿。', 'size':14, 'color':L.WHITE, 'bold':True, 'font':L.HEI}])
    page_num(s)

# ---------- P9 感动中国·颁奖词 ----------
def s_award(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    px = L.W - L.M - Inches(4.0)
    place_photo(s, PH('real_yuanlongping.jpg'), px, L.M, Inches(4.0), Inches(5.4))
    caption(s, '袁隆平（图源：Wikimedia Commons）', px, L.M+Inches(5.45), Inches(4.0))
    kicker(s, '柒 · 感动中国', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(7.8), Inches(0.8),
            [{'text':'一介农夫，毕生远离饥饿', 'size':27, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    quote_block(s, L.M, Inches(2.3), Inches(7.9),
                '他是一位真正的耕耘者。当名满天下，却仍只专注于田畴；淡泊名利，一介农夫，播撒智慧，收获富足。他毕生的梦想，就是让所有的人远离饥饿。',
                '2004 感动中国年度人物 · 颁奖词', L.FROST)
    textbox(s, L.M, Inches(5.7), Inches(7.8), Inches(1.0), [
        {'text':'从「乡村教师」到「名满天下」，变的是声名，不变的是田畴里的专注。', 'size':14.5, 'color':L.INK, 'font':L.HEI, 'line':1.4}])
    page_num(s)

# ---------- P10 实用文三步法 ----------
def s_method(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    kicker(s, '捌 · 阅读方法', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(11), Inches(0.8),
            [{'text':'实用文阅读三步法', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    steps = [
        ('1','圈关键信息','读长文不「全部划线」。抓时间、人物、事件、数据，只圈最关键的。', L.XIANG),
        ('2','理事件脉络','把零散信息连成一条线：起因→经过→结果，如科研时间线。', L.FROST),
        ('3','析人物精神','从事件里读出人的精神：尊重事实、质疑权威、心系苍生。', L.GOLD),
    ]
    cw = Inches(3.71); gap = Inches(0.4); x0 = L.M
    for i,(n,t,b,acc) in enumerate(steps):
        x = int(x0) + i*(int(cw)+int(gap))
        step_card(s, x, Inches(2.4), cw, Inches(3.2), n, t, [b], acc)
    textbox(s, L.M, Inches(5.9), L.CW, Inches(0.6),
            [{'text':'这套方法，下一篇通讯《心有一团火》（张秉贵）照样用得上。', 'size':14.5, 'color':L.MUTED, 'italic':True, 'font':L.KAI}])
    page_num(s)

# ---------- P11 单元呼应·劳动光荣 ----------
def s_unit_echo(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    place_photo(s, PH('real_rice_field.jpg'), 0, 0, L.W, L.H)
    scrim(s, 0, 0, L.W, Inches(2.0), L.INK, 0.5)
    scrim(s, 0, L.H-Inches(5.2), L.W, Inches(5.2), L.INK, 0.66)
    kicker(s, '单元呼应 · 劳动光荣', L.M, Inches(0.55), L.GOLD)
    textbox(s, L.M, Inches(2.5), Inches(11.5), Inches(0.9),
            [{'text':'科学劳动，也是劳动', 'size':32, 'color':L.WHITE, 'bold':True, 'font':L.KAI}])
    textbox(s, L.M, Inches(3.5), Inches(11.5), Inches(2.6), [
        {'text':'袁隆平在稻田里劳动，张秉贵在柜台前劳动，钟扬在高原上劳动。', 'size':17, 'color':L.WHITE, 'font':L.HEI, 'space_after':10},
        {'text':'劳动光荣，不分工种。', 'size':20, 'color':L.GOLD, 'bold':True, 'font':L.KAI, 'space_after':10},
        {'text':'下节课，我们读另外两双手——售货员的一抓准，教授的高原采种。', 'size':15, 'color':L.SOFT, 'font':L.HEI, 'line':1.4}])
    page_num(s, dark=True)

# ---------- P12 板书 ----------
def s_blackboard(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.INK)
    kicker(s, '板书 · 课堂框架', L.M, L.M, L.GOLD)
    textbox(s, L.M, Inches(1.2), Inches(11), Inches(0.7),
            [{'text':'喜看稻菽千重浪', 'size':28, 'color':L.WHITE, 'bold':True, 'font':L.KAI}])
    cards = [
        ('科研时间线','1960 发现天然杂交稻\n1964 找到雄性不育株\n1973 三系杂交稻成功\n（14万株 → 1株）', L.XIANG),
        ('科学精神','尊重事实 → 质疑权威\n→ 实践检验\n初心：让所有人吃饱饭', L.FROST),
        ('实用文三步法','圈关键信息\n→ 理事件脉络\n→ 析人物精神', L.GOLD),
    ]
    cw = Inches(3.71); gap = Inches(0.4); x0 = L.M
    for i,(t,b,acc) in enumerate(cards):
        x = int(x0) + i*(int(cw)+int(gap))
        step_card(s, x, Inches(2.3), cw, Inches(3.4), i+1, t, b.split('\n'), acc)
    rule(s, L.M, Inches(6.1), L.CW, L.GOLD, 2.2)
    textbox(s, L.M, Inches(6.25), L.CW, Inches(0.6),
            [{'text':'劳动光荣：科学劳动者的坚守与创新', 'size':15, 'color':L.GOLD, 'bold':True, 'font':L.HEI}])
    page_num(s, dark=True)

# ---------- P13 分层作业 ----------
def s_homework(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    kicker(s, '玖 · 分层作业', L.M, L.M, L.FROST)
    textbox(s, L.M, Inches(1.25), Inches(11), Inches(0.8),
            [{'text':'把课文读厚，也读薄', 'size':30, 'color':L.INK, 'bold':True, 'font':L.KAI}])
    cards = [
        ('基础作业',[
            '1. 完成袁隆平科研时间线（时间｜事件），至少 3 个节点。',
            '2. 解释「通讯」的文体特征，至少说出 3 点。'], L.FROST),
        ('提高作业',[
            '写一段 150 字分析：',
            '《从 14 万株到 1 株看袁隆平的科学精神》。',
            '要求：引原文数据＋析劳动量＋点明精神。'], L.XIANG),
        ('拓展作业',[
            '用「三步法」试读下一篇通讯',
            '《心有一团火》（张秉贵），',
            '圈出关键信息，理出他的「绝活」。'], L.GOLD),
    ]
    cw = Inches(3.71); gap = Inches(0.4); x0 = L.M
    for i,(t,body,acc) in enumerate(cards):
        x = int(x0) + i*(int(cw)+int(gap))
        step_card(s, x, Inches(2.35), cw, Inches(3.6), '·', t, body, acc)
    page_num(s)

# ---------- P14 青春·劳动价值 ----------
def s_closing(prs, BLANK):
    s = new_slide(prs, BLANK)
    bg(s, L.PAPER)
    place_photo(s, PH('real_golden_field.jpg'), 0, 0, L.W, L.H)
    scrim(s, 0, 0, L.W, Inches(1.7), L.INK, 0.5)
    scrim(s, 0, L.H-Inches(3.4), L.W, Inches(3.4), L.INK, 0.64)
    kicker(s, '青春 · 劳动价值', L.M, Inches(0.55), L.GOLD)
    textbox(s, L.M, Inches(2.6), Inches(11.5), Inches(0.9),
            [{'text':'把论文写在大地上', 'size':34, 'color':L.WHITE, 'bold':True, 'font':L.KAI}])
    textbox(s, L.M, Inches(3.7), Inches(11.5), Inches(2.4), [
        {'text':'真正的学问，不在书斋里，而在祖国需要的土地上。', 'size':18, 'color':L.WHITE, 'font':L.HEI, 'space_after':12},
        {'text':'无论将来做什么——科学、服务、教育——', 'size':17, 'color':L.SOFT, 'font':L.HEI, 'space_after':8},
        {'text':'认真干的劳动，都光荣。', 'size':20, 'color':L.GOLD, 'bold':True, 'font':L.KAI}])
    page_num(s, dark=True)

# ---------- BUILD ----------
def build():
    prs, BLANK = new_presentation()
    s_cover(prs, BLANK)
    s_contents(prs, BLANK)
    s_import(prs, BLANK)
    s_genre(prs, BLANK)
    s_timeline(prs, BLANK)
    s_authority(prs, BLANK)
    s_spirit(prs, BLANK)
    s_background(prs, BLANK)
    s_award(prs, BLANK)
    s_method(prs, BLANK)
    s_unit_echo(prs, BLANK)
    s_blackboard(prs, BLANK)
    s_homework(prs, BLANK)
    s_closing(prs, BLANK)
    out = os.path.join(ROOT, 'preview_v7', 'yuanlongping.pptx')
    prs.save(out)
    print('SAVED', out, 'pages=', len(prs.slides._sldIdLst))

if __name__ == '__main__':
    build()
