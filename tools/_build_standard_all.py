# -*- coding: utf-8 -*-
"""逐课「手精标准 + 四层审计门禁」驱动（857 课统一执行）。

对 _lessons_all.json 中每一课：
  1) 用升级版 mapper（学生向结构 + 设计系统 + 照片地图接口）渲染 PPTX
  2) 四层审计：教师口吻(BLOCK) / 文字溢出 / 照片遮罩 / 越界
  3) 溢出则自动缩放字号(1.0→0.92→0.85)重渲再审
  4) 记录 PASS / WARN / FAIL + 具体问题 到 JSONL（断点续跑）

输出: preview_v7/<subject>/<id>.pptx  (覆盖旧懒人批量版)
日志: preview_v7/_STANDARD_AUDIT.jsonl
"""
import os, sys, json, time
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

from _classroom_ppt import render_lesson
from _audit_text import audit as audit_text
from _audit_layout import audit as audit_layout
from _audit_scrim import audit as audit_scrim
from _photo_map import get_photos
from pptx import Presentation

BASE = os.path.dirname(HERE)
JSON_PATH = os.path.join(BASE, 'preview_v7', '_lessons_all.json')
LOG = os.path.join(BASE, 'preview_v7', '_STANDARD_AUDIT.jsonl')

W_IN, H_IN, EPS = 13.333, 7.5, 0.02
EMU = 914400

def check_bounds(path):
    prs = Presentation(path); bad = []
    for i, sl in enumerate(prs.slides, 1):
        over = False
        for sh in sl.shapes:
            if sh.left is None or sh.width is None:
                continue
            r = (sh.left + sh.width) / EMU
            b = (sh.top + sh.height) / EMU if (sh.top is not None and sh.height is not None) else 0
            if r > W_IN + EPS or b > H_IN + EPS:
                over = True; break
        if over:
            bad.append(i)
    return bad

def audit_one(path):
    rows, bh, wh = audit_text(path)
    lay = audit_layout(path)
    scr = audit_scrim(path)
    bnd = check_bounds(path)
    overflow = [s for s, iss in lay if any(x.startswith('TEXT_OVERFLOW') for x in iss)]
    photo = [s for s, iss in lay if any('PHOTO+TEXT' in x for x in iss)]
    return {
        'slides': len(rows),
        'block': len(bh),
        'warn': len(wh),
        'overflow': overflow,
        'photo_txt': photo,
        'scrim_unprotected': len(scr),
        'bounds_over': bnd,
    }

def render_audit(lesson, path, photos):
    """渲染并尝试自动修复溢出（缩放字号）。返回最终审计结果。"""
    last = None
    for scale in (1.0, 0.92, 0.85, 0.78, 0.70):
        render_lesson(lesson, path, photos=photos, scale=scale)
        a = audit_one(path)
        last = a
        if a['overflow'] or a['scrim_unprotected']:
            continue  # 还有溢出/遮罩问题 → 缩小再试
        break
    return last

def load_done():
    done = {}
    if os.path.exists(LOG):
        with open(LOG, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    rec = json.loads(line)
                except Exception:
                    continue
                # 仅跳过已 PASS 的课；WARN/FAIL 需以修复后的 mapper 重渲
                if rec.get('verdict') == 'PASS':
                    done[rec['id']] = rec['verdict']
    return done

def main():
    data = json.load(open(JSON_PATH, 'r', encoding='utf-8'))
    done = load_done()
    f = open(LOG, 'a', encoding='utf-8')
    total = passed = warned = failed = errored = 0
    t0 = time.time()
    for subj, lessons in data.items():
        for les in lessons:
            lid = les.get('id') or ('lesson_%d' % total)
            if lid in done:
                total += 1
                continue
            out = os.path.join(BASE, 'preview_v7', subj, lid + '.pptx')
            photos = get_photos(les)
            try:
                a = render_audit(les, out, photos)
            except Exception as e:
                rec = {'id': lid, 'subj': subj, 'verdict': 'ERROR',
                       'error': repr(e)[:240]}
                f.write(json.dumps(rec, ensure_ascii=False) + '\n'); f.flush()
                errored += 1; total += 1
                print('[%d] %s ERROR %s' % (total, lid, rec['error']), flush=True)
                continue
            if a['block'] > 0:
                verdict = 'FAIL'; failed += 1
            elif a['overflow'] or a['scrim_unprotected'] or a['bounds_over']:
                verdict = 'WARN'; warned += 1
            else:
                verdict = 'PASS'; passed += 1
            rec = {'id': lid, 'subj': subj, **a, 'verdict': verdict}
            f.write(json.dumps(rec, ensure_ascii=False) + '\n'); f.flush()
            total += 1
            el = time.time() - t0
            print('[%d/866] %s %s block=%d overflow=%s scrim=%d bounds=%s (%.1fs)' % (
                total, lid, verdict, a['block'], a['overflow'],
                a['scrim_unprotected'], a['bounds_over'], el), flush=True)
    f.close()
    print('==== DONE total=%d PASS=%d WARN=%d FAIL=%d ERROR=%d ====' % (
        total, passed, warned, failed, errored), flush=True)

if __name__ == '__main__':
    main()
