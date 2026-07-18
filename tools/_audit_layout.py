# tools/_audit_layout.py
# Highest-standard LAYOUT audit: catches what _check_bounds cannot.
#  1) TEXT_OVERFLOW: estimates rendered text height (PIL wrap, same as _pptx_to_png)
#     and compares to the shape's box height. Flags text spilling out of its box.
#  2) PHOTO_TEXT: flags slides where a text shape overlaps a picture (needs scrim/panel).
#  3) EMPTY: slides with neither text nor picture.
import os, glob
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

DPI = 110
SCALE = DPI / 72.0
EMU_PER_IN = 914400
FONT_DIR = r"C:\Windows\Fonts"
# Cross-platform CJK font fallbacks so the overflow estimate is accurate off-Windows.
_MAC_FONTS = {
    "msyh":   "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "simhei": "/System/Library/Fonts/STHeiti Medium.ttc",
    "hei":    "/System/Library/Fonts/STHeiti Medium.ttc",
    "simsun": "/System/Library/Fonts/Supplemental/Songti.ttc",
    "song":   "/System/Library/Fonts/Supplemental/Songti.ttc",
    "simkai": "/System/Library/Fonts/Supplemental/Songti.ttc",
    "kai":    "/System/Library/Fonts/Supplemental/Songti.ttc",
    "simfang":"/System/Library/Fonts/Supplemental/Songti.ttc",
    "fang":   "/System/Library/Fonts/Supplemental/Songti.ttc",
}
_MAC_DEFAULT = "/System/Library/Fonts/STHeiti Medium.ttc"
def _f(name, size_px, bold=False):
    path = os.path.join(FONT_DIR, name + (".ttc" if name in ("msyh","simsun") else ".ttf"))
    if not os.path.exists(path): path = os.path.join(FONT_DIR, "simhei.ttf")
    if not os.path.exists(path):
        # off-Windows: resolve a real CJK font by family key
        key = name.lower()
        path = None
        for k, p in _MAC_FONTS.items():
            if k in key and os.path.exists(p):
                path = p; break
        if not path:
            path = _MAC_DEFAULT if os.path.exists(_MAC_DEFAULT) else None
    try: return ImageFont.truetype(path, size_px)
    except Exception: return ImageFont.load_default()
def emu_px(v): return 0 if v is None else int(round(v/EMU_PER_IN*DPI))
def run_font(run):
    n=(run.font.name or "").lower()
    for k in ("msyh","simhei","simkai","simsun","simfang","kai","hei","song","fang"):
        if k in n: return k
    return "simhei"
def run_size(run):
    if run.font.size is not None: return max(8, int(round(run.font.size.pt*SCALE)))
    return max(8, int(round(18*SCALE)))
def est_text_h(tf, box_w_px):
    """Estimate total rendered height (px) of a text frame at given box width."""
    d = ImageDraw.Draw(Image.new("RGB",(10,10)))
    total = 0
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            total += int(20*SCALE); continue
        size = run_size(runs[0]); bold=bool(runs[0].font.bold)
        fnt = _f(run_font(runs[0]), size, bold)
        ptext = "".join(r.text for r in runs)
        sb = int((para.space_before.pt if para.space_before else 0)*SCALE)
        sa = int((para.space_after.pt if para.space_after else 0)*SCALE)
        ls = para.line_spacing
        mul = ls if isinstance(ls,float) else (1.2 if ls is None else 1.2)
        # wrap
        lines, cur = [], ""
        for ch in ptext:
            if ch=="\n":
                if cur: lines.append(cur)
                cur=""
                continue
            trial=cur+ch
            if d.textlength(trial,font=fnt)<=box_w_px-8*SCALE or not cur: cur=trial
            else: lines.append(cur); cur=ch
        if cur: lines.append(cur)
        lh = int(size*mul)
        total += sb + len(lines)*lh + sa
    return total

def audit(pptx):
    prs = Presentation(pptx)
    findings = []
    for i, slide in enumerate(prs.slides, 1):
        has_pic=False; pic_boxes=[]; text_shapes=[]; any_text=False
        over=[]
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_pic=True
                pic_boxes.append((emu_px(sh.left),emu_px(sh.top),emu_px(sh.width),emu_px(sh.height)))
            if sh.has_text_frame and sh.text_frame.text.strip():
                any_text=True
                bw=emu_px(sh.width); bh=emu_px(sh.height)
                if bw>0 and bh>0:
                    h=est_text_h(sh.text_frame, bw)
                    if h > bh*1.08:
                        over.append((round(h,0), round(bh,0), sh.text_frame.text[:24]))
        # photo-text overlap triage
        photo_txt=False
        if has_pic:
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    sl,st,sw,shh=emu_px(sh.left),emu_px(sh.top),emu_px(sh.width),emu_px(sh.height)
                    for (pl,pt,pw,ph) in pic_boxes:
                        if not(sl+sw<=pl or sl>=pl+pw or st+shh<=pt or st>=pt+ph):
                            photo_txt=True; break
                    if photo_txt: break
        issues=[]
        if over: issues.append("TEXT_OVERFLOW:"+";".join(f"h{thi}>box{bi}({tx})" for thi,bi,tx in over))
        if photo_txt: issues.append("PHOTO+TEXT(verify scrim/panel)")
        if not any_text and not has_pic: issues.append("EMPTY")
        if issues:
            findings.append((i, issues))
    return findings

if __name__ == '__main__':
    for f in sorted(glob.glob('preview_v7/*.pptx')):
        fn=os.path.basename(f)
        try:
            fnd=audit(f)
        except Exception as e:
            print(f"{fn}: ERROR {e}"); continue
        if fnd:
            print(f"\n### {fn}  ({len(fnd)} slide(s) with layout issues)")
            for i,iss in fnd:
                print(f"   slide {i:>2}: " + " | ".join(iss))
        else:
            print(f"{fn}: LAYOUT OK ✅")
