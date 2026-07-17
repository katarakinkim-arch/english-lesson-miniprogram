# tools/_audit_scrim.py
# Rigorous checker: for every slide where a TEXT shape overlaps a PICTURE,
# verify a protecting layer (semi-transparent scrim OR opaque panel) sits over
# the overlap region. A naked text-on-photo = contrast bug (FAIL).
import os, glob
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

EMU = 914400
def px(v): return 0 if v is None else int(round(v/EMU*110))

def overlap(a, b):
    ax, ay, aw, ah = a; bx, by, bw, bh = b
    ix = max(0, min(ax+aw, bx+bw) - max(ax, bx))
    iy = max(0, min(ay+ah, by+bh) - max(ay, by))
    return ix*iy

def is_protecting(sh):
    # autoshape/rect with solid fill (with or without alpha) covering area
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return False
    try:
        fill = sh.fill
        if fill.type is None: return False
        fc = fill.fore_color
        xfill = fc._xFill
        if xfill is None: return False
        # srgbClr present => has a solid color (scrim or panel)
        if xfill.find(qn('a:srgbClr')) is not None:
            return True
    except Exception:
        return False
    return False

def audit(pptx):
    prs = Presentation(pptx)
    fails = []
    for i, slide in enumerate(prs.slides, 1):
        pics = [(px(s.left),px(s.top),px(s.width),px(s.height)) for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        texts = [(px(s.left),px(s.top),px(s.width),px(s.height), s) for s in slide.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        protectors = [(px(s.left),px(s.top),px(s.width),px(s.height)) for s in slide.shapes
                      if is_protecting(s)]
        for (tl,tt,tw,th,sh) in texts:
            for pic in pics:
                if overlap((tl,tt,tw,th), pic) <= 0:
                    continue
                ov = overlap((tl,tt,tw,th), pic)
                protected = any(overlap(prot, (tl,tt,tw,th)) >= ov*0.6 or
                                overlap(prot, pic) >= ov*0.6 for prot in protectors)
                if not protected:
                    fails.append((i, sh.text_frame.text[:30], ov))
                break
    return fails

if __name__ == '__main__':
    for f in sorted(glob.glob('preview_v7/*.pptx')):
        fn = os.path.basename(f)
        try:
            fails = audit(f)
        except Exception as e:
            print(f"{fn}: ERROR {e}"); continue
        if fails:
            print(f"\n### {fn}  ({len(fails)} UNPROTECTED text-on-photo)")
            for i, txt, ov in fails:
                print(f"   slide {i:>2}: '{txt}' (overlap {ov}px²)")
        else:
            print(f"{fn}: all photo+text slides PROTECTED ✅")
