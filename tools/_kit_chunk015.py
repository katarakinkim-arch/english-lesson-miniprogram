# -*- coding: utf-8 -*-
"""chunk015 共享渲染工具：生物（及非语文）学生版 9 页课堂 PPT。

设计：仅依赖 _classroom_lib；不引入任何照片（生物无适配自由授权实拍图，
按规范§2.3 用学科色块兜底并标注「本课件未使用无关配图」）。
关键保证：所有文本框经 PIL 自适应字号（与 _audit_layout 度量一致），
杜绝 TEXT_OVERFLOW；无照片 => 遮罩审计天然 PASS。
"""
import os, json, re
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont

from _classroom_lib import (
    PAPER, INK, FROST, XIANG, MUTED, WHITE, GOLD, SOFT,
    KAI, HEI, W, H, M, CW, Inches, Pt,
    bg, textbox, rule, kicker, new_slide, page_num, caption, new_presentation,
)

# ---------- PIL 度量（与 _audit_layout 完全一致） ----------
DPI = 110
SCALE = DPI / 72.0
EMU_PER_IN = 914400
_MAC_FONTS = {
    "msyh":   "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "simhei": "/System/Library/Fonts/STHeiti Medium.ttc",
    "hei":    "/System/Library/Fonts/STHeiti Medium.ttc",
    "simsun": "/System/Library/Fonts/Supplemental/Songti.ttc",
    "song":   "/System/Library/Fonts/Supplemental/Songti.ttc",
    "simkai": "/System/Library/Fonts/Supplemental/Songti.ttc",
    "kai":    "/System/Library/Fonts/Supplemental/Songti.ttc",
    "simfang":"/System/Library/Fonts/Supplemental/Songti.ttc",
    "fang":   "/System/Library/Fonts/Supplemental/Songti.ttc",
}
_MAC_DEFAULT = "/System/Library/Fonts/STHeiti Medium.ttc"

def _f(name, size_px, bold=False):
    path = os.path.join("/System/Library/Fonts", name + (".ttc" if name in ("msyh", "simsun") else ".ttf"))
    if not os.path.exists(path):
        path = os.path.join("/System/Library/Fonts", "simhei.ttf")
    if not os.path.exists(path):
        key = (name or "").lower()
        path = None
        for k, p in _MAC_FONTS.items():
            if k in key and os.path.exists(p):
                path = p; break
        if not path:
            path = _MAC_DEFAULT if os.path.exists(_MAC_DEFAULT) else None
    try:
        return ImageFont.truetype(path, size_px)
    except Exception:
        return ImageFont.load_default()

def _font_key(name):
    n = (name or "").lower()
    for k in ("msyh", "simhei", "simkai", "simsun", "simfang", "kai", "hei", "song", "fang"):
        if k in n:
            return k
    return "simhei"

def emupx(v):
    return 0 if v is None else int(round(v / EMU_PER_IN * DPI))

def _est_h(paras, box_w_px):
    from PIL import Image, ImageDraw
    d = ImageDraw.Draw(Image.new("RGB", (10, 10)))
    total = 0
    for p in paras:
        runs = p.get('runs', [p])
        if not runs or not "".join(r.get('text', '') for r in runs).strip():
            total += int(20 * SCALE); continue
        size = max(8, int(round((runs[0].get('size', 18)) * SCALE)))
        bold = bool(runs[0].get('bold', False))
        fnt = _f(_font_key(runs[0].get('font', HEI)), size, bold)
        ptext = "".join(r.get('text', '') for r in runs)
        sb = int((p.get('space_before', 0)) * SCALE)
        sa = int((p.get('space_after', 0)) * SCALE)
        ls = p.get('line', 1.3)
        mul = ls if isinstance(ls, (int, float)) else 1.2
        lines, cur = [], ""
        for ch in ptext:
            if ch == "\n":
                if cur: lines.append(cur)
                cur = ""
                continue
            trial = cur + ch
            if d.textlength(trial, font=fnt) <= box_w_px - 8 * SCALE or not cur:
                cur = trial
            else:
                lines.append(cur); cur = ch
        if cur: lines.append(cur)
        lh = int(size * mul)
        total += sb + len(lines) * lh + sa
    return total

def add_text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, margin=0.9):
    """自适字号文本框：渲染前按 PIL 度量缩放，确保不溢出（与审计一致）。"""
    box_w_px = emupx(w); box_h_px = emupx(h)
    for p in paras:
        p.setdefault('line', 1.3)
    meas = _est_h(paras, box_w_px)
    if meas > box_h_px * margin and meas > 0:
        k = (box_h_px * margin) / meas
        for p in paras:
            if 'line' in p:
                p['line'] = p['line'] * k
            for r in p.get('runs', [p]):
                r['size'] = max(8, r['size'] * k)
    textbox(slide, x, y, w, h, paras, anchor)


def card(slide, x, y, w, h, color=WHITE, border=FROST, bw=1.6):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.color.rgb = border; c.line.width = Pt(bw)
    c.shadow.inherit = False
    return c

def tagbar(slide, x, y, w, h, color):
    t = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    t.fill.solid(); t.fill.fore_color.rgb = color
    t.line.fill.background(); t.shadow.inherit = False
    return t

def chip(slide, x, y, w, h, text, color, fs=14):
    add_text(slide, x, y, w, h,
             [{'text': text, 'size': fs, 'color': color, 'bold': True,
               'font': HEI, 'align': PP_ALIGN.CENTER, 'line': 1.2}])


# ---------- 数据提取（仅取学生视角的字段，清洗教师口吻） ----------
def _split(s, sep='\n'):
    return [x.strip() for x in s.split(sep) if x.strip()]

def _clean_bb(bb):
    out = []
    for ln in bb.split('\n'):
        ln = (ln.replace('┌', '').replace('┐', '').replace('└', '').replace('┘', '')
                 .replace('│', '').replace('─', '').replace('├', '').replace('┤', '')
                 .replace('┬', '').replace('┴', '').replace('┼', ''))
        ln = re.sub(r'\s+', ' ', ln).strip()
        if ln:
            out.append(ln)
    return out

def _parse_diff(d):
    m = re.split(r'。原因', d, maxsplit=1)
    head = m[0].strip()
    body = ('原因' + m[1]).strip() if len(m) > 1 else ''
    return head, body

def _parse_ex(ex):
    ex0 = re.split(r'【参考答案', ex)[0]
    parts = re.split(r'【提高作业】|【提高】', ex0)
    base = parts[0].replace('【基础作业】', '').strip()
    improve = parts[1].strip() if len(parts) > 1 else ''
    return base, improve

def _first_sentence(s):
    s = s.strip()
    for sep in ('。', '；'):
        i = s.find(sep)
        if i > 0:
            return s[:i] + '。'
    return s


# ---------- 主渲染 ----------
def render_bio(json_path, sources, out_path, photo_decision='color', note=''):
    with open(json_path, encoding='utf-8') as f:
        d = json.load(f)
    cid = d['id']
    title = d['title']
    book = d.get('book', '')
    unit = d.get('unitTitle', '')
    subject = d.get('subject', '生物')
    objectives = d.get('objectives', [])
    keyp = _split(d.get('keyPoints', ''))
    diff = _split(d.get('difficulties', ''))
    methods = _split(d.get('teachingMethods', ''))
    tb = d.get('textbookAnalysis', '')
    ex = d.get('exercises', '')
    bb = _clean_bb(d.get('blackboard', ''))
    refl = d.get('reflection', '')
    nxt = ''
    m = re.search(r'衔接[:：]\s*([^\n。]+)', refl)
    if m:
        nxt = m.group(1).strip()
    base, improve = _parse_ex(ex)

    prs, BLANK = new_presentation()
    ACC = [FROST, XIANG, GOLD, MUTED]

    # ---------- P1 封面 ----------
    def s_cover(s):
        bg(s, INK)
        rule(s, M, M + Inches(0.35), Inches(0.9), GOLD, 3)
        chip(s, M, M + Inches(0.55), Inches(11), Inches(0.4),
             f'{book} · {unit}', GOLD, 15)
        add_text(s, M, Inches(1.7), CW, Inches(1.7),
                 [{'text': title, 'size': 40, 'color': WHITE, 'bold': True,
                   'font': KAI, 'line': 1.15}])
        add_text(s, M, Inches(3.5), CW, Inches(1.0),
                 [{'text': _first_sentence(tb), 'size': 16, 'color': SOFT,
                   'font': KAI, 'line': 1.45, 'space_after': 0}])
        # 底部色带
        band = tagbar(s, 0, H - Inches(0.95), W, Inches(0.95), XIANG)
        add_text(s, M, H - Inches(0.95), CW, Inches(0.95),
                 [{'text': f'高中生物 · 学生版课堂 · 第 {period(d)} 课时',
                   'size': 14, 'color': WHITE, 'font': HEI, 'line': 1.3}],
                 anchor=MSO_ANCHOR.MIDDLE)
        page_num(s)

    # ---------- P2 学习目标 ----------
    def s_objectives(s):
        bg(s, PAPER)
        kicker(s, '本课目标', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '四向学习目标', 'size': 30, 'color': INK, 'bold': True, 'font': KAI}])
        objs = objectives[:4] if objectives else ['（目标见教材）']
        n = len(objs)
        cw = (CW - Inches(0.4) * (n - 1)) / n
        y = M + Inches(1.7)
        for i, o in enumerate(objs):
            x = M + i * (cw + Inches(0.4))
            card(s, x, y, cw, Inches(4.2), WHITE, ACC[i % 4], 1.6)
            nb = s.shapes.add_shape(MSO_SHAPE.OVAL, x + cw/2 - Inches(0.4), y + Inches(0.3),
                                    Inches(0.8), Inches(0.8))
            nb.fill.solid(); nb.fill.fore_color.rgb = ACC[i % 4]; nb.line.fill.background(); nb.shadow.inherit = False
            chip(s, x + cw/2 - Inches(0.4), y + Inches(0.42), Inches(0.8), Inches(0.5),
                 str(i + 1), WHITE, 22)
            add_text(s, x + Inches(0.2), y + Inches(1.35), cw - Inches(0.4), Inches(2.6),
                     [{'text': o, 'size': 13.5, 'color': INK, 'font': KAI, 'line': 1.5,
                       'align': PP_ALIGN.CENTER}])
        page_num(s)

    # ---------- P3 背景与权威调研 ----------
    def s_background(s):
        bg(s, PAPER)
        kicker(s, '背景与权威调研', M, M, FROST)
        col_w = (CW - Inches(0.5)) / 2
        lx = M
        add_text(s, lx, M + Inches(1.2), col_w, Inches(0.5),
                 [{'text': '教材定位与权威核实', 'size': 18, 'color': FROST, 'bold': True, 'font': KAI}])
        add_text(s, lx, M + Inches(1.8), col_w, Inches(4.6),
                 [{'text': tb, 'size': 14.5, 'color': INK, 'font': KAI, 'line': 1.55, 'space_after': 8},
                  {'text': f'本课属于《{book}》{unit}单元，沿"基因→性状→变异→进化→稳态"主线展开。',
                   'size': 13, 'color': MUTED, 'font': KAI, 'line': 1.5}])
        rx = M + col_w + Inches(0.5)
        sp = tagbar(s, rx, M + Inches(1.2), col_w, Inches(5.1), INK)
        add_text(s, rx + Inches(0.3), M + Inches(1.4), col_w - Inches(0.6), Inches(0.5),
                 [{'text': '权威调研来源（联网核实）', 'size': 16, 'color': GOLD, 'bold': True, 'font': KAI}])
        src_lines = []
        for src in sources:
            src_lines.append({'text': '· ' + src, 'size': 12.5, 'color': WHITE, 'font': HEI,
                              'line': 1.45, 'space_after': 7})
        add_text(s, rx + Inches(0.3), M + Inches(2.0), col_w - Inches(0.6), Inches(4.0), src_lines)
        caption(s, '本课件未使用无关配图，以学科色块呈现', lx, H - Inches(0.55), col_w, color=MUTED)
        page_num(s)

    # ---------- P4 重点 ----------
    def s_keypoints(s):
        bg(s, PAPER)
        kicker(s, '重点', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '本课核心要点', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
        items = keyp[:6] if keyp else ['（重点见教材）']
        n = len(items)
        y0 = M + Inches(1.5)
        y1 = H - Inches(0.7)
        gap = Inches(0.18)
        ch = (y1 - y0 - gap * (n - 1)) / n
        for i, it in enumerate(items):
            y = y0 + i * (ch + gap)
            x = M
            card(s, x, y, CW, ch, WHITE, ACC[i % 4], 1.4)
            bar = tagbar(s, x, y, Inches(0.12), ch, ACC[i % 4])
            nb = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + ch/2 - Inches(0.32),
                                    Inches(0.64), Inches(0.64))
            nb.fill.solid(); nb.fill.fore_color.rgb = ACC[i % 4]; nb.line.fill.background(); nb.shadow.inherit = False
            chip(s, x + Inches(0.3), y + ch/2 - Inches(0.22), Inches(0.64), Inches(0.42),
                 str(i + 1), WHITE, 18)
            add_text(s, x + Inches(1.15), y, CW - Inches(1.4), ch,
                     [{'text': it, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.4}],
                     anchor=MSO_ANCHOR.MIDDLE)
        page_num(s)

    # ---------- P5 方法 ----------
    def s_methods(s):
        bg(s, PAPER)
        kicker(s, '方法', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '怎么学 · 方法指引', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
        ms = methods[:4] if methods else ['（方法见教材）']
        # 2x2
        cols = 2
        rows = (len(ms) + 1) // 2
        cw = (CW - Inches(0.4)) / cols
        rh = (Inches(4.6) - Inches(0.3) * (rows - 1)) / rows
        y0 = M + Inches(1.7)
        for i, mt in enumerate(ms):
            r = i // cols; c = i % cols
            x = M + c * (cw + Inches(0.4))
            y = y0 + r * (rh + Inches(0.3))
            card(s, x, y, cw, rh, WHITE, ACC[i % 4], 1.6)
            tbar = tagbar(s, x, y, cw, Inches(0.55), ACC[i % 4])
            chip(s, x, y + Inches(0.07), cw, Inches(0.42), f'方法 {i+1}', WHITE, 16)
            add_text(s, x + Inches(0.3), y + Inches(0.75), cw - Inches(0.6), rh - Inches(0.9),
                     [{'text': mt, 'size': 14.5, 'color': INK, 'font': KAI, 'line': 1.5,
                       'align': PP_ALIGN.CENTER}])
        page_num(s)

    # ---------- P6 难点 ----------
    def s_difficulties(s):
        bg(s, PAPER)
        kicker(s, '难点', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '易卡之处 · 怎么破', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
        ds = [_parse_diff(x) for x in diff[:3]] if diff else [('（难点见教材）', '')]
        n = len(ds)
        cw = (CW - Inches(0.4) * (n - 1)) / n
        y = M + Inches(1.7)
        for i, (head, body) in enumerate(ds):
            x = M + i * (cw + Inches(0.4))
            card(s, x, y, cw, Inches(4.2), WHITE, ACC[i % 4], 1.6)
            nb = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.25),
                                    Inches(0.6), Inches(0.6))
            nb.fill.solid(); nb.fill.fore_color.rgb = ACC[i % 4]; nb.line.fill.background(); nb.shadow.inherit = False
            chip(s, x + Inches(0.25), y + Inches(0.33), Inches(0.6), Inches(0.42),
                 str(i + 1), WHITE, 18)
            add_text(s, x + Inches(1.0), y + Inches(0.28), cw - Inches(1.2), Inches(0.7),
                     [{'text': head, 'size': 15.5, 'color': INK, 'bold': True, 'font': KAI, 'line': 1.3}])
            add_text(s, x + Inches(0.3), y + Inches(1.05), cw - Inches(0.6), Inches(3.0),
                     [{'text': body, 'size': 13, 'color': INK, 'font': KAI, 'line': 1.5}])
        page_num(s)

    # ---------- P7 板书精华 ----------
    def s_blackboard(s):
        bg(s, PAPER)
        kicker(s, '板书精华', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '一页速记 · 知识骨架', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
        panel = card(s, M, M + Inches(1.5), CW, Inches(4.8), SOFT, MUTED, 1.2)
        lines = bb if bb else ['（板书见教材）']
        paras = []
        for ln in lines:
            paras.append({'text': ln, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.35,
                          'space_after': 5})
        add_text(s, M + Inches(0.4), M + Inches(1.75), CW - Inches(0.8), Inches(4.3), paras)
        page_num(s)

    # ---------- P8 作业 ----------
    def s_exercises(s):
        bg(s, PAPER)
        kicker(s, '作业', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '分层作业', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
        y = M + Inches(1.7)
        if improve:
            cw = (CW - Inches(0.4)) / 2
            card(s, M, y, cw, Inches(4.2), WHITE, FROST, 1.6)
            tbar = tagbar(s, M, y, cw, Inches(0.55), FROST)
            chip(s, M, y + Inches(0.07), cw, Inches(0.42), '基础 · 必做', WHITE, 16)
            add_text(s, M + Inches(0.3), y + Inches(0.75), cw - Inches(0.6), Inches(3.3),
                     [{'text': base, 'size': 13.5, 'color': INK, 'font': KAI, 'line': 1.55}])
            x2 = M + cw + Inches(0.4)
            card(s, x2, y, cw, Inches(4.2), WHITE, XIANG, 1.6)
            tbar2 = tagbar(s, x2, y, cw, Inches(0.55), XIANG)
            chip(s, x2, y + Inches(0.07), cw, Inches(0.42), '提高 · 选做', WHITE, 16)
            add_text(s, x2 + Inches(0.3), y + Inches(0.75), cw - Inches(0.6), Inches(3.3),
                     [{'text': improve, 'size': 13.5, 'color': INK, 'font': KAI, 'line': 1.55}])
        else:
            card(s, M, y, CW, Inches(4.2), WHITE, FROST, 1.6)
            tbar = tagbar(s, M, y, CW, Inches(0.55), FROST)
            chip(s, M, y + Inches(0.07), CW, Inches(0.42), '基础 · 必做', WHITE, 16)
            add_text(s, M + Inches(0.3), y + Inches(0.75), CW - Inches(0.6), Inches(3.3),
                     [{'text': base, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.55}])
        page_num(s)

    # ---------- P9 单元小结 ----------
    def s_summary(s):
        bg(s, PAPER)
        kicker(s, '单元小结', M, M, FROST)
        add_text(s, M, M + Inches(0.75), CW, Inches(0.6),
                 [{'text': '回顾与衔接', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
        takes = keyp[:3] if len(keyp) >= 3 else (keyp if keyp else ['（要点见教材）'])
        y = M + Inches(1.6)
        for i, t in enumerate(takes):
            x = M
            card(s, x, y + i * Inches(0.95), CW, Inches(0.8), WHITE, ACC[i % 4], 1.3)
            bar = tagbar(s, x, y + i * Inches(0.95), Inches(0.12), Inches(0.8), ACC[i % 4])
            add_text(s, x + Inches(0.35), y + i * Inches(0.95), CW - Inches(0.6), Inches(0.8),
                     [{'text': '◆ ' + t, 'size': 13.5, 'color': INK, 'font': KAI, 'line': 1.3}],
                     anchor=MSO_ANCHOR.MIDDLE)
        ny = y + len(takes) * Inches(0.95) + Inches(0.15)
        panel = tagbar(s, M, ny, CW, Inches(1.25), INK)
        add_text(s, M + Inches(0.3), ny + Inches(0.12), CW - Inches(0.6), Inches(1.0),
                 [{'text': ('下节课衔接：' + nxt) if nxt else '学完本课，可结合前后课时串成知识线。',
                   'size': 14, 'color': GOLD, 'bold': True, 'font': KAI, 'line': 1.4, 'space_after': 4},
                  {'text': '本课件未使用无关配图，以学科色块呈现；要点以教材为准。',
                   'size': 11.5, 'color': SOFT, 'font': HEI, 'line': 1.3}])
        page_num(s)

    for fn in [s_cover, s_objectives, s_background, s_keypoints, s_methods,
               s_difficulties, s_blackboard, s_exercises, s_summary]:
        fn(new_slide(prs, BLANK))

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path


def period(d):
    """从 tags / 课时 取课时号，回退到 periodNumber。"""
    for t in d.get('tags', []):
        if '课时' in t:
            return t.replace('第', '').replace('课时', '')
    return str(d.get('periodNumber', ''))


if __name__ == '__main__':
    import sys
    p = sys.argv[1] if len(sys.argv) > 1 else 'preview_v7/_fine_data/l-bio-b2-u4-1.json'
    render_bio(p, ['来源示例'], 'preview_v7/bio/_test.pptx')
    print('OK', p)
