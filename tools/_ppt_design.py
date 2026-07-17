# -*- coding: utf-8 -*-
"""通用设计系统：暖纸杂志风（课堂学生版 PPT 共享模块）。
从沁园春范本 + 百合花修正提取。所有 866 课 renderer 统一调用。
设计令牌：PAPER/INK/FROST/XIANG/MUTED/WHITE/GOLD/SOFT + KAI/HEI。
助手：bg / place_photo / scrim / textbox / rule / kicker / caption /
      card / pill / summary_bar / make_prs / page_num。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as ET

# ---------- design tokens ----------
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
INK   = RGBColor(0x1C, 0x2A, 0x33)
FROST = RGBColor(0xB2, 0x3A, 0x2A)   # 霜红
XIANG = RGBColor(0x2E, 0x7D, 0x6B)   # 湘碧
MUTED = RGBColor(0x6B, 0x62, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xC8, 0xA8, 0x6B)
SOFT  = RGBColor(0xE7, 0xDF, 0xD2)

KAI = '楷体'
HEI = '微软雅黑'

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.7)
CW = Inches(13.333 - 1.4)
CH = Inches(7.5 - 1.4)

# 全局字号缩放系数：render_lesson 在溢出自动修复时临时下调，渲染后还原。
FONT_SCALE = 1.0

# ---------- low-level ----------
def set_ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def make_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=PAPER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def place_photo(slide, path, x, y, w, h):
    """铺满指定框并居中裁剪；包围盒严格等于 (x,y,w,h)，绝不溢出。
    int() 防御：新版 python-pptx 不接受 float 坐标。"""
    if not path or not os.path.exists(path):
        return None
    x, y, w, h = int(x), int(y), int(w), int(h)
    try:
        pic = slide.shapes.add_picture(path, x, y)
    except Exception:
        return None
    iw, ih = pic.width, pic.height
    if not iw or not ih:
        return pic
    img_ratio = iw / ih
    box_ratio = w / h
    pic.width = w
    pic.height = h
    if img_ratio > box_ratio:
        vis = box_ratio / img_ratio
        crop = (1 - vis) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        vis = (h * img_ratio) / w
        crop = (1 - vis) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    pic.left = x; pic.top = y
    return pic

def scrim(slide, x, y, w, h, color, alpha_pct):
    x, y, w, h = int(x), int(y), int(w), int(h)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    srgb = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    alpha = ET.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(int(alpha_pct * 100000)))
    return shp

def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = int(x), int(y), int(w), int(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'space_after' in p: para.space_after = Pt(p['space_after'])
        if 'space_before' in p: para.space_before = Pt(p['space_before'])
        if 'line' in p: para.line_spacing = p['line']
        runs = p['runs'] if 'runs' in p else [p]
        for r in runs:
            run = para.add_run()
            run.text = r['text']
            run.font.size = Pt(r['size'] * FONT_SCALE)
            run.font.bold = r.get('bold', False)
            run.font.name = r.get('font', HEI)
            run.font.color.rgb = r['color']
            set_ea(run, r.get('font', HEI))
    return tb

def rule(slide, x, y, w, color=FROST, thick=2.2):
    x, y, w = int(x), int(y), int(w)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thick / 72.0))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def kicker(slide, text, x=M, y=M, color=FROST):
    textbox(slide, x, y, Inches(8), Inches(0.4),
            [{'text': text, 'size': 13, 'color': color, 'bold': True, 'font': HEI, 'space_after': 0}])
    rule(slide, x, y + Inches(0.5), Inches(0.9), color, 2.4)

def caption(slide, text, x, y, w, color=MUTED):
    textbox(slide, x, y, w, Inches(0.5),
            [{'text': text, 'size': 11, 'color': color, 'font': HEI}])

def card(slide, x, y, w, h, fill=WHITE, line=SOFT, line_w=1, radius=True):
    x, y, w, h = int(x), int(y), int(w), int(h)
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def pill(slide, x, y, w, h, text, fill, txt_color=WHITE, size=14, font=HEI):
    x, y, w, h = int(x), int(y), int(w), int(h)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    textbox(slide, x, y, w, h,
            [{'text': text, 'size': size, 'color': txt_color, 'bold': True,
              'font': font, 'align': PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
    return shp

def summary_bar(slide, text, x=M, color=INK, text_color=GOLD, size=18, bar_h=0.62):
    """底部醒目小结色块（百合花修正范式）：墨蓝底 + 金色居中大字。"""
    y = H - Inches(bar_h + 0.3)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(CW), Inches(bar_h))
    band.fill.solid(); band.fill.fore_color.rgb = color
    band.line.fill.background(); band.shadow.inherit = False
    textbox(slide, int(x), int(y + Inches(0.09)), int(CW), Inches(bar_h - 0.18),
            [{'text': text, 'size': size, 'color': text_color, 'bold': True,
              'font': HEI, 'align': PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
    return band

def page_num(slide, n, dark=False):
    c = WHITE if dark else MUTED
    textbox(slide, W - Inches(1.3), H - Inches(0.55), Inches(0.8), Inches(0.35),
            [{'text': f'{n:02d}', 'size': 11, 'color': c, 'font': HEI,
              'align': PP_ALIGN.RIGHT}])

import os
