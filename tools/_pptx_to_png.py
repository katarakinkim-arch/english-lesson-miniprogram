# -*- coding: utf-8 -*-
"""Generic PPTX -> PNG/PDF previewer.
Reads the REAL pptx (shapes, fills, text, embedded photos) and redraws each
slide with PIL at ~110 dpi so the user can "open" the deck inside chat.
Coordinates/colors/photos are taken straight from the file, so the preview
mirrors the actual PowerPoint layout (wrapping is approximated).
"""
import io, os, sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

# SRC / OUTDIR can be overridden via argv[1] / argv[2]
SRC = sys.argv[1] if len(sys.argv) > 1 else r"C:\Users\1\WorkBuddy\2026-07-08-15-47-48\miniprogram\preview_v7\qinyuanchun_consensus.pptx"
OUTDIR = sys.argv[2] if len(sys.argv) > 2 else r"C:\Users\1\WorkBuddy\2026-07-08-15-47-48\miniprogram\preview_v7\full"
DPI = 110
EMU_PER_IN = 914400
SCALE = DPI / 72.0

FONT_DIR = r"C:\Windows\Fonts"
FONT_MAP = {
    "msyh":   os.path.join(FONT_DIR, "msyh.ttc"),
    "msyhbd": os.path.join(FONT_DIR, "msyhbd.ttc"),
    "simhei": os.path.join(FONT_DIR, "simhei.ttf"),
    "simkai": os.path.join(FONT_DIR, "simkai.ttf"),
    "simsun": os.path.join(FONT_DIR, "simsun.ttc"),
    "simfang":os.path.join(FONT_DIR, "simfang.ttf"),
    "kai":    os.path.join(FONT_DIR, "simkai.ttf"),
    "fang":   os.path.join(FONT_DIR, "simfang.ttf"),
    "song":   os.path.join(FONT_DIR, "simsun.ttc"),
    "hei":    os.path.join(FONT_DIR, "simhei.ttf"),
}
_fc = {}
def font(name, size_px, bold=False):
    key = (name, size_px, bold)
    if key in _fc: return _fc[key]
    path = FONT_MAP.get(name.lower())
    if not path or not os.path.exists(path):
        path = FONT_MAP["simhei"]
    try:
        f = ImageFont.truetype(path, size_px)
    except Exception:
        f = ImageFont.load_default()
    _fc[key] = f
    return f

def emu_px(v):
    if v is None: return 0
    return int(round(v / EMU_PER_IN * DPI))

def rgb_of(color):
    try:
        if color and color.type is not None and hasattr(color, 'rgb'):
            return (color.rgb[0], color.rgb[1], color.rgb[2])
    except Exception:
        pass
    return None

def wrap_text(draw, text, fnt, max_w):
    """Greedy char-by-char wrap. Chinese chars ~ full em width; latin uses measure.
    Embedded '\\n' forces a hard line break."""
    lines, cur = [], ""
    for ch in text:
        if ch == "\n":
            if cur:
                lines.append(cur)
            lines.append("")
            cur = ""
            continue
        trial = cur + ch
        w = draw.textlength(trial, font=fnt)
        if w <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur); cur = ch
    if cur:
        lines.append(cur)
    return lines or [""]

def run_color(run):
    try:
        if run.font.color and run.font.color.type is not None:
            c = run.font.color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return (28, 43, 51)

def run_size(run):
    if run.font.size is not None:
        return max(8, int(round(run.font.size.pt * SCALE)))
    return max(8, int(round(18 * SCALE)))

def run_fontname(run):
    n = (run.font.name or "").lower()
    for k in FONT_MAP:
        if k in n: return k
    return "simhei"

def draw_text_frame(draw, shape, W, H):
    tf = shape.text_frame
    x0 = emu_px(shape.left); y0 = emu_px(shape.top)
    bw = emu_px(shape.width); bh = emu_px(shape.height)
    # anchor
    anchor = tf.vertical_anchor
    y = y0
    line_hs = []
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            y += int(20*SCALE); continue
        size = run_size(runs[0])
        col = run_color(runs[0])
        nm = run_fontname(runs[0])
        bold = bool(runs[0].font.bold)
        fnt = font(nm, size, bold)
        # gather full paragraph text (simplify: use first run style for whole line)
        ptext = "".join(r.text for r in runs)
        space_before = int((para.space_before.pt if para.space_before else 0) * SCALE)
        space_after = int((para.space_after.pt if para.space_after else 0) * SCALE)
        ls = para.line_spacing
        if isinstance(ls, float): line_mul = ls
        elif ls is None: line_mul = 1.2
        else: line_mul = 1.2
        lines = wrap_text(draw, ptext, fnt, bw - int(8*SCALE))
        lh = int(size * line_mul)
        y += space_before
        for i, ln in enumerate(lines):
            tw = draw.textlength(ln, font=fnt)
            if para.alignment == PP_ALIGN.CENTER:
                tx = x0 + (bw - tw)//2
            elif para.alignment == PP_ALIGN.RIGHT:
                tx = x0 + bw - tw - int(4*SCALE)
            else:
                tx = x0 + int(4*SCALE)
            draw.text((tx, y), ln, font=fnt, fill=col)
            y += lh
        y += space_after
    return y

def render_slide(slide, idx, prs_w, prs_h):
    W = emu_px(prs_w); H = emu_px(prs_h)
    img = Image.new("RGB", (W, H), (247, 242, 234))
    d = ImageDraw.Draw(img)
    # background rect (first full-size solid shape) already drawn by shapes loop;
    # but draw a default paper bg first.
    for sh in slide.shapes:
        try:
            st = sh.shape_type
        except Exception:
            st = None
        left = emu_px(sh.left); top = emu_px(sh.top)
        w = emu_px(sh.width); h = emu_px(sh.height)
        if st == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = sh.image.blob
                pic = Image.open(io.BytesIO(blob)).convert("RGB")
                pic = pic.resize((max(1,w), max(1,h)))
                img.paste(pic, (left, top))
            except Exception as e:
                d.rectangle([left, top, left+w, top+h], fill=(220,220,220))
            continue
        # autoshape / rectangle
        fill = None
        alpha = 1.0
        try:
            if sh.fill.type is not None and sh.fill.type == 1:  # solid
                c = sh.fill.fore_color.rgb
                fill = (c[0], c[1], c[2])
                # read a:alpha if present (OOXML: val is in 1/100000 of full)
                try:
                    xf = sh.fill.fore_color._xFill
                    srgb = xf.find(qn('a:srgbClr'))
                    if srgb is not None:
                        ael = srgb.find(qn('a:alpha'))
                        if ael is not None:
                            alpha = min(1.0, max(0.0, int(ael.get('val')) / 100000.0))
                except Exception:
                    pass
        except Exception:
            fill = None
        if fill is not None:
            is_round = False
            try:
                if 'ROUND' in str(sh.auto_shape_type): is_round = True
            except Exception:
                pass
            if alpha < 0.999:
                # semi-transparent: blend over whatever is already on the canvas
                region = img.crop((left, top, left + w, top + h))
                ov = Image.new("RGB", region.size, fill)
                region = Image.blend(region, ov, alpha)
                img.paste(region, (left, top))
            elif is_round:
                r = min(w, h)//6
                d.rounded_rectangle([left, top, left+w, top+h], radius=r, fill=fill)
            else:
                d.rectangle([left, top, left+w, top+h], fill=fill)
        # text
        if sh.has_text_frame and sh.text_frame.text.strip():
            draw_text_frame(d, sh, W, H)
    return img

def main():
    os.makedirs(OUTDIR, exist_ok=True)
    prs = Presentation(SRC)
    pw = prs.slide_width; ph = prs.slide_height
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        im = render_slide(slide, i, pw, ph)
        p = os.path.join(OUTDIR, f"p{i:02d}.png")
        im.save(p, "PNG")
        pages.append(im)
        print(f"  slide {i}: {im.size}")
    # build PDF
    _base = os.path.basename(OUTDIR.rstrip(os.sep)) or "preview"
    pdf_path = os.path.abspath(os.path.join(OUTDIR, "..", _base + "_preview.pdf"))
    if pages:
        first = pages[0].convert("RGB")
        rest = [p.convert("RGB") for p in pages[1:]]
        first.save(pdf_path, "PDF", save_all=True, append_images=rest)
        print("PDF ->", pdf_path)
    print("TOTAL", len(pages))

if __name__ == "__main__":
    main()
