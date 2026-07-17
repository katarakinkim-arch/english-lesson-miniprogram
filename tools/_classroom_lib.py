# -*- coding: utf-8 -*-
"""课堂版 PPT 共享库 — 设计系统 + 通用 helper（magazine 风）。

所有手工精排课共用，避免每课重复基础设施。
设计系统：暖纸 PAPER / 墨蓝 INK / 霜红 FROST / 湘碧 XIANG / 金 GOLD
字体：楷体 KAI（标题/引文）+ 微软雅黑 HEI（正文/标签）
画布：16:9 (13.333in x 7.5in)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as ET

# ---------- design system ----------
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
INK   = RGBColor(0x1C, 0x2A, 0x33)
FROST = RGBColor(0xB2, 0x3A, 0x2A)
XIANG = RGBColor(0x2E, 0x7D, 0x6B)
MUTED = RGBColor(0x6B, 0x62, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xC8, 0xA8, 0x6B)
SOFT  = RGBColor(0xE7, 0xDF, 0xD2)

KAI = '楷体'
HEI = '微软雅黑'

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.7)
CW = Inches(13.333 - 1.4)

def new_presentation():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs, prs.slide_layouts[6]

# ---------- helpers ----------
def set_ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def bg(slide, color=PAPER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def place_photo(slide, path, x, y, w, h):
    x = int(x); y = int(y); w = int(w); h = int(h)
    pic = slide.shapes.add_picture(path, x, y)
    iw, ih = pic.width, pic.height
    ir = iw / ih
    br = w / h
    pic.width = w
    pic.height = h
    if ir > br:
        vis = br / ir
        crop = (1 - vis) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        vis = (h * ir) / w
        crop = (1 - vis) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    pic.left = x
    pic.top = y
    return pic

def scrim(slide, x, y, w, h, color, alpha_pct):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    srgb = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    alpha = ET.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(int(alpha_pct * 100000)))
    return shp

def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'space_after' in p:
            para.space_after = Pt(p['space_after'])
        if 'space_before' in p:
            para.space_before = Pt(p['space_before'])
        if 'line' in p:
            para.line_spacing = p['line']
        runs = p['runs'] if 'runs' in p else [p]
        for r in runs:
            run = para.add_run()
            run.text = r['text']
            run.font.size = Pt(r['size'])
            run.font.bold = r.get('bold', False)
            run.font.name = r.get('font', HEI)
            run.font.color.rgb = r['color']
            set_ea(run, r.get('font', HEI))
    return tb

def rule(slide, x, y, w, color=FROST, thick=2.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thick / 72.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def kicker(slide, text, x=M, y=M, color=FROST):
    textbox(slide, x, y, Inches(7), Inches(0.4),
            [{'text': text, 'size': 13, 'color': color, 'bold': True, 'font': HEI, 'space_after': 0}])
    rule(slide, x, y + Inches(0.5), Inches(0.9), color, 2.4)

def new_slide(prs, BLANK):
    return prs.slides.add_slide(BLANK)

def page_num(slide, dark=False):
    PAGE[0] += 1
    c = WHITE if dark else MUTED
    textbox(slide, W - Inches(1.3), H - Inches(0.55), Inches(0.8), Inches(0.35),
            [{'text': f'{PAGE[0]:02d}', 'size': 11, 'color': c, 'font': HEI, 'align': PP_ALIGN.RIGHT}])

PAGE = [0]

def caption(slide, text, x, y, w, color=MUTED):
    textbox(slide, x, y, w, Inches(0.4),
            [{'text': text, 'size': 11, 'color': color, 'font': HEI, 'align': PP_ALIGN.CENTER}])

def quote_block(slide, x, y, w, text, source, color=FROST):
    """引用块：左侧竖条 + 楷体引文 + 来源。"""
    rule(slide, x, y, Inches(0.06), color, 30)
    textbox(slide, x + Inches(0.25), y, w - Inches(0.25), Inches(2.2),
            [{'text': text, 'size': 18, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 8},
             {'text': '—— ' + source, 'size': 12, 'color': MUTED, 'font': HEI}])

def step_card(slide, x, y, w, h, num, title, body_lines, accent=FROST):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(1.8)
    card.shadow.inherit = False
    nb = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.5), Inches(0.5))
    nb.fill.solid()
    nb.fill.fore_color.rgb = accent
    nb.line.fill.background()
    nb.shadow.inherit = False
    # badge sized for 1-2 chars; long labels (e.g. 4-digit year) shrink to fit one line
    nb_size = 16 if len(str(num)) <= 2 else 12
    textbox(slide, x + Inches(0.18), y + Inches(0.24), Inches(0.5), Inches(0.42),
            [{'text': str(num), 'size': nb_size, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    textbox(slide, x + Inches(0.8), y + Inches(0.18), w - Inches(1.0), Inches(0.5),
            [{'text': title, 'size': 17, 'color': INK, 'bold': True, 'font': HEI}])
    bl = [{'text': t, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.4, 'space_after': 4} for t in body_lines]
    textbox(slide, x + Inches(0.8), y + Inches(0.72), w - Inches(1.0), h - Inches(0.85), bl)
