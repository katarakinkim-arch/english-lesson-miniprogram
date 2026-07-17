# -*- coding: utf-8 -*-
"""数据驱动课堂 PPT 渲染器 v2 — 语文 128 课，统一杂志风设计系统
（与《沁园春·长沙》锁定标准一致的配色/字体/网格/kicker/rule + 全幅封面大图 + 导览页 + 章节分隔）
读取 tools/_lessons_cn.json，为每课生成独立 .pptx 到 preview_cn/<id>.pptx。
内容 = 每课自己的教材/教案数据（不做视频融合）；皮 = 沁园春打磨度；每课保持自己的结构。
封面大图用 PIL 手绘（渐变+淡楷体水印），无水印、不联网、不研究。
"""
import os, json, re, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as ET
from PIL import Image, ImageDraw, ImageFont

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(BASE, 'tools', '_lessons_cn.json')
OUT_DIR = os.path.join(BASE, 'preview_cn')
COVER_DIR = os.path.join(OUT_DIR, '_covers')
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(COVER_DIR, exist_ok=True)

FONT_DIR = r'C:\Windows\Fonts'

# ---------- design system（锁定标准）----------
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
INK   = RGBColor(0x1C, 0x2A, 0x33)
FROST = RGBColor(0xB2, 0x3A, 0x2A)
XIANG = RGBColor(0x2E, 0x7D, 0x6B)
MUTED = RGBColor(0x6B, 0x62, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xC8, 0xA8, 0x6B)
SOFT  = RGBColor(0xE7, 0xDF, 0xD2)
BLUE  = RGBColor(0x2E, 0x5B, 0x8A)

KAI = '楷体'
HEI = '微软雅黑'

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.7)
CW = Inches(13.333 - 1.4)

prs = None
PAGE = [0]

# ---------- helpers ----------
def set_ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def bg(slide, color=PAPER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def place_photo(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(path, x, y)
    iw, ih = pic.width, pic.height
    img_ratio = iw / ih; box_ratio = w / h
    pic.width = w; pic.height = h
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - (h * img_ratio) / w) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    pic.left = x; pic.top = y
    return pic

def scrim(slide, x, y, w, h, color, alpha_pct):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    srgb = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    alpha = ET.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(int(alpha_pct * 100000)))
    return shp

def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'space_after' in p: para.space_after = Pt(p['space_after'])
        if 'space_before' in p: para.space_before = Pt(p['space_before'])
        if 'line' in p: para.line_spacing = p['line']
        runs = p['runs'] if 'runs' in p else [p]
        for r in runs:
            run = para.add_run(); run.text = r['text']
            run.font.size = Pt(r['size']); run.font.bold = r.get('bold', False)
            run.font.name = r.get('font', HEI); run.font.color.rgb = r['color']
            set_ea(run, r.get('font', HEI))
    return tb

def rule(slide, x, y, w, color=FROST, thick=2.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thick / 72.0))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def kicker(slide, text, x=M, y=M, color=FROST):
    textbox(slide, x, y, Inches(7), Inches(0.4),
            [{'text': text, 'size': 13, 'color': color, 'bold': True, 'font': HEI}])
    rule(slide, x, y + Inches(0.5), Inches(0.9), color, 2.4)

def new_slide():
    global prs
    return prs.slides.add_slide(prs.slide_layouts[6])

def page_num(slide, dark=False):
    PAGE[0] += 1
    c = WHITE if dark else MUTED
    textbox(slide, W - Inches(1.3), H - Inches(0.55), Inches(0.8), Inches(0.35),
            [{'text': f'{PAGE[0]:02d}', 'size': 11, 'color': c, 'font': HEI,
              'align': PP_ALIGN.RIGHT}])

def mk_paras(text, size=16, color=INK, font=HEI, line=1.6, space_after=8, bold=False):
    out = []
    for ln in (text or '').split('\n'):
        ln = ln.strip()
        if not ln: continue
        out.append({'text': ln, 'size': size, 'color': color, 'font': font,
                    'line': line, 'space_after': space_after, 'bold': bold})
    return out

def card(slide, x, y, w, h, fill=SOFT, line_color=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

# ---------- PIL 封面大图（手绘，无联网/无研究）----------
def _load_font(name, size):
    try:
        return ImageFont.truetype(os.path.join(FONT_DIR, name), size)
    except OSError:
        return ImageFont.load_default()

def _theme_bottom(lesson):
    tg = lesson.get('taskGroup', '') or ''
    if '文学' in tg or '写作' in tg: return (0x7A, 0x26, 0x1E)   # 深霜红
    if '思辨' in tg or '议论' in tg: return (0x8A, 0x6E, 0x2E)   # 深金
    if '实用' in tg or '新闻' in tg: return (0x1E, 0x3F, 0x5E)   # 深蓝
    if '整本书' in tg: return (0x1E, 0x52, 0x46)                # 深湘碧
    return (0x6E, 0x24, 0x1C)

def gen_cover_bg(lesson):
    tid = lesson.get('id') or 'x'
    out = os.path.join(COVER_DIR, tid + '.jpg')
    if os.path.exists(out):
        return out
    Wp, Hp = 1600, 900
    top = (0x1C, 0x2A, 0x33)
    bot = _theme_bottom(lesson)
    img = Image.new('RGB', (Wp, Hp))
    draw = ImageDraw.Draw(img)
    for y in range(Hp):
        t = y / Hp
        r = int(top[0] + (bot[0] - top[0]) * t)
        g = int(top[1] + (bot[1] - top[1]) * t)
        b = int(top[2] + (bot[2] - top[2]) * t)
        draw.line([(0, y), (Wp, y)], fill=(r, g, b))
    # 淡竖向纹理线（杂志风）
    for i in range(0, Wp, 28):
        draw.line([(i, 0), (i, Hp)], fill=(0xDD, 0xD5, 0xCA), width=1)
    # 底部淡楷体水印（标题，低对比）
    title = lesson.get('title', '') or ''
    if title:
        size = min(230, int((Wp - 120) / max(len(title), 1) * 1.02))
        size = max(120, size)
        fbig = _load_font('simkai.ttf', size)
        bbox = draw.textbbox((0, 0), title, font=fbig)
        tw = bbox[2] - bbox[0]
        tx = Wp - tw - 70
        ty = Hp - size - 90
        draw.text((tx, ty), title, fill=(0x49, 0x53, 0x59), font=fbig)
    # 淡金细线装饰（左下）
    draw.line([(70, Hp - 60), (Wp - 70, Hp - 60)], fill=(0xC8, 0xA8, 0x6B), width=2)
    img.save(out, 'JPEG', quality=86)
    return out

# ---------- content parsers ----------
def parse_step(content):
    parts = {'stage': '', 'teacher': '', 'preset': '', 'board': '', 'diff': '', 'pitfall': ''}
    m = re.search(r'【([^】]+)】', content or '')
    if m: parts['stage'] = m.group(1)
    rest = re.sub(r'【[^】]*】', '', content or '')
    segs = re.split(r'(?=教师：|预设回答：|板书时机：|差异化提示：|易错点提醒：)', rest)
    for seg in segs:
        seg = seg.strip()
        if not seg: continue
        if seg.startswith('教师：'): parts['teacher'] = seg[3:].strip()
        elif seg.startswith('预设回答：'): parts['preset'] = seg[4:].strip()
        elif seg.startswith('板书时机：'): parts['board'] = seg[4:].strip()
        elif seg.startswith('差异化提示：'): parts['diff'] = seg[5:].strip()
        elif seg.startswith('易错点提醒：'): parts['pitfall'] = seg[5:].strip()
    return parts

def strip_prefix(text, prefix):
    if text and text.startswith(prefix):
        return text[len(prefix):].strip()
    return text

def split_marker(text, markers):
    out = {mk: '' for mk in markers}
    if not text: return out
    pat = '(' + '|'.join(re.escape(mk) for mk in markers) + ')'
    segs = re.split(pat, text)
    for i in range(1, len(segs) - 1, 2):
        mk = segs[i]; body = segs[i + 1].strip() if i + 1 < len(segs) else ''
        out[mk] = body
    return out

def obj_color(label):
    d = {'语言能力': FROST, '文化意识': XIANG, '思维品质': GOLD, '学习能力': BLUE}
    for k, v in d.items():
        if label.startswith(k): return v
    return INK

def first_line(text, n=46):
    text = (text or '').strip()
    if not text: return '—'
    line = text.split('\n')[0].strip()
    return line[:n] + ('…' if len(line) > n else '')

# ---------- slide builders ----------
def s_cover(lesson, bg_path):
    s = new_slide()
    place_photo(s, bg_path, 0, 0, W, H)
    scrim(s, 0, H - Inches(3.5), W, Inches(3.5), INK, 0.66)
    scrim(s, 0, 0, W, Inches(1.15), INK, 0.32)
    kicker(s, f"{lesson.get('book','')} · 第{lesson.get('unitNumber','')}单元", M, M, GOLD)
    title = lesson.get('title', '')
    textbox(s, M, H - Inches(3.0), CW, Inches(1.7),
            [{'text': title, 'size': 58, 'color': WHITE, 'bold': True, 'font': KAI, 'line': 1.1}])
    sub = []
    if lesson.get('author'):
        sub.append({'text': lesson.get('author'), 'size': 22, 'color': WHITE, 'font': HEI, 'bold': True})
    ltn = lesson.get('lessonTypeName', '') or lesson.get('lessonType', '')
    if ltn:
        sub.append({'text': ('    ' if lesson.get('author') else '') + ltn,
                    'size': 16, 'color': SOFT, 'font': HEI})
    if sub:
        textbox(s, M, H - Inches(1.25), CW, Inches(0.6), [{'runs': sub}])
    foot = f"{lesson.get('taskGroup','')} · 第{lesson.get('periodNumber','')}课时 · {lesson.get('duration','')}分钟"
    textbox(s, M, H - Inches(0.72), CW, Inches(0.4),
            [{'text': foot, 'size': 13, 'color': GOLD, 'font': HEI, 'bold': True}])
    page_num(s, dark=True)

def s_contents(lesson):
    s = new_slide(); bg(s)
    kicker(s, 'CONTENTS · 本课导览')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': f"{lesson.get('title','')} · 课堂导览", 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    process = lesson.get('process', []) or []
    items = []
    if lesson.get('textbookAnalysis'): items.append(('壹', '教材分析', first_line(lesson.get('textbookAnalysis'))))
    if lesson.get('overview'): items.append(('贰', '学情分析', first_line(lesson.get('overview'))))
    if lesson.get('objectives'): items.append(('叁', '教学目标', '核心素养 · 四维度目标'))
    if lesson.get('keyPoints') or lesson.get('difficulties'): items.append(('肆', '教学重难点', '重点突破 与 难点化解'))
    if process: items.append(('伍', '教学过程', f'{len(process)} 个教学环节 · 逐步推进'))
    if lesson.get('blackboard'): items.append(('陆', '板书设计', '主板书结构'))
    if lesson.get('exercises'): items.append(('柒', '作业设计', '基础 · 提高 · 参考答案'))
    if lesson.get('reflection'): items.append(('捌', '教学反思', '亮点 · 改进 · 衔接'))
    y = M + Inches(1.75)
    rowh = Inches(0.62)
    for num, t, sub in items:
        textbox(s, M, y, Inches(1.0), rowh,
                [{'text': num, 'size': 27, 'color': FROST, 'bold': True, 'font': KAI}])
        textbox(s, M + Inches(1.1), y, Inches(4.5), rowh,
                [{'text': t, 'size': 20, 'color': INK, 'bold': True, 'font': HEI, 'space_after': 0}])
        textbox(s, M + Inches(5.7), y + Inches(0.05), Inches(6.0), rowh,
                [{'text': sub, 'size': 13, 'color': MUTED, 'font': HEI}])
        rule(s, M, y + rowh - Inches(0.03), CW, SOFT, 1.0)
        y = y + rowh + Inches(0.02)
    page_num(s)

def s_divider(num, title, sub, accent=FROST):
    s = new_slide(); bg(s, INK)
    # 大号淡色序号水印
    textbox(s, W - Inches(3.6), M - Inches(0.1), Inches(3.2), Inches(3.2),
            [{'text': num, 'size': 200, 'color': RGBColor(0x2E, 0x3A, 0x44), 'bold': True, 'font': KAI, 'align': PP_ALIGN.RIGHT}])
    kicker(s, sub, M, M, GOLD)
    textbox(s, M, Inches(3.0), CW, Inches(1.2),
            [{'text': title, 'size': 44, 'color': WHITE, 'bold': True, 'font': KAI, 'line': 1.1}])
    rule(s, M, Inches(4.35), Inches(1.4), accent, 3.0)
    page_num(s, dark=True)

def s_text_block(lesson, kicker_txt, title_txt, body_text, color=INK, src_key=None):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': title_txt, 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])
    textbox(s, M, M + Inches(1.65), CW, Inches(4.7),
            mk_paras(strip_prefix(body_text, '【学情分析】') if src_key == 'overview' else body_text,
                     size=16, color=color, line=1.6, space_after=9))
    page_num(s)

def s_objectives(lesson):
    s = new_slide(); bg(s)
    kicker(s, '教学目标 · OBJECTIVES')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': '教学目标', 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])
    objs = lesson.get('objectives', []) or []
    y = M + Inches(1.55)
    ch = Inches(1.12); gap = Inches(0.14)
    for obj in objs:
        idx = obj.find('：')
        label = obj[:idx] if idx > 0 else '目标'
        body = obj[idx + 1:] if idx > 0 else obj
        col = obj_color(label)
        card(s, M, y, CW, ch, fill=SOFT)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, y, Inches(0.12), ch)
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
        textbox(s, M + Inches(0.32), y + Inches(0.14), CW - Inches(0.5), ch - Inches(0.28),
                [{'runs': [
                    {'text': label, 'size': 15, 'color': col, 'bold': True, 'font': HEI},
                    {'text': '   ' + body, 'size': 13.5, 'color': INK, 'font': HEI, 'line': 1.35},
                ]}])
        y = y + ch + gap
    page_num(s)

def s_keypoints(lesson):
    s = new_slide(); bg(s)
    kicker(s, '教学重难点 · FOCUS')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': '教学重点 与 难点', 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])
    colw = (CW - Inches(0.5)) / 2
    hdr1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, M + Inches(1.55), colw, Inches(0.55))
    hdr1.fill.solid(); hdr1.fill.fore_color.rgb = FROST; hdr1.line.fill.background()
    textbox(s, M + Inches(0.2), M + Inches(1.62), colw - Inches(0.4), Inches(0.42),
            [{'text': '教学重点', 'size': 17, 'color': WHITE, 'bold': True, 'font': HEI}])
    textbox(s, M + Inches(0.05), M + Inches(2.25), colw - Inches(0.15), Inches(4.0),
            mk_paras(lesson.get('keyPoints', ''), size=14.5, color=INK, line=1.5, space_after=9))
    rx = M + colw + Inches(0.5)
    hdr2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, M + Inches(1.55), colw, Inches(0.55))
    hdr2.fill.solid(); hdr2.fill.fore_color.rgb = XIANG; hdr2.line.fill.background()
    textbox(s, rx + Inches(0.2), M + Inches(1.62), colw - Inches(0.4), Inches(0.42),
            [{'text': '教学难点', 'size': 17, 'color': WHITE, 'bold': True, 'font': HEI}])
    textbox(s, rx + Inches(0.05), M + Inches(2.25), colw - Inches(0.15), Inches(4.0),
            mk_paras(lesson.get('difficulties', ''), size=14.5, color=INK, line=1.5, space_after=9))
    page_num(s)

def s_methods(lesson):
    s = new_slide(); bg(s)
    kicker(s, '教学方法 · METHODS')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': '教学方法', 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])
    textbox(s, M, M + Inches(1.6), CW, Inches(4.8),
            mk_paras(lesson.get('teachingMethods', ''), size=15.5, color=INK, line=1.6, space_after=10))
    page_num(s)

def s_step(lesson, step, idx, total):
    s = new_slide(); bg(s)
    kicker(s, f"教学过程 · STEP {idx}/{total}")
    textbox(s, M, M + Inches(0.55), CW - Inches(2.0), Inches(0.7),
            [{'text': step.get('step', ''), 'size': 22, 'color': INK, 'bold': True, 'font': HEI, 'line': 1.1}])
    tb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, W - M - Inches(1.6), M + Inches(0.55), Inches(1.5), Inches(0.55))
    tb.fill.solid(); tb.fill.fore_color.rgb = XIANG; tb.line.fill.background()
    textbox(s, W - M - Inches(1.6), M + Inches(0.66), Inches(1.5), Inches(0.4),
            [{'text': step.get('time', ''), 'size': 15, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    parts = parse_step(step.get('content', ''))
    cy = M + Inches(1.4); chh = Inches(2.35)
    card(s, M, cy, CW, chh, fill=SOFT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, cy, Inches(0.12), chh)
    bar.fill.solid(); bar.fill.fore_color.rgb = FROST; bar.line.fill.background()
    textbox(s, M + Inches(0.3), cy + Inches(0.14), Inches(2.0), Inches(0.4),
            [{'text': '教师台词', 'size': 14, 'color': FROST, 'bold': True, 'font': HEI}])
    textbox(s, M + Inches(0.3), cy + Inches(0.58), CW - Inches(0.55), chh - Inches(0.75),
            [{'text': parts['teacher'] or '—', 'size': 15, 'color': INK, 'font': HEI, 'line': 1.5}])
    py = cy + chh + Inches(0.18); ph = Inches(1.15)
    card(s, M, py, CW, ph, fill=WHITE, line_color=XIANG, line_w=1.5)
    textbox(s, M + Inches(0.25), py + Inches(0.1), Inches(2.0), Inches(0.35),
            [{'text': '预设回答', 'size': 13, 'color': XIANG, 'bold': True, 'font': HEI}])
    textbox(s, M + Inches(0.25), py + Inches(0.45), CW - Inches(0.5), ph - Inches(0.6),
            [{'text': parts['preset'] or '—', 'size': 13, 'color': INK, 'font': HEI, 'line': 1.35}])
    ny = H - M - Inches(0.85)
    nw = (CW - Inches(0.6)) / 3
    notes = [('板书时机', parts['board'], FROST),
             ('差异化提示', parts['diff'], XIANG),
             ('易错点提醒', parts['pitfall'], GOLD)]
    for i, (lab, val, col) in enumerate(notes):
        nx = M + i * (nw + Inches(0.3))
        textbox(s, nx, ny, nw, Inches(0.8),
                [{'text': lab, 'size': 12, 'color': col, 'bold': True, 'font': HEI, 'space_after': 2},
                 {'text': val or '—', 'size': 11.5, 'color': MUTED, 'font': HEI, 'line': 1.25}])
    page_num(s)

def s_blackboard(lesson):
    s = new_slide(); bg(s, INK)
    kicker(s, '板书设计 · BLACKBOARD', color=GOLD)
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '板书设计', 'size': 26, 'color': WHITE, 'bold': True, 'font': HEI}])
    card(s, M, M + Inches(1.55), CW, Inches(4.7), fill=INK, line_color=GOLD, line_w=1.5)
    textbox(s, M + Inches(0.35), M + Inches(1.8), CW - Inches(0.7), Inches(4.2),
            [{'text': lesson.get('blackboard', ''), 'size': 12.5, 'color': WHITE, 'font': HEI, 'line': 1.15}])
    page_num(s, dark=True)

def s_exercises(lesson):
    s = new_slide(); bg(s)
    kicker(s, '作业设计 · ASSIGNMENT', color=XIANG)
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '作业与练习', 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])
    secs = split_marker(lesson.get('exercises', ''),
                        ['【基础作业】', '【提高作业】', '【参考答案——教师用】'])
    cols = [('【基础作业】', FROST), ('【提高作业】', XIANG), ('【参考答案——教师用】', MUTED)]
    colw = (CW - Inches(0.6)) / 3
    for i, (mk, col) in enumerate(cols):
        x = M + i * (colw + Inches(0.3))
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, M + Inches(1.55), colw, Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col; hdr.line.fill.background()
        label = mk.replace('【', '').replace('】', '')
        textbox(s, x + Inches(0.12), M + Inches(1.6), colw - Inches(0.2), Inches(0.4),
                [{'text': label, 'size': 13, 'color': WHITE, 'bold': True, 'font': HEI}])
        textbox(s, x + Inches(0.05), M + Inches(2.2), colw - Inches(0.1), Inches(4.0),
                mk_paras(secs.get(mk, ''), size=12.5, color=INK, line=1.45, space_after=7))
    page_num(s)

def s_reflection(lesson):
    s = new_slide(); bg(s)
    kicker(s, '教学反思 · REFLECTION')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '教学反思', 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])
    t = lesson.get('reflection', '') or ''
    for e in ['✅', '⚠️', '📌']:
        t = t.replace(e, '')
    parts = {'highlight': '', 'improve': '', 'next': ''}
    m1 = re.search(r'亮点：(.*?)(?=需改进：|$)', t, re.S)
    m2 = re.search(r'需改进：(.*?)(?=下节课衔接：|$)', t, re.S)
    m3 = re.search(r'下节课衔接：(.*)$', t, re.S)
    if m1: parts['highlight'] = m1.group(1).strip()
    if m2: parts['improve'] = m2.group(1).strip()
    if m3: parts['next'] = m3.group(1).strip()
    cards = [('亮点', parts['highlight'], FROST),
             ('需改进', parts['improve'], GOLD),
             ('下节课衔接', parts['next'], XIANG)]
    colw = (CW - Inches(0.6)) / 3
    for i, (lab, val, col) in enumerate(cards):
        x = M + i * (colw + Inches(0.3))
        y = M + Inches(1.6)
        ch = Inches(4.3)
        card(s, x, y, colw, ch, fill=SOFT)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, colw, Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
        textbox(s, x + Inches(0.15), y + Inches(0.12), colw - Inches(0.3), Inches(0.4),
                [{'text': lab, 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI}])
        textbox(s, x + Inches(0.18), y + Inches(0.72), colw - Inches(0.36), ch - Inches(0.9),
                [{'text': val or '—', 'size': 13.5, 'color': INK, 'font': HEI, 'line': 1.5}])
    page_num(s)

# ---------- build one lesson ----------
def build_lesson(lesson):
    global prs, PAGE
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    PAGE[0] = 0
    bg_p = gen_cover_bg(lesson)
    s_cover(lesson, bg_p)
    s_contents(lesson)
    s_text_block(lesson, '教材分析 · TEXTBOOK', '教材分析',
                 lesson.get('textbookAnalysis', ''), color=INK)
    s_text_block(lesson, '学情分析 · STUDENTS', '学情分析',
                 lesson.get('overview', ''), color=XIANG, src_key='overview')
    s_objectives(lesson)
    s_keypoints(lesson)
    s_methods(lesson)
    process = lesson.get('process', []) or []
    total = len(process)
    if process:
        s_divider('伍', '教学过程', 'TEACHING PROCESS · 课堂推进', FROST)
        for i, step in enumerate(process, 1):
            s_step(lesson, step, i, total)
    s_blackboard(lesson)
    s_divider('陆', '巩固与反思', 'ASSIGNMENT & REFLECTION', XIANG)
    s_exercises(lesson)
    s_reflection(lesson)
    out = os.path.join(OUT_DIR, f"{lesson.get('id','lesson')}.pptx")
    prs.save(out)
    return out

# ---------- main ----------
def main():
    with open(DATA, encoding='utf-8') as f:
        lessons = json.load(f)
    n = len(lessons)
    for i, lesson in enumerate(lessons, 1):
        try:
            out = build_lesson(lesson)
            if i % 10 == 0 or i == n:
                print(f"  [{i}/{n}] {lesson.get('id')} -> {os.path.basename(out)}")
        except Exception as e:
            print(f"  [ERR {i}] {lesson.get('id')}: {e}")
    print(f"DONE: {n} lessons -> {OUT_DIR}")

if __name__ == '__main__':
    main()
