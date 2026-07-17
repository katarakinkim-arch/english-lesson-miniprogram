# -*- coding: utf-8 -*-
"""Bounds checker: flag any shape that extends beyond the 13.333x7.5in canvas."""
import sys
from pptx import Presentation
from pptx.util import Emu

SRC = sys.argv[1] if len(sys.argv) > 1 else r"preview_v7/oxiangxue.pptx"
EPS = 0.02  # inches tolerance
prs = Presentation(SRC)
W_IN = prs.slide_width / 914400.0
H_IN = prs.slide_height / 914400.0
print(f"canvas = {W_IN:.3f} x {H_IN:.3f} in")
problems = 0
for i, slide in enumerate(prs.slides, 1):
    worst_r = 0.0
    worst_b = 0.0
    for sh in slide.shapes:
        l = sh.left / 914400.0 if sh.left is not None else 0
        t = sh.top / 914400.0 if sh.top is not None else 0
        w = sh.width / 914400.0 if sh.width is not None else 0
        h = sh.height / 914400.0 if sh.height is not None else 0
        r = l + w
        b = t + h
        if r > W_IN + EPS:
            worst_r = max(worst_r, r - W_IN)
        if b > H_IN + EPS:
            worst_b = max(worst_b, b - H_IN)
    if worst_r > 0 or worst_b > 0:
        problems += 1
        print(f"  slide {i:02d}: OVERFLOW right+{worst_r:.2f}in bottom+{worst_b:.2f}in")
    else:
        print(f"  slide {i:02d}: ok")
print("PROBLEMS" if problems else "ALL WITHIN BOUNDS")
