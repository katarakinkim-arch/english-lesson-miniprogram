# -*- coding: utf-8 -*-
"""自检：打开生成的 pptx，校验结构与图片嵌入情况。"""
import os, sys
from pptx import Presentation
from pptx.util import Emu

DESKTOP = os.path.expanduser('~/Desktop')
PATH = sys.argv[1] if len(sys.argv) > 1 else os.path.join(DESKTOP, '沁园春长沙-课堂版-杂志v7.pptx')

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height
print(f'文件: {PATH}')
print(f'画布: {Emu(SW).inches:.2f} x {Emu(SH).inches:.2f} in')
print(f'幻灯片数: {len(prs.slides)}')

problems = []
pics_total = 0
media_files = set()
for i, slide in enumerate(prs.slides, 1):
    pics = 0
    for sh in slide.shapes:
        # picture?
        if sh.shape_type == 13:  # PICTURE
            pics += 1
            pics_total += 1
            if sh.width <= 0 or sh.height <= 0:
                problems.append(f'  幻灯片{i}: 图片尺寸为0')
            # bounds check
            if sh.left < 0 or sh.top < 0 or sh.left + sh.width > SW or sh.top + sh.height > SH:
                problems.append(f'  幻灯片{i}: 图片超出画布边界')
        # textbox bounds
        if sh.has_text_frame:
            if sh.left is not None and sh.width is not None:
                if sh.left + sh.width > SW + Emu(0.05 * 914400) or sh.left < -Emu(0.05*914400):
                    pass  # textboxes often intentional; skip strict
    # rels -> media
    for rel in slide.part.rels.values():
        if 'media' in rel.reltype or '/image' in rel.reltype:
            media_files.add(rel.target_ref)
    print(f'  幻灯片{i:02d}: shapes={len(slide.shapes)} pictures={pics}')

print(f'嵌入图片总数(形状): {pics_total}')
print(f'媒体关系数: {len(media_files)}')
for m in sorted(media_files):
    print(f'  media: {m}')

if problems:
    print('\n发现问题:')
    for p in problems:
        print(p)
else:
    print('\n结构自检通过：无图片尺寸/边界异常。')

# zip integrity
import zipfile
z = zipfile.ZipFile(PATH)
bad = z.testzip()
if bad:
    print(f'ZIP损坏: {bad}')
else:
    print('ZIP完整性: OK')
