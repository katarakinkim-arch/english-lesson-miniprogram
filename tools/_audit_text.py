# tools/_audit_text.py
# Highest-standard text-layer audit for classroom PPTs.
# Extracts ALL text (incl. table cells) per slide via python-pptx (decompresses XML),
# so Chinese teacher-tone patterns are actually visible (grep on the .pptx binary is blind).
import sys, os, glob, re
from pptx import Presentation

# ---- BLOCK: must NEVER appear in a student-facing classroom PPT ----
BLOCK = [
    '老师','教师','老师们','教师们',
    '同学们','请大家','请同学','下面请','我们说','同学们看','同学们想','同学们思考','请大家思考',
    '预设回答','板书时机','易错点提醒','易错点','参考答案·教师用','参考答案(教师用)','教师用参考答案',
    '教师讲解','教师导语','教师讲','老师讲','跟我读','一起读','齐读','跟读',
    '板书设计','教学过程','学情分析','教学目标','教学重点','教学难点','教学设想','教法','学法指导',
    '课堂小结','授课','讲授','授课对象','教学反思','备课',
]
# ---- WARN: likely teacher-facing but may be acceptable in context; human review ----
WARN = [
    '注意看','我们一起来','下面我们','请翻开','请打开','请大家','我们学习','这节课','本课时',
    '想一想','思考','我们来',
]
# Phrases that contain a BLOCK substring but are LEGITIMATE in a student PPT
# (literary quote / biographical fact / fictional character). Local-context whitelist.
WHITELIST = ['同学少年','恰同学','公社同学','乡村教师','乡村教员','人民教师']

def slide_texts(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if t and t.strip():
                out.append(t.strip())
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        out.append(cell.text.strip())
        if shape.has_text_frame and shape.text_frame.text:
            out.append(shape.text_frame.text.strip())
    return out

def audit(pptx):
    prs = Presentation(pptx)
    rows = []
    block_hits = []
    warn_hits = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [t for t in slide_texts(slide) if t]
        blob = '\n'.join(texts)
        total_chars = sum(len(t) for t in texts)
        # block scan (whole-pattern, not substring of other words where ambiguous)
        for pat in BLOCK:
            start = 0
            while True:
                idx = blob.find(pat, start)
                if idx < 0:
                    break
                # local-context whitelist: skip if pat sits inside a legit phrase
                window = blob[max(0, idx-6):idx+len(pat)+6]
                if any(wp in window for wp in WHITELIST):
                    start = idx + len(pat)
                    continue
                ctx = blob[max(0, idx-25):idx+len(pat)+25].replace('\n', ' / ')
                block_hits.append((i, pat, ctx))
                start = idx + len(pat)
        for pat in WARN:
            if pat in blob:
                idx = blob.find(pat)
                ctx = blob[max(0, idx-20):idx+len(pat)+20].replace('\n', ' / ')
                warn_hits.append((i, pat, ctx))
        rows.append((i, len(texts), total_chars))
    return rows, block_hits, warn_hits

if __name__ == '__main__':
    files = sorted(glob.glob('preview_v7/*.pptx'))
    os.chdir(os.path.dirname(os.path.abspath(__file__)) + '/..') if False else None
    for f in files:
        rows, bh, wh = audit(f)
        print(f"\n{'='*70}")
        print(f"FILE: {os.path.basename(f)}  | slides={len(rows)}")
        print(f"{'='*70}")
        if bh:
            print(f"  >>> [BLOCK TONE] {len(bh)} hit(s) — FAIL:")
            for i, pat, ctx in bh:
                print(f"      slide {i:>2}: '{pat}'  …{ctx}…")
        else:
            print("  [BLOCK TONE] CLEAN ✅")
        if wh:
            print(f"  >>> [WARN TONE] {len(wh)} hit(s) — review:")
            for i, pat, ctx in wh:
                print(f"      slide {i:>2}: '{pat}'  …{ctx}…")
        else:
            print("  [WARN TONE] none")
        # overflow-risk heuristic
        risky = [r for r in rows if r[2] > 700]
        if risky:
            print(f"  >>> [DENSITY] slides with >700 chars (overflow risk): " +
                  ", ".join(f"s{r[0]}({r[2]})" for r in risky))
        empties = [r for r in rows if r[1] == 0]
        if empties:
            print(f"  >>> [EMPTY] slides with no text (photo-only?): " +
                  ", ".join(f"s{r[0]}" for r in empties))
        print(f"  [SLIDE CHARS] " + " ".join(f"s{r[0]}:{r[2]}" for r in rows))
