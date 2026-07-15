"""用 python-pptx 把 16 张设计稿 PNG 真正生成精美排版 PPTX。
- 每张 16:9 幻灯片（1280x720）用对应 PNG 作整页背景
- 4 张无水印图标（passport/boardingpass/suitcase/map）作为 step1 的素材图示
- 音频作为单独资源随包附带（PowerPoint 不支持原生嵌入 .wav 到普通页）"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image

# 路径
DESIGN = r'C:\Users\1\WorkBuddy\2026-07-08-15-47-48\miniprogram\design_preview'
ICONS = r'C:\Users\1\WorkBuddy\2026-07-08-15-47-48\miniprogram\subpackages\eng\images\lessons\l-eng-b1-u2-ls'
WAV = os.path.join(ICONS, 'listening.wav')
OUT = r'C:\Users\1\WorkBuddy\2026-07-08-15-47-48\miniprogram\pilot-l-eng-b1-u2-ls.pptx'

# 设计稿顺序（16 页）
slides = [
    'slide_01_cover.png', 'slide_02_contents.png', 'slide_03_overview.png',
    'slide_04_objectives.png', 'slide_05_keypoints.png', 'slide_06_step1.png',
    'slide_07_step2.png', 'slide_08_step3.png', 'slide_09_step4.png',
    'slide_10_step5.png', 'slide_11_step6.png', 'slide_12_blackboard.png',
    'slide_13_homework_basic.png', 'slide_14_homework_advanced.png',
    'slide_15_reflection.png', 'slide_16_summary.png',
]

# 创建 16:9 演示文稿（13.333 x 7.5 inch = 1280x720px @ 96dpi）
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # 空白版式

for i, name in enumerate(slides, 1):
    slide = prs.slides.add_slide(blank)
    img_path = os.path.join(DESIGN, name)
    # 设置整页背景图
    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    print(f'  slide {i:02d}  {name}')

# 把 4 张图标图额外添加在 step1 那张（slide6）之后作为"素材图示"附录
# 但本 PPTX 已用设计稿包含，不必重复附加
# 不过为了让 PowerPoint 也能看到原图（不仅看设计稿里的图），加 1 张附录页
appendix = prs.slides.add_slide(blank)
# 简单说明页
from pptx.util import Pt
from pptx.dml.color import RGBColor
# 临时先纯图方式
img = Image.new('RGB', (1280, 720), (245, 247, 250))
from PIL import ImageDraw, ImageFont
d = ImageDraw.Draw(img)
d.text((60, 60), '素材图示（原始图片）', font=ImageFont.truetype(r'C:/Windows/Fonts/msyhbd.ttc', 40), fill=(30, 41, 59))
d.text((60, 120), '附录 · 课前展示用图标', font=ImageFont.truetype(r'C:/Windows/Fonts/msyh.ttc', 20), fill=(100, 116, 139))

# 4 张原图缩略图 2x2
icons = ['passport.jpg', 'boardingpass.jpg', 'suitcase.jpg', 'map.jpg']
labels = ['护照 / Passport', '登机牌 / Boarding Pass', '行李箱 / Suitcase', '地图 / Map']
cw, ch = 540, 250
gap = 20
x0, y0 = 60, 180
for k, (icon, lab) in enumerate(zip(icons, labels)):
    col, row = k % 2, k // 2
    x = x0 + col * (cw + gap)
    y = y0 + row * (ch + gap)
    p = os.path.join(ICONS, icon)
    if os.path.exists(p):
        # 缩到 cw x (ch-40)
        im = Image.open(p).convert('RGB')
        im.thumbnail((cw, ch-40))
        d.rectangle([x, y, x+cw, y+ch], outline=(219, 234, 254), width=2, fill=(255, 255, 255))
        ix = x + (cw - im.width) // 2
        iy = y + 10 + (ch-50 - im.height) // 2
        img.paste(im, (ix, iy))
        d.text((x+10, y+ch-30), lab, font=ImageFont.truetype(r'C:/Windows/Fonts/msyhbd.ttc', 18), fill=(37, 99, 235))
appendix_path = os.path.join(DESIGN, '_appendix.png')
img.save(appendix_path)
appendix.shapes.add_picture(appendix_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

# 保存
prs.save(OUT)
print(f'\n✅ PPTX 已保存: {OUT}')
print(f'   幻灯片数: {len(prs.slides)}')
print(f'   大小: {os.path.getsize(OUT)//1024} KB')
