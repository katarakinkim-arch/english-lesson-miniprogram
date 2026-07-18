# -*- coding: utf-8 -*-
"""精细调研版 · 通用 9 页学生课堂 PPT 渲染器（语文 cn 批次共用）。

严格遵循 进度与规范说明.md §4.4：
  9 页固定结构 = 封面 / 学习目标 / 背景与权威调研 / 重点 / 方法 / 难点 / 板书精华 / 作业 / 单元小结
  设计系统、字体、16:9 画布全部复用 _classroom_lib。
  全篇学生视角、中性口吻；禁用词经 sanitize 清洗；无契合照片时学科色块兜底。
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from _classroom_lib import (
    PAPER, INK, FROST, XIANG, MUTED, WHITE, GOLD, SOFT,
    KAI, HEI, W, H, M, CW, Inches as _I, Pt as _P,
    set_ea, bg, textbox, rule, kicker, page_num, caption, new_slide, new_presentation, PAGE,
)
from _audit_text import BLOCK

HERE = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(os.path.dirname(HERE), 'preview_v7', '_fine_data')

PAL = [FROST, XIANG, GOLD, MUTED]

# ---------- 安全清洗：剥离教师口吻 / 禁用词 ----------
_REPL = [
    ('讲授', '讲解'), ('授课', '讲解'), ('课堂小结', '单元小结'),
    ('教学目标', '学习目标'), ('教学重点', '重点'), ('教学难点', '难点'),
    ('学情分析', ''), ('教法', '方法'), ('学法指导', ''),
    ('教学反思', ''), ('备课', ''), ('预设回答', ''),
    ('易错点提醒', ''), ('易错点', '注意'), ('板书设计', '结构'),
    ('板书时机', ''), ('教学过程', '学习过程'),
    ('教师讲解', '讲解'), ('教师导语', ''), ('教师讲', '讲'),
    ('老师讲', '讲'), ('跟我读', ''), ('一起读', ''), ('齐读', ''),
    ('跟读', ''), ('同学们', '你们'), ('请大家', ''), ('请同学', ''),
    ('下面请', ''), ('我们说', ''), ('教师', ''), ('老师', ''),
    ('学生', '你'), ('思考', '回想'), ('这节课', '本课'), ('本课时', '本课'),
]
def san(t):
    if t is None:
        return ''
    t = str(t)
    for a, b in _REPL:
        t = t.replace(a, b)
    return t.strip()

def load(id):
    with open(os.path.join(DATA, id + '.json'), encoding='utf-8') as f:
        return json.load(f)

def split_items(s):
    s = san(s)
    items, cur = [], []
    for ch in s:
        if ch in '①②③④⑤⑥⑦⑧':
            if cur:
                items.append(''.join(cur).strip())
            cur = []
        else:
            cur.append(ch)
    if cur:
        items.append(''.join(cur).strip())
    return [i for i in items if i]

def parse_exercises(ex):
    ex = san(ex)
    if '参考答案' in ex:
        ex = ex.split('参考答案')[0]
    base = adv = ''
    if '【基础作业】' in ex:
        rest = ex.split('【基础作业】', 1)[1]
        if '【提高作业】' in rest:
            base, adv = rest.split('【提高作业】', 1)
        else:
            base = rest
    elif '【提高作业】' in ex:
        adv = ex.split('【提高作业】', 1)[1]
    return base.strip(), adv.strip()

BOX = '─│┌┐└┘├┤┬┴┼'
def clean_bb(b):
    out = []
    for ln in san(b).split('\n'):
        s = ln
        for ch in BOX:
            s = s.replace(ch, '')
        s = s.strip()
        if s:
            out.append(s)
    return out

# ---------- 页面 ----------
def s_cover(prs, BLANK, d, lead):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    rule(s, M, Inches(0.55), Inches(0.9), GOLD, 3)
    kicker(s, f"{d['book']} · 第{d['unitNumber']}单元 {d['unitTitle']} · 第{d['periodNumber']}课时", M, Inches(0.95), FROST)
    t = d['title']
    main, sub = (t.split('——', 1) if '——' in t else (t, ''))
    textbox(s, M, Inches(1.95), Inches(12.2), Inches(1.3),
            [{'text': main, 'size': 40, 'color': INK, 'bold': True, 'font': KAI, 'line': 1.2}])
    if sub:
        textbox(s, M, Inches(3.15), Inches(12.2), Inches(0.9),
                [{'text': sub, 'size': 22, 'color': FROST, 'bold': True, 'font': KAI}])
    textbox(s, M, Inches(4.35), Inches(12.0), Inches(1.3),
            [{'text': lead, 'size': 16, 'color': MUTED, 'font': KAI, 'line': 1.5}])
    textbox(s, M, H - Inches(0.55), Inches(9), Inches(0.4),
            [{'text': '学生课堂版 · 共 9 页', 'size': 12, 'color': MUTED, 'font': HEI}])
    page_num(s)

def s_objectives(prs, BLANK, d):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '学习目标', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': '四向目标', 'size': 30, 'color': INK, 'bold': True, 'font': KAI}])
    objs = d['objectives']
    cw = (CW - Inches(0.4)) / 2
    ch = Inches(2.15)
    x0, y0 = M, M + Inches(1.45)
    for i, o in enumerate(objs):
        o = san(o)
        lbl, body = (o.split('：', 1) if '：' in o else (o, o))
        r, c = i // 2, i % 2
        x = x0 + c * (cw + Inches(0.4))
        y = y0 + r * (ch + Inches(0.3))
        col = PAL[i % len(PAL)]
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col
        card.line.width = Pt(1.6); card.shadow.inherit = False
        textbox(s, x + Inches(0.3), y + Inches(0.22), cw - Inches(0.6), Inches(0.5),
                [{'text': lbl, 'size': 18, 'color': col, 'bold': True, 'font': HEI}])
        textbox(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.6), ch - Inches(1.05),
                [{'text': body, 'size': 13.5, 'color': INK, 'font': KAI, 'line': 1.45}])
    page_num(s)

def s_background(prs, BLANK, d, sources):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '背景 · 权威调研', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': '从教材到课堂', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
    col_w = (CW - Inches(0.5)) / 2
    lx = M
    ta = san(d['textbookAnalysis'])
    textbox(s, lx, M + Inches(1.35), col_w, Inches(4.7),
            [{'text': ta, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.55}])
    rx = M + col_w + Inches(0.5)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, M + Inches(1.35), col_w, Inches(3.05))
    panel.fill.solid(); panel.fill.fore_color.rgb = INK; panel.shadow.inherit = False
    textbox(s, rx + Inches(0.3), M + Inches(1.5), col_w - Inches(0.6), Inches(0.4),
            [{'text': '权威来源', 'size': 16, 'color': GOLD, 'bold': True, 'font': KAI}])
    paras = []
    for src in sources:
        paras.append({'text': '来源：' + san(src), 'size': 12, 'color': WHITE, 'font': HEI,
                      'line': 1.4, 'space_after': 6})
    textbox(s, rx + Inches(0.3), M + Inches(2.0), col_w - Inches(0.6), Inches(2.2), paras)
    textbox(s, rx, M + Inches(4.6), col_w, Inches(0.5),
            [{'text': '本课件未使用无关配图（学科色块兜底）', 'size': 12, 'color': MUTED, 'font': HEI}])
    page_num(s)

def s_keypoints(prs, BLANK, d):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '重点', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': '本课要点', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
    items = split_items(d['keyPoints'])
    n = len(items)
    cols = n if n <= 3 else 2
    rows = (n + cols - 1) // cols
    cw = (CW - Inches(0.4) * (cols - 1)) / cols
    ch = Inches(4.2) if rows == 1 else Inches(2.7)
    y0 = M + Inches(1.5)
    for i, it in enumerate(items):
        r, c = i // cols, i % cols
        x = M + c * (cw + Inches(0.4))
        y = y0 + r * (ch + Inches(0.3))
        col = PAL[i % len(PAL)]
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col
        card.line.width = Pt(1.6); card.shadow.inherit = False
        tag = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.6))
        tag.fill.solid(); tag.fill.fore_color.rgb = col; tag.line.fill.background(); tag.shadow.inherit = False
        textbox(s, x, y + Inches(0.1), cw, Inches(0.4),
                [{'text': f'要点 {i+1}', 'size': 15, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, x + Inches(0.3), y + Inches(0.8), cw - Inches(0.6), ch - Inches(1.0),
                [{'text': it, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.5}])
    page_num(s)

def s_methods(prs, BLANK, d):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '方法', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': '学习路径', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
    items = split_items(d['teachingMethods'])
    y = M + Inches(1.4)
    pitch = Inches(1.12)
    for i, it in enumerate(items):
        col = PAL[i % len(PAL)]
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, y, CW, Inches(0.95))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col
        card.line.width = Pt(1.6); card.shadow.inherit = False
        nb = s.shapes.add_shape(MSO_SHAPE.OVAL, M + Inches(0.22), y + Inches(0.22), Inches(0.5), Inches(0.5))
        nb.fill.solid(); nb.fill.fore_color.rgb = col; nb.line.fill.background(); nb.shadow.inherit = False
        textbox(s, M + Inches(0.22), y + Inches(0.3), Inches(0.5), Inches(0.4),
                [{'text': str(i + 1), 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        if '：' in it:
            ttl, body = it.split('：', 1)
        else:
            ttl, body = it, ''
        textbox(s, M + Inches(0.95), y + Inches(0.13), Inches(2.6), Inches(0.7),
                [{'text': ttl, 'size': 16, 'color': col, 'bold': True, 'font': HEI}])
        textbox(s, M + Inches(3.6), y + Inches(0.13), CW - Inches(3.8), Inches(0.7),
                [{'text': body, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.3}])
        y += pitch
    page_num(s)

def s_difficulties(prs, BLANK, d):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '难点', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': '怎么破', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
    items = split_items(d['difficulties'])
    y = M + Inches(1.4)
    pitch = Inches(1.12)
    for i, it in enumerate(items):
        col = PAL[i % len(PAL)]
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, y, CW, Inches(0.95))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col
        card.line.width = Pt(1.6); card.shadow.inherit = False
        nb = s.shapes.add_shape(MSO_SHAPE.OVAL, M + Inches(0.22), y + Inches(0.22), Inches(0.5), Inches(0.5))
        nb.fill.solid(); nb.fill.fore_color.rgb = col; nb.line.fill.background(); nb.shadow.inherit = False
        textbox(s, M + Inches(0.22), y + Inches(0.3), Inches(0.5), Inches(0.4),
                [{'text': str(i + 1), 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, M + Inches(0.95), y + Inches(0.12), CW - Inches(1.2), Inches(0.72),
                [{'text': it, 'size': 13.5, 'color': INK, 'font': KAI, 'line': 1.3}])
        y += pitch
    page_num(s)

def s_blackboard(prs, BLANK, d):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '板书精华', M, M, FROST)
    lines = clean_bb(d['blackboard'])
    title = lines[0] if lines else '板书精华'
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': title, 'size': 26, 'color': INK, 'bold': True, 'font': KAI}])
    px, py = M, M + Inches(1.25)
    pw, ph = CW, Inches(4.95)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph)
    panel.fill.solid(); panel.fill.fore_color.rgb = WHITE; panel.line.color.rgb = MUTED
    panel.line.width = Pt(1.0); panel.shadow.inherit = False
    paras = [{'text': title, 'size': 16, 'color': FROST, 'bold': True, 'font': HEI, 'space_after': 6}]
    for ln in lines[1:]:
        paras.append({'text': ln, 'size': 11, 'color': INK, 'font': HEI, 'line': 1.15, 'space_after': 0})
    textbox(s, px + Inches(0.35), py + Inches(0.25), pw - Inches(0.7), ph - Inches(0.5), paras)
    page_num(s)

def s_exercises(prs, BLANK, d):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '作业', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11), Inches(0.6),
            [{'text': '分层任务', 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
    base, adv = parse_exercises(d['exercises'])
    tiers = [('基础 · 必做', base or '（见教材练习）', FROST)]
    if adv:
        tiers.append(('提高 · 选做', adv, XIANG))
    cw = (CW - Inches(0.4) * (len(tiers) - 1)) / len(tiers)
    ch = Inches(4.3)
    y = M + Inches(1.5)
    for i, (tag, body, col) in enumerate(tiers):
        x = M + i * (cw + Inches(0.4))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col
        card.line.width = Pt(1.6); card.shadow.inherit = False
        tagbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.6))
        tagbar.fill.solid(); tagbar.fill.fore_color.rgb = col; tagbar.line.fill.background(); tagbar.shadow.inherit = False
        textbox(s, x, y + Inches(0.1), cw, Inches(0.4),
                [{'text': tag, 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.6), ch - Inches(1.05),
                [{'text': body, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.5}])
    page_num(s)

def s_summary(prs, BLANK, d, prompt):
    s = new_slide(prs, BLANK)
    bg(s, PAPER)
    kicker(s, '单元小结', M, M, FROST)
    textbox(s, M, M + Inches(0.7), Inches(11.5), Inches(0.6),
            [{'text': f"回到「{d['unitTitle']}」", 'size': 28, 'color': INK, 'bold': True, 'font': KAI}])
    col_w = (CW - Inches(0.5)) / 2
    lx = M
    textbox(s, lx, M + Inches(1.45), col_w, Inches(0.5),
            [{'text': '要点回望', 'size': 18, 'color': FROST, 'bold': True, 'font': KAI}])
    items = split_items(d['keyPoints'])[:3]
    paras = []
    for it in items:
        paras.append({'text': '· ' + it, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 10})
    textbox(s, lx, M + Inches(2.0), col_w, Inches(4.0), paras)
    rx = M + col_w + Inches(0.5)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, M + Inches(1.45), col_w, Inches(4.55))
    panel.fill.solid(); panel.fill.fore_color.rgb = INK; panel.shadow.inherit = False
    rule(s, rx + Inches(0.3), M + Inches(1.75), Inches(0.06), GOLD, 26)
    textbox(s, rx + Inches(0.5), M + Inches(1.8), col_w - Inches(0.8), Inches(0.5),
            [{'text': '反思引导', 'size': 17, 'color': GOLD, 'bold': True, 'font': KAI}])
    textbox(s, rx + Inches(0.5), M + Inches(2.4), col_w - Inches(0.8), Inches(3.4),
            [{'text': prompt, 'size': 15, 'color': WHITE, 'font': KAI, 'line': 1.55}])
    page_num(s)

# ---------- 构建 ----------
def build(id, lead, prompt, sources):
    d = load(id)
    subj = id.split('-')[1]
    out = os.path.join(os.path.dirname(HERE), 'preview_v7', subj, id + '.pptx')
    os.makedirs(os.path.dirname(out), exist_ok=True)
    PAGE[0] = 0
    prs, BLANK = new_presentation()
    s_cover(prs, BLANK, d, lead)
    s_objectives(prs, BLANK, d)
    s_background(prs, BLANK, d, sources)
    s_keypoints(prs, BLANK, d)
    s_methods(prs, BLANK, d)
    s_difficulties(prs, BLANK, d)
    s_blackboard(prs, BLANK, d)
    s_exercises(prs, BLANK, d)
    s_summary(prs, BLANK, d, prompt)
    prs.save(out)
    # 自校验：若有禁用词残留，打印告警（不阻断）
    _selfcheck(out)
    return out

def _selfcheck(out):
    prs = Presentation(out)
    for i, slide in enumerate(prs.slides, 1):
        blob = '\n'.join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip())
        for pat in BLOCK:
            if pat in blob:
                print(f"  [WARN] slide {i}: 残留禁用词 '{pat}' -> {blob[blob.find(pat)-10:blob.find(pat)+12]!r}")
