# -*- coding: utf-8 -*-
"""边界 + 文字溢出 校验：客观检查每张幻灯片是否有形状超出页面、文字被裁。"""
import sys
from pptx import Presentation
from pptx.util import Emu, Pt

PATH = sys.argv[1] if len(sys.argv) > 1 else r'preview_v7\qinyuanchun_fused.pptx'
TOL = Emu(int(0.02 * 914400))  # 0.02" 容差

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height

def emu_in(v):
    return round(v / 914400, 2)

issues = []
for si, slide in enumerate(prs.slides, 1):
    # 1) 形状越界
    for sh in slide.shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:
            continue
        if l is None or t is None or w is None or h is None:
            continue
        r, b = l + w, t + h
        if r > SW + TOL:
            issues.append(f'[P{si:02d}] 右溢 {emu_in(r-SW)}"  {sh.shape_type}')
        if b > SH + TOL:
            issues.append(f'[P{si:02d}] 底溢 {emu_in(b-SH)}"  {sh.shape_type}')
        if l < -TOL:
            issues.append(f'[P{si:02d}] 左溢 {emu_in(-l)}"')
        if t < -TOL:
            issues.append(f'[P{si:02d}] 顶溢 {emu_in(-t)}"')
    # 2) 文字估算是否超出文本框高度（粗略：CJK 字宽≈字号pt）
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        try:
            bw = sh.width / 12700  # pt
            bh = sh.height / 12700  # pt
        except Exception:
            continue
        if bw <= 0 or bh <= 0:
            continue
        total_h = 0.0
        for p in tf.paragraphs:
            txt = ''.join(r.text for r in p.runs)
            if not txt:
                continue
            sizes = [r.font.size.pt for r in p.runs if r.font.size]
            fs = max(sizes) if sizes else 18
            ls = p.line_spacing if p.line_spacing else 1.0
            if isinstance(ls, float):
                lh = fs * ls
            else:
                lh = fs * 1.2
            cpl = max(1, int(bw / fs))
            lines = max(1, (len(txt) + cpl - 1) // cpl)
            sa = p.space_after.pt if p.space_after else 0
            sb = p.space_before.pt if p.space_before else 0
            total_h += lines * lh + sa + sb
        if total_h > bh + 6:  # 超出文本框 6pt 以上报警
            issues.append(f'[P{si:02d}] 文字估算溢出 需~{round(total_h)}pt / 框{bh:.0f}pt  "{tf.text[:18]}..."')

if issues:
    print('发现问题 %d 处：' % len(issues))
    for i in issues:
        print('  ', i)
else:
    print('OK：无越界、无文字溢出估算。')
print('slide size = %.2f x %.2f in' % (emu_in(SW), emu_in(SH)))
