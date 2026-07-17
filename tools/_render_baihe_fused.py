# -*- coding: utf-8 -*-
"""百合花 课堂版 PPT — 融合版（fused from top high-praise lessons）
融合两堂高赞课的招牌招式，嵌入本课教材结构：
  · 郁雪琳省一等奖课：以"细节"为抓手 + "有人情味的英雄" + 新媳妇=那朵纯洁百合花
    + 以小见大 + 茅盾评"清新俊逸" + "半瓣花上说人情"
  · 10436字公开课逐字稿：年龄钩子勾连单元青春 + 反右1958背景"肝胆相照生死与共"
    + 学生提问链（死封建/平常/缝补）+ 身份平常vs精神不平常
教材结构（本课6步）：四幕脉络/被子意象演变/三个反复细节/限知视角
设计系统沿用沁园春杂志风（暖纸/墨蓝/霜红/湘碧/金 + 楷体/雅黑）。
python-pptx 生成可编辑 .pptx。图片：PIL 手绘百合花专属插图（内容专属，非通用封面）。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as ET
from PIL import Image, ImageDraw, ImageFilter
import math

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_PPTX = os.path.join(BASE, 'preview_v7', 'baihe_fused.pptx')
IMG_CACHE = os.path.join(BASE, 'preview_v7', '_baihe_imgs')
os.makedirs(IMG_CACHE, exist_ok=True)
os.makedirs(os.path.join(BASE, 'preview_v7'), exist_ok=True)

# ---------- design system (same as 沁园春) ----------
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
INK   = RGBColor(0x1C, 0x2A, 0x33)
FROST = RGBColor(0xB2, 0x3A, 0x2A)
XIANG = RGBColor(0x2E, 0x7D, 0x6B)
MUTED = RGBColor(0x6B, 0x62, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xC8, 0xA8, 0x6B)
SOFT  = RGBColor(0xE7, 0xDF, 0xD2)

KAI = '楷体'
HEI = '微软雅黑'
FKAI = r'C:\Windows\Fonts\simkai.ttf'
FHEI = r'C:\Windows\Fonts\msyh.ttc'

W = Inches(13.333); H = Inches(7.5)
M = Inches(0.7); CW = Inches(13.333 - 1.4)

# ---------- PIL illustrations (content-specific 百合花) ----------
def _lily(draw, cx, cy, scale, petal_col=(252,250,245), stamen=(0xD4,0xA0,0x3A)):
    """画一朵百合花：6瓣 + 雄蕊。"""
    # 花瓣（椭圆，6片旋转）
    for i in range(6):
        ang = math.radians(i * 60)
        px = cx + math.cos(ang) * scale * 0.55
        py = cy + math.sin(ang) * scale * 0.55
        bbox = [px - scale*0.42, py - scale*0.85, px + scale*0.42, py + scale*0.85]
        draw.ellipse(bbox, fill=petal_col, outline=(0xE0,0xD8,0xC8))
    # 花心
    draw.ellipse([cx-scale*0.22, cy-scale*0.22, cx+scale*0.22, cy+scale*0.22], fill=stamen)
    # 雄蕊点
    for i in range(6):
        ang = math.radians(i*60 + 30)
        draw.ellipse([cx+math.cos(ang)*scale*0.15-3, cy+math.sin(ang)*scale*0.15-3,
                      cx+math.cos(ang)*scale*0.15+3, cy+math.sin(ang)*scale*0.15+3],
                     fill=(0xB2,0x3A,0x2A))

def gen_cover():
    p = os.path.join(IMG_CACHE, 'cover.jpg')
    if os.path.exists(p): return p
    img = Image.new('RGB', (1600, 900), (0x1C, 0x2A, 0x33))
    d = ImageDraw.Draw(img)
    # 垂直渐变 深墨蓝→暗红
    for y in range(900):
        t = y / 900.0
        r = int(0x1C + (0x5A - 0x1C) * t)
        g = int(0x2A + (0x28 - 0x2A) * t)
        b = int(0x33 + (0x32 - 0x33) * t)
        d.line([(0, y), (1600, y)], fill=(r, g, b))
    # 淡竖纹理
    for x in range(0, 1600, 80):
        d.line([(x, 0), (x, 900)], fill=(255, 255, 255), width=1)
        # 降透明度近似：叠一层暗色
    overlay = Image.new('RGB', (1600, 900), (0, 0, 0))
    img = Image.blend(img, overlay, 0.0)  # noop placeholder
    d = ImageDraw.Draw(img)
    # 三朵百合花（右上、中下、左下），柔和
    for (cx, cy, sc) in [(1280, 230, 150), (300, 640, 120), (1100, 700, 95)]:
        layer = Image.new('RGBA', (1600, 900), (0, 0, 0, 0))
        ld = ImageDraw.Draw(layer)
        _lily(ld, cx, cy, sc)
        layer = layer.filter(ImageFilter.GaussianBlur(1.2))
        img.paste(layer, (0, 0), layer)
    # 底部压暗
    ov = Image.new('RGBA', (1600, 900), (0, 0, 0, 0))
    od = ImageDraw.Draw(ov)
    for y in range(500, 900):
        a = int(180 * (y - 500) / 400)
        od.line([(0, y), (1600, y)], fill=(0, 0, 0, a))
    img = Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')
    img.save(p, 'JPEG', quality=82)
    return p

def gen_flower():
    p = os.path.join(IMG_CACHE, 'flower.jpg')
    if os.path.exists(p): return p
    img = Image.new('RGB', (1000, 1200), (0xF4, 0xEF, 0xE6))
    d = ImageDraw.Draw(img)
    # 暖底渐变
    for y in range(1200):
        t = y / 1200.0
        r = int(0xF4 - 0x10 * t); g = int(0xEF - 0x14 * t); b = int(0xE6 - 0x18 * t)
        d.line([(0, y), (1000, y)], fill=(r, g, b))
    # 主花
    _lily(d, 500, 520, 320, petal_col=(253, 251, 246))
    # 茎叶
    d.line([(500, 760), (500, 1150)], fill=(0x2E, 0x7D, 0x6B), width=14)
    d.ellipse([430, 950, 620, 1010], fill=(0x2E, 0x7D, 0x6B))
    img.save(p, 'JPEG', quality=82)
    return p

def gen_quilt():
    p = os.path.join(IMG_CACHE, 'quilt.jpg')
    if os.path.exists(p): return p
    img = Image.new('RGB', (1200, 800), (0xF4, 0xEF, 0xE6))
    d = ImageDraw.Draw(img)
    # 被子纹理：菱格 + 百合花纹
    for y in range(0, 800, 100):
        for x in range(0, 1200, 100):
            col = (0xE7, 0xDF, 0xD2) if ((x//100 + y//100) % 2 == 0) else (0xEF, 0xE8, 0xDC)
            d.polygon([(x+50, y), (x+100, y+50), (x+50, y+100), (x, y+50)], fill=col, outline=(0xD4, 0xC8, 0xB0))
    # 角落百合花
    for (cx, cy, sc) in [(300, 200, 90), (900, 600, 70), (600, 400, 60)]:
        _lily(d, cx, cy, sc, petal_col=(253, 250, 244))
    img.save(p, 'JPEG', quality=82)
    return p

def gen_moon():
    p = os.path.join(IMG_CACHE, 'moon.jpg')
    if os.path.exists(p): return p
    img = Image.new('RGB', (1400, 800), (0x35, 0x4A, 0x60))
    d = ImageDraw.Draw(img)
    # 夜空渐变（上深下浅，月下有微光）
    for y in range(800):
        t = y / 800.0
        r = int(0x35 + 0x30 * t); g = int(0x4A + 0x2C * t); b = int(0x60 + 0x28 * t)
        d.line([(0, y), (1400, y)], fill=(r, g, b))
    # 月光晕散（多层）
    for r in range(80, 280, 14):
        a = max(0, 60 - r // 6)
        col = (0xF4, 0xEF, 0xE6)
        d.ellipse([1120 - r, 240 - r, 1120 + r, 240 + r], outline=col)
    # 大月亮
    d.ellipse([1040, 160, 1200, 320], fill=(0xFA, 0xF6, 0xEC), outline=(0xE8, 0xDC, 0xB8))
    # 月面阴影
    d.ellipse([1080, 200, 1150, 260], fill=(0xE0, 0xD0, 0xB0))
    # 远山（提亮成清晰剪影）
    d.polygon([(0, 620), (180, 500), (380, 580), (600, 480), (820, 560), (1050, 490), (1250, 570), (1400, 510), (1400, 800), (0, 800)],
              fill=(0x1E, 0x2C, 0x3A))
    # 中景
    d.polygon([(0, 720), (300, 660), (600, 700), (900, 650), (1200, 690), (1400, 660), (1400, 800), (0, 800)],
              fill=(0x14, 0x1E, 0x28))
    img.save(p, 'JPEG', quality=85)
    return p

# ---------- 内容专属图标（扁平线描，400x400 暖底） ----------
INK_RGB=(0x1C,0x2A,0x33); FROST_RGB=(0xB2,0x3A,0x2A); XIANG_RGB=(0x2E,0x7D,0x6B)
GOLD_RGB=(0xC8,0xA8,0x6B); SOFT_RGB=(0xE7,0xDF,0xD2); PAPER_RGB=(0xF4,0xEF,0xE6)
def _icon_bg(d):
    d.rectangle([0,0,400,400], fill=PAPER_RGB)
    d.rounded_rectangle([8,8,392,392], radius=24, outline=SOFT_RGB, width=3)
def _person(d, cx, cy, col, cap=False, bun=False, flower=False):
    # 头
    d.ellipse([cx-26, cy-70, cx+26, cy-18], fill=col)
    # 帽/发髻
    if cap: d.rectangle([cx-30, cy-78, cx+30, cy-60], fill=col)
    if bun: d.ellipse([cx-16, cy-92, cx+16, cy-64], fill=col)
    # 身（三角）
    d.polygon([(cx-44, cy+90), (cx+44, cy+90), (cx, cy-10)], fill=col)
    if flower:
        d.ellipse([cx+30, cy-30, cx+48, cy-12], fill=FROST_RGB)
def gen_icon(name):
    p = os.path.join(IMG_CACHE, name + '.jpg')
    if os.path.exists(p): return p
    img = Image.new('RGB', (400, 400), PAPER_RGB)
    d = ImageDraw.Draw(img); _icon_bg(d)
    c = INK_RGB
    if name == 'borrow':       # 借被：菱形被子 + 两手
        for i in range(3):
            d.polygon([(200, 120+i*50), (280, 150+i*50), (200, 180+i*50), (120, 150+i*50)],
                      outline=FROST_RGB, width=4)
        d.ellipse([60, 180, 110, 230], outline=INK_RGB, width=6)   # 左手
        d.ellipse([290, 180, 340, 230], outline=INK_RGB, width=6)  # 右手
    elif name == 'aid':        # 包扎：红十字
        d.rectangle([170, 110, 230, 290], fill=FROST_RGB)
        d.rectangle([110, 170, 290, 230], fill=FROST_RGB)
        d.line([(150, 320), (250, 320)], fill=INK_RGB, width=8)
    elif name == 'sacrifice':  # 牺牲：倒影 + 星
        _person(d, 200, 200, INK_RGB, cap=True)
        d.line([(120, 300), (280, 300)], fill=FROST_RGB, width=6)
        for (sx, sy) in [(150, 90), (260, 110), (300, 70)]:
            d.line([(sx, sy-12), (sx, sy+12)], fill=GOLD_RGB, width=4)
            d.line([(sx-12, sy), (sx+12, sy)], fill=GOLD_RGB, width=4)
    elif name == 'cover':      # 盖被：人形 + 被子覆盖
        _person(d, 200, 210, SOFT_RGB)
        d.polygon([(110, 250), (290, 250), (270, 340), (130, 340)], fill=XIANG_RGB)
        for i in range(3):
            d.line([(130+i*40, 250), (130+i*40, 340)], fill=PAPER_RGB, width=3)
    elif name == 'gun_flower': # 枪筒野菊
        d.line([(200, 320), (200, 120)], fill=INK_RGB, width=10)  # 枪筒
        d.line([(180, 300), (220, 300)], fill=INK_RGB, width=6)
        for (fx, fy) in [(200, 110), (175, 130), (225, 130)]:
            for ang in range(0, 360, 72):
                a = math.radians(ang)
                d.ellipse([fx+math.cos(a)*10-4, fy+math.sin(a)*10-4, fx+math.cos(a)*10+4, fy+math.sin(a)*10+4], fill=FROST_RGB)
            d.ellipse([fx-6, fy-6, fx+6, fy+6], fill=GOLD_RGB)
    elif name == 'tear':       # 衣服破洞
        d.rectangle([120, 130, 280, 300], outline=INK_RGB, width=6)
        d.ellipse([185, 195, 225, 235], fill=FROST_RGB)            # 破洞
        d.line([(185, 215), (170, 240)], fill=INK_RGB, width=3)    # 缝线
        d.line([(225, 215), (240, 240)], fill=INK_RGB, width=3)
    elif name == 'bun':        # 两个馒头
        d.ellipse([110, 200, 210, 290], fill=SOFT_RGB, outline=INK_RGB, width=4)
        d.ellipse([190, 200, 290, 290], fill=SOFT_RGB, outline=INK_RGB, width=4)
        d.arc([130, 210, 190, 270], 200, 340, fill=INK_RGB, width=3)
        d.arc([210, 210, 270, 270], 200, 340, fill=INK_RGB, width=3)
    elif name == 'soldier':    # 通讯员剪影
        _person(d, 200, 200, INK_RGB, cap=True)
    elif name == 'bride':      # 新媳妇剪影
        _person(d, 200, 200, FROST_RGB, bun=True, flower=True)
    elif name == 'eye':        # 眼睛视野（限知）
        d.ellipse([80, 150, 320, 250], fill=WHITE if False else PAPER_RGB, outline=INK_RGB, width=8)
        d.ellipse([170, 170, 230, 230], fill=INK_RGB)              # 瞳
        d.ellipse([188, 188, 202, 202], fill=PAPER_RGB)            # 高光
        # 视野扇形
        d.pieslice([100, 130, 300, 270], 40, 140, fill=SOFT_RGB, outline=GOLD_RGB)
    elif name == 'women':      # 两位女战士剪影（导入）
        _person(d, 140, 210, FROST_RGB, cap=True)
        _person(d, 260, 210, INK_RGB, bun=True, flower=True)
    elif name == 'paper':      # 旧报纸（1958背景）
        d.rectangle([80, 80, 320, 340], fill=(0xF0,0xE8,0xD0), outline=INK_RGB, width=3)
        d.rectangle([90, 95, 310, 115], fill=INK_RGB)  # 报头条
        d.rectangle([90, 130, 200, 150], fill=FROST_RGB)  # 期号
        for y in [170, 195, 220, 245, 270, 295, 315]:
            d.line([(100, y), (300, y)], fill=INK_RGB, width=2)
    elif name == 'ask':        # 问号对话气泡（提问链）
        # 大气泡
        d.ellipse([80, 90, 260, 270], fill=FROST_RGB, outline=INK_RGB, width=4)
        d.polygon([(140, 270), (160, 270), (130, 320), (150, 270)], fill=FROST_RGB, outline=INK_RGB)
        # 问号
        d.arc([140, 130, 200, 210], 180, 360, fill=WHITE if False else PAPER_RGB, width=12)
        d.ellipse([158, 215, 182, 240], fill=PAPER_RGB)
        # 小气泡
        d.ellipse([260, 260, 340, 340], fill=XIANG_RGB, outline=INK_RGB, width=3)
    elif name == 'pen':        # 钢笔+本子（作业）
        d.rectangle([90, 220, 310, 340], fill=WHITE if False else PAPER_RGB, outline=INK_RGB, width=3)
        for y in [245, 265, 285, 305]:
            d.line([(110, y), (290, y)], fill=SOFT_RGB, width=2)
        # 笔
        d.polygon([(200, 60), (240, 60), (250, 200), (190, 200)], fill=INK_RGB)
        d.polygon([(190, 200), (250, 200), (245, 215), (195, 215)], fill=GOLD_RGB)
    elif name == 'loop':       # 循环箭头（反思）
        d.arc([80, 100, 320, 320], 20, 160, fill=GOLD_RGB, width=10)
        d.polygon([(310, 130), (340, 150), (300, 170)], fill=GOLD_RGB)
        # 中心星
        d.ellipse([180, 180, 220, 220], fill=FROST_RGB)
    img.save(p, 'JPEG', quality=85)
    return p

# 预生成所有图标
for _n in ['borrow','aid','sacrifice','cover','gun_flower','tear','bun','soldier','bride','eye',
           'women','paper','ask','pen','loop']:
    gen_icon(_n)
ICON = {n: os.path.join(IMG_CACHE, n + '.jpg') for n in
        ['borrow','aid','sacrifice','cover','gun_flower','tear','bun','soldier','bride','eye',
         'women','paper','ask','pen','loop']}

# 真实照片（从 Wikimedia Commons 下载，替换部分 PIL 图为真实照片）
PHOTO_DIR = os.path.join(BASE, 'preview_v7', '_baihe_photos')
PHOTO = {
    'lily':      os.path.join(PHOTO_DIR, 'real_lily.jpg'),
    'lily2':     os.path.join(PHOTO_DIR, 'real_lily2.jpg'),
    'moon':      os.path.join(PHOTO_DIR, 'real_moon.jpg'),
    'newspaper': os.path.join(PHOTO_DIR, 'real_newspaper.jpg'),
    'quilt':     os.path.join(PHOTO_DIR, 'real_quilt.jpg'),
    'pen':       os.path.join(PHOTO_DIR, 'real_pen.jpg'),
}

COVER = gen_cover(); FLOWER = gen_flower(); QUILT = gen_quilt(); MOON = gen_moon()

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def set_ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', name)

def bg(slide, color=PAPER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def place_photo(slide, path, x, y, w, h):
    x = int(x); y = int(y); w = int(w); h = int(h)
    pic = slide.shapes.add_picture(path, x, y)
    iw, ih = pic.width, pic.height
    ir = iw / ih; br = w / h
    pic.width = w; pic.height = h
    if ir > br:
        vis = br / ir; crop = (1 - vis) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        vis = (h * ir) / w; crop = (1 - vis) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    pic.left = x; pic.top = y
    return pic

def scrim(slide, x, y, w, h, color, alpha_pct):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    srgb = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    alpha = ET.SubElement(srgb, qn('a:alpha')); alpha.set('val', str(int(alpha_pct * 100000)))
    return shp

def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'space_after' in p: para.space_after = Pt(p['space_after'])
        if 'space_before' in p: para.space_before = Pt(p['space_before'])
        if 'line' in p: para.line_spacing = p['line']
        runs = p['runs'] if 'runs' in p else [p]
        for r in runs:
            run = para.add_run(); run.text = r['text']
            run.font.size = Pt(r['size']); run.font.bold = r.get('bold', False)
            run.font.name = r.get('font', HEI); run.font.color.rgb = r['color']
            set_ea(run, r.get('font', HEI))
    return tb

def rule(slide, x, y, w, color=FROST, thick=2.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thick / 72.0))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def kicker(slide, text, x=M, y=M, color=FROST):
    textbox(slide, x, y, Inches(7), Inches(0.4),
            [{'text': text, 'size': 13, 'color': color, 'bold': True, 'font': HEI, 'space_after': 0}])
    rule(slide, x, y + Inches(0.5), Inches(0.9), color, 2.4)

def new_slide(): return prs.slides.add_slide(BLANK)

def caption(slide, text, x, y, w, color=MUTED):
    textbox(slide, x, y, w, Inches(0.4),
            [{'text': text, 'size': 11, 'color': color, 'font': HEI, 'align': PP_ALIGN.CENTER}])

PAGE = [0]
def page_num(slide, dark=False):
    PAGE[0] += 1
    c = WHITE if dark else MUTED
    textbox(slide, W - Inches(1.3), H - Inches(0.55), Inches(0.8), Inches(0.35),
            [{'text': f'{PAGE[0]:02d}', 'size': 11, 'color': c, 'font': HEI, 'align': PP_ALIGN.RIGHT}])

def quote_block(slide, x, y, w, text, source, color=FROST):
    """引用块：左侧竖条 + 楷体引文 + 来源。"""
    rule(slide, x, y, Inches(0.06), color, 30)
    textbox(slide, x + Inches(0.25), y, w - Inches(0.25), Inches(2.2),
            [{'text': text, 'size': 18, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 8},
             {'text': '—— ' + source, 'size': 12, 'color': MUTED, 'font': HEI}])

def step_card(slide, x, y, w, h, num, title, body_lines, accent=FROST):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent; card.line.width = Pt(1.8); card.shadow.inherit = False
    # 序号圆
    nb = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.5), Inches(0.5))
    nb.fill.solid(); nb.fill.fore_color.rgb = accent; nb.line.fill.background(); nb.shadow.inherit = False
    textbox(slide, x + Inches(0.18), y + Inches(0.24), Inches(0.5), Inches(0.42),
            [{'text': str(num), 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    textbox(slide, x + Inches(0.8), y + Inches(0.18), w - Inches(1.0), Inches(0.5),
            [{'text': title, 'size': 17, 'color': INK, 'bold': True, 'font': HEI}])
    bl = [{'text': t, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.4, 'space_after': 4} for t in body_lines]
    textbox(slide, x + Inches(0.8), y + Inches(0.72), w - Inches(1.0), h - Inches(0.85), bl)

# ===================================================================
# SLIDES
# ===================================================================
def s_cover():
    s = new_slide()
    place_photo(s, PHOTO['lily'], 0, 0, W, H)
    scrim(s, 0, H - Inches(3.6), W, Inches(3.6), INK, 0.62)
    scrim(s, 0, 0, W, Inches(1.2), INK, 0.35)
    kicker(s, '高中语文 · 必修上册 · 第三课', M, M, GOLD)
    textbox(s, M, H - Inches(3.0), CW, Inches(1.7),
            [{'text': '百 合 花', 'size': 64, 'color': WHITE, 'bold': True, 'font': KAI}])
    textbox(s, M, H - Inches(1.4), CW, Inches(0.9),
            [{'runs': [
                {'text': '茹志鹃', 'size': 22, 'color': WHITE, 'font': HEI, 'bold': True},
                {'text': '    1958 年 · 细节里的人性之光', 'size': 16, 'color': SOFT, 'font': HEI}]}])
    textbox(s, M, H - Inches(0.62), CW, Inches(0.4),
            [{'text': '贯穿线索：以「细节」为抓手，读战火中最圣洁的青春', 'size': 13, 'color': GOLD, 'font': HEI, 'bold': True}])
    page_num(s, dark=True)

def s_contents():
    s = new_slide(); bg(s)
    kicker(s, 'CONTENTS · 本课导览')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': '九个篇章，读懂一朵百合花', 'size': 30, 'color': INK, 'bold': True, 'font': HEI}])
    items = [
        ('壹', '导入 · 青春勾连', '年龄钩子，从"恰同学少年"走进战时青春'),
        ('贰', '背景 · 肝胆相照', '1958反右与"生死与共"的同志之情'),
        ('叁', '情节 · 四幕剧', '借→帮→牺→盖，一物串全文'),
        ('肆', '意象 · 被子演变', '不愿借→共用→献出，情感在变'),
        ('伍', '细读 · 三个反复', '枪筒树枝·衣服破洞·两个馒头'),
        ('陆', '人物 · 两人一花', '有人情味的英雄 · 纯洁的百合花'),
        ('柒', '视角 · 限知辨析', '"我"只看见一部分的真实与悬念'),
        ('捌', '提问链 · 以小见大', '学生发问：死封建·平常·缝补'),
        ('玖', '风格 · 清新俊逸', '茅盾之评与"半瓣花上说人情"'),
    ]
    col_w = CW / 3
    for i, (num, t, d) in enumerate(items):
        col = i % 3; row = i // 3
        x = M + col * col_w; y = M + Inches(1.7) + row * Inches(1.7)
        textbox(s, x, y, Inches(0.7), Inches(0.7),
                [{'text': num, 'size': 30, 'color': FROST, 'bold': True, 'font': KAI}])
        textbox(s, x + Inches(0.75), y, col_w - Inches(0.9), Inches(0.5),
                [{'text': t, 'size': 16, 'color': INK, 'bold': True, 'font': HEI}])
        textbox(s, x + Inches(0.75), y + Inches(0.5), col_w - Inches(0.9), Inches(0.9),
                [{'text': d, 'size': 12, 'color': MUTED, 'font': HEI, 'line': 1.4}])
    page_num(s)

def s_intro():
    s = new_slide(); bg(s)
    kicker(s, '壹 · 导入——青春勾连')
    textbox(s, M, M + Inches(0.7), Inches(10.5), Inches(0.7),
            [{'text': '同是风华正茂的青春：今天的我们，与战火中的她们', 'size': 24, 'color': INK, 'bold': True, 'font': HEI}])
    place_photo(s, ICON['women'], W - M - Inches(1.4), M + Inches(0.55), Inches(1.3), Inches(1.3))
    # 左：作家作品卡（学生应知的基础信息）
    lc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, M + Inches(1.7), Inches(5.6), Inches(4.4))
    lc.fill.solid(); lc.fill.fore_color.rgb = SOFT; lc.line.color.rgb = FROST; lc.line.width = Pt(1.5); lc.shadow.inherit = False
    textbox(s, M + Inches(0.3), M + Inches(1.9), Inches(5.0), Inches(0.4),
            [{'text': '作家 · 作品', 'size': 13, 'color': FROST, 'bold': True, 'font': HEI}])
    textbox(s, M + Inches(0.3), M + Inches(2.35), Inches(5.0), Inches(3.6),
            [{'text': '茹志鹃（1925—1998）', 'size': 17, 'color': INK, 'bold': True, 'font': HEI, 'space_after': 8},
             {'text': '当代女作家。代表作《百合花》《剪辑错了的故事》。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55, 'space_after': 10},
             {'text': '《百合花》1958 年发表，是"诗化小说"代表——不写惨烈，专写温柔人性。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55, 'space_after': 10},
             {'text': '文体：短篇小说。以"我"的所见所闻展开——这叫限知视角。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55}])
    # 右：青春对照 + 学习重点
    rx = M + Inches(7.1)
    textbox(s, rx, M + Inches(1.7), Inches(5.2), Inches(0.5),
            [{'text': '青春对照', 'size': 15, 'color': INK, 'bold': True, 'font': HEI}])
    compare = [
        ('我们', '15 · 16 岁', '和平校园，恰同学少年'),
        ('她们', '20 · 28 岁', '战时前线，风华正茂'),
    ]
    for i, (who, age, note) in enumerate(compare):
        y = M + Inches(2.25) + i * Inches(1.05)
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, Inches(5.2), Inches(0.9))
        c.fill.solid(); c.fill.fore_color.rgb = WHITE; c.line.color.rgb = XIANG; c.line.width = Pt(1.2); c.shadow.inherit = False
        textbox(s, rx + Inches(0.25), y + Inches(0.12), Inches(4.7), Inches(0.7),
                [{'runs': [{'text': who + '  ', 'size': 16, 'color': XIANG, 'bold': True, 'font': HEI},
                           {'text': age + '   ', 'size': 16, 'color': FROST, 'bold': True, 'font': HEI},
                           {'text': note, 'size': 13, 'color': MUTED, 'font': KAI}]}])
    # 学习重点条
    hb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, M + Inches(4.55), Inches(5.2), Inches(0.55))
    hb.fill.solid(); hb.fill.fore_color.rgb = INK; hb.line.fill.background(); hb.shadow.inherit = False
    textbox(s, rx, M + Inches(4.62), Inches(5.2), Inches(0.4),
            [{'text': '本课重点：读细节 · 析意象 · 辨视角 · 悟风格', 'size': 13, 'color': GOLD, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    page_num(s)

def s_background():
    s = new_slide(); bg(s)
    kicker(s, '贰 · 背景——肝胆相照，生死与共')
    textbox(s, M, M + Inches(0.7), Inches(10.5), Inches(0.7),
            [{'text': '1958 年反右之后，茹志鹃悲凉地思念起战时的同志关系', 'size': 24, 'color': INK, 'bold': True, 'font': HEI}])
    place_photo(s, PHOTO['newspaper'], W - M - Inches(1.4), M + Inches(0.55), Inches(1.3), Inches(1.3))
    # 左：背景卡
    lc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, M + Inches(1.7), Inches(5.6), Inches(4.4))
    lc.fill.solid(); lc.fill.fore_color.rgb = SOFT; lc.line.fill.background(); lc.shadow.inherit = False
    textbox(s, M + Inches(0.3), M + Inches(1.9), Inches(5.0), Inches(4.0),
            [{'text': '【写作背景】', 'size': 12, 'color': FROST, 'bold': True, 'font': HEI, 'space_after': 8},
             {'text': '《百合花》写于 1958 年 3 月，正值反右斗争（1957–1958）之后。茹志鹃的家庭成员是这场扩大化运动的受害者。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55, 'space_after': 10},
             {'text': '冷峻的现实使人不无悲凉地思念起战时生活——那时的同志关系。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55, 'space_after': 10},
             {'text': '什么叫"同志"？志同道合，有共同的志向与理想。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55}])
    # 右：茹志鹃自述引文
    rx = M + Inches(6.1)
    quote_block(s, rx, M + Inches(1.7), Inches(6.2),
                '"战争使人不能有长谈的机会，但战争却能使人深交。有时仅几十分钟，甚至只来得及瞥了一眼便一闪而过。然而人与人之间，就在这一刹那间，能够肝胆相照，生死与共。"',
                '茹志鹃《我写〈百合花〉的经过》', FROST)
    textbox(s, rx, M + Inches(4.3), Inches(6.2), Inches(1.6),
            [{'text': '这一句，是理解全篇的钥匙：', 'size': 15, 'color': GOLD, 'bold': True, 'font': HEI, 'space_after': 8},
             {'text': '小说写的不是战争的惨烈，而是"一刹那间肝胆相照"的人性之光。读细节，就是读这份生死与共的同志之情。', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.55}])
    page_num(s)

def s_plot():
    s = new_slide(); bg(s)
    kicker(s, '叁 · 情节——四幕剧')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '一床被子，串起四个场景', 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    acts = [
        ('第一幕', '借被子', '通讯员送"我"去包扎所\n青涩腼腆，向新媳妇借被', FROST, 'borrow'),
        ('第二幕', '包扎所', '新媳妇主动来帮忙\n细心照料伤员', XIANG, 'aid'),
        ('第三幕', '牺牲', '通讯员舍身救担架员\n壮烈牺牲', INK, 'sacrifice'),
        ('第四幕', '盖被子', '新媳妇含泪将百合花被\n盖在烈士身上', GOLD, 'cover'),
    ]
    cw = (CW - Inches(0.6)) / 4
    y0 = M + Inches(1.65)
    for i, (act, title, desc, col, icon) in enumerate(acts):
        x = M + i * (cw + Inches(0.2))
        y = y0
        # 顶部色条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
        # 卡片
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.12), cw, Inches(3.95))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = SOFT; card.line.width = Pt(1); card.shadow.inherit = False
        textbox(s, x, y + Inches(0.28), cw, Inches(0.4),
                [{'text': act, 'size': 12, 'color': col, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, x, y + Inches(0.72), cw, Inches(0.7),
                [{'text': title, 'size': 23, 'color': INK, 'bold': True, 'font': KAI, 'align': PP_ALIGN.CENTER}])
        # 场景图标
        isz = Inches(1.15)
        place_photo(s, ICON[icon], x + cw/2 - isz/2, y + Inches(1.45), isz, isz)
        textbox(s, x + Inches(0.2), y + Inches(2.75), cw - Inches(0.4), Inches(1.15),
                [{'text': desc, 'size': 13, 'color': MUTED, 'font': KAI, 'line': 1.45, 'align': PP_ALIGN.CENTER}])
        # 箭头
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw - Inches(0.05), y + Inches(1.95), Inches(0.3), Inches(0.3))
            ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit = False
    # 底部小结条（醒目色块）
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, H - Inches(0.92), CW, Inches(0.62))
    band.fill.solid(); band.fill.fore_color.rgb = INK; band.line.fill.background(); band.shadow.inherit = False
    textbox(s, M, H - Inches(0.83), CW, Inches(0.45),
            [{'text': '一物串全文 · 一物塑两人 · 一物点主旨 ｜ 情节链：借被 → 用被 → 盖被', 'size': 18, 'color': GOLD, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    page_num(s)

def s_quilt_imagery():
    s = new_slide(); bg(s)
    kicker(s, '肆 · 意象——被子的演变')
    textbox(s, M, M + Inches(0.7), Inches(9.0), Inches(0.7),
            [{'text': '同一床被子，三种姿态——人的情感在变', 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    # 被子纹理图（本页核心意象）
    place_photo(s, PHOTO['quilt'], W - M - Inches(2.6), M + Inches(0.55), Inches(2.6), Inches(1.5))
    stages = [
        ('借', '不愿', '私有', '新媳妇刚过门三天，百合花被是嫁妆、是"百年好合"的念想，舍不得借给陌生男人。', FROST),
        ('铺', '共用', '公共', '被子铺在包扎所，给伤员用。从"我的"变成"大家的"，私有的隔阂被战时需要消融。', XIANG),
        ('盖', '献出', '牺牲', '通讯员牺牲，新媳妇含泪把被子盖在他身上。从"共用"到"献出"——纯洁、崇敬、永别。', GOLD),
    ]
    cw = (CW - Inches(1.0)) / 3
    y0 = M + Inches(1.65)
    for i, (act, attitude, nature, desc, col) in enumerate(stages):
        x = M + i * (cw + Inches(0.5))
        y = y0
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(3.95))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col; card.line.width = Pt(2); card.shadow.inherit = False
        # 大字
        textbox(s, x, y + Inches(0.28), cw, Inches(1.1),
                [{'text': act, 'size': 54, 'color': col, 'bold': True, 'font': KAI, 'align': PP_ALIGN.CENTER}])
        textbox(s, x, y + Inches(1.45), cw, Inches(0.5),
                [{'text': attitude + ' · ' + nature, 'size': 18, 'color': INK, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, x + Inches(0.25), y + Inches(2.1), cw - Inches(0.5), Inches(1.7),
                [{'text': desc, 'size': 14, 'color': MUTED, 'font': KAI, 'line': 1.5, 'align': PP_ALIGN.CENTER}])
    # 底部小结条（醒目色块）
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, H - Inches(0.92), CW, Inches(0.62))
    band.fill.solid(); band.fill.fore_color.rgb = INK; band.line.fill.background(); band.shadow.inherit = False
    textbox(s, M, H - Inches(0.83), CW, Inches(0.45),
            [{'text': '意象随情节演变：私有 → 公共 → 牺牲 ｜ 被子没变，变的是人心', 'size': 18, 'color': GOLD, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    page_num(s)

def s_detail_repeats():
    s = new_slide(); bg(s)
    kicker(s, '伍 · 细读——三个反复')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '"半瓣花上说人情"——细节的反复，就是情感的递进', 'size': 24, 'color': INK, 'bold': True, 'font': HEI}])
    details = [
        ('枪筒树枝', '3 次', '纯真', ['第一次：去包扎所路上，枪筒插树枝', '第二次：回来时多了野菊花', '第三次：牺牲后，枪筒的花还在', '→ 战火中依然向往自然、热爱生活的少年'], FROST, 'gun_flower'),
        ('衣服破洞', '3 次', '腼腆→伤痕', ['第一次：借被时挂破衣肩，腼腆', '第二次：包扎所新媳妇要缝，他不让', '第三次：牺牲后破洞仍在，新媳妇含泪缝', '→ 同一个洞，从羞涩到永恒的遗憾'], XIANG, 'tear'),
        ('两个馒头', '2 次', '以物写人', ['第一次：临别留给"我"当晚饭', '第二次：牺牲后，馒头还在身边', '→ 没写一句不舍，馒头替他说话'], GOLD, 'bun'),
    ]
    cw = (CW - Inches(1.0)) / 3
    for i, (name, times, theme, lines, col, icon) in enumerate(details):
        x = M + i * (cw + Inches(0.5))
        y = M + Inches(1.8)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(4.6))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col; card.line.width = Pt(2); card.shadow.inherit = False
        textbox(s, x, y + Inches(0.22), cw, Inches(0.5),
                [{'text': name, 'size': 20, 'color': INK, 'bold': True, 'font': KAI, 'align': PP_ALIGN.CENTER}])
        # 细节图标
        isz = Inches(1.0)
        place_photo(s, ICON[icon], x + cw/2 - isz/2, y + Inches(0.75), isz, isz)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + cw/2 - Inches(0.4), y + Inches(1.85), Inches(0.8), Inches(0.8))
        badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background(); badge.shadow.inherit = False
        textbox(s, x + cw/2 - Inches(0.4), y + Inches(1.97), Inches(0.8), Inches(0.6),
                [{'text': times, 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        bl = [{'text': t, 'size': 13, 'color': INK if not t.startswith('→') else col, 'font': KAI, 'line': 1.4, 'space_after': 5,
               'bold': t.startswith('→')} for t in lines]
        textbox(s, x + Inches(0.3), y + Inches(2.8), cw - Inches(0.6), Inches(1.7), bl)
    page_num(s)

def s_characters():
    s = new_slide(); bg(s)
    kicker(s, '陆 · 人物——两人一花')
    # 左：通讯员
    lx = M; lw = Inches(5.6)
    c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, M + Inches(1.5), lw, Inches(5.0))
    c1.fill.solid(); c1.fill.fore_color.rgb = WHITE; c1.line.color.rgb = FROST; c1.line.width = Pt(2); c1.shadow.inherit = False
    hdr1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, M + Inches(1.5), lw, Inches(0.7))
    hdr1.fill.solid(); hdr1.fill.fore_color.rgb = FROST; hdr1.line.fill.background(); hdr1.shadow.inherit = False
    textbox(s, lx + Inches(0.2), M + Inches(1.62), lw, Inches(0.5),
            [{'text': '小通讯员 · 有人情味的英雄', 'size': 18, 'color': WHITE, 'bold': True, 'font': HEI}])
    textbox(s, lx + Inches(0.3), M + Inches(2.4), lw - Inches(2.6), Inches(4.0),
            [{'text': '他是英雄吗？', 'size': 15, 'color': FROST, 'bold': True, 'font': HEI, 'space_after': 6},
             {'text': '他出一头大汗，见"我"挨近就张皇——来自竹海的男孩子，腼腆、纯真。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8},
             {'text': '枪筒插树枝野菊花——向往自然，阳光乐观。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8},
             {'text': '留两个馒头给"我"——周到暖心的大男孩。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8},
             {'text': '舍身救担架员——英勇无畏。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 10},
             {'text': '→ 不是高大全的英雄，是"有人情味的英雄"——自然的人，有生活实感。', 'size': 14, 'color': FROST, 'bold': True, 'font': KAI, 'line': 1.45}])
    # 通讯员剪影
    place_photo(s, ICON['soldier'], lx + lw - Inches(2.3), M + Inches(2.3), Inches(2.0), Inches(2.0))
    # 右：新媳妇
    rx = M + Inches(6.1); rw = Inches(5.6)
    c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, M + Inches(1.5), rw, Inches(5.0))
    c2.fill.solid(); c2.fill.fore_color.rgb = WHITE; c2.line.color.rgb = XIANG; c2.line.width = Pt(2); c2.shadow.inherit = False
    hdr2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, M + Inches(1.5), rw, Inches(0.7))
    hdr2.fill.solid(); hdr2.fill.fore_color.rgb = XIANG; hdr2.line.fill.background(); hdr2.shadow.inherit = False
    textbox(s, rx + Inches(0.2), M + Inches(1.62), rw, Inches(0.5),
            [{'text': '新媳妇 · 那朵纯洁的百合花', 'size': 18, 'color': WHITE, 'bold': True, 'font': HEI}])
    textbox(s, rx + Inches(0.3), M + Inches(2.4), rw - Inches(2.6), Inches(4.0),
            [{'text': '她是一个怎样的新媳妇？', 'size': 15, 'color': XIANG, 'bold': True, 'font': HEI, 'space_after': 6},
             {'text': '高鼻梁弯眉，粗布却是新的——美丽、淳朴。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8},
             {'text': '把被子朝"我"一送"抱去吧"——有点小脾气、调皮。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8},
             {'text': '经思想斗争才借出被子——真实，不虚假。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8},
             {'text': '牺牲后庄严虔诚地缝衣、拭身、盖被——崇高。', 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 10},
             {'text': '→ 腼腆羞涩，又深明大义重情重义——她就是那一朵纯洁的百合花。', 'size': 14, 'color': XIANG, 'bold': True, 'font': KAI, 'line': 1.45}])
    # 新媳妇剪影
    place_photo(s, ICON['bride'], rx + rw - Inches(2.3), M + Inches(2.3), Inches(2.0), Inches(2.0))
    page_num(s)

def s_perspective():
    s = new_slide(); bg(s)
    kicker(s, '柒 · 视角——限知辨析')
    textbox(s, M, M + Inches(0.7), Inches(9.5), Inches(0.7),
            [{'text': '"我"只看见一部分——这是策略，不是缺陷', 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    # 眼睛/视野示意图
    place_photo(s, ICON['eye'], W - M - Inches(1.6), M + Inches(0.55), Inches(1.6), Inches(1.6))
    # 对比：全知 vs 限知
    step_card(s, M, M + Inches(1.8), Inches(5.6), Inches(2.4), 'A', '全知视角',
              ['叙述者无所不知，可写所有人的内心', '如《三国演义》——上帝俯瞰', '好处：全面；代价：少悬念、少真实感'], MUTED)
    step_card(s, M + Inches(6.1), M + Inches(1.8), Inches(5.6), Inches(2.4), 'B', '限知视角（本篇）',
              ['"我"只在场，只写"我"看见的', '借被子的"受气"现场"我"没亲历→留白', '好处：真实亲切 + 悬念 + 便于抒情'], FROST)
    # 底部追问
    qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, M + Inches(4.5), CW, Inches(1.8))
    qb.fill.solid(); qb.fill.fore_color.rgb = SOFT; qb.line.color.rgb = GOLD; qb.line.width = Pt(1.5); qb.shadow.inherit = False
    textbox(s, M + Inches(0.3), M + Inches(4.7), CW - Inches(0.6), Inches(1.5),
            [{'text': '问题探究', 'size': 12, 'color': GOLD, 'bold': True, 'font': HEI, 'space_after': 6},
             {'text': '新媳妇借被时"受气"的现场，"我"并不在场，小说为什么偏不直接写？', 'size': 16, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 6},
             {'text': '→ 正因没写，才让学生去"还原情境"——这就是限知视角制造的参与感：读者用想象补全，人情才更浓。', 'size': 14, 'color': FROST, 'bold': True, 'font': KAI, 'line': 1.45}])
    page_num(s)

def s_questions():
    s = new_slide(); bg(s)
    kicker(s, '捌 · 提问链——学生发问，以小见大')
    textbox(s, M, M + Inches(0.7), Inches(10.5), Inches(0.5),
            [{'text': '好课，是把提问权交给学生', 'size': 24, 'color': INK, 'bold': True, 'font': HEI}])
    place_photo(s, ICON['ask'], W - M - Inches(1.4), M + Inches(0.4), Inches(1.4), Inches(1.4))
    textbox(s, M, M + Inches(1.25), CW, Inches(0.4),
            [{'text': '手法点睛：被子、百合花本普通——却写尽战争宏大主题。这叫"以小见大"。', 'size': 13, 'color': GOLD, 'bold': True, 'font': HEI}])
    qs = [
        ('Q1', '为什么骂新媳妇"死封建"？', ['嫁妆是"百年好合"的念想，借给伤员怕不吉利', '通讯员没说"为老百姓"的意义，只说"借被子"', '新媳妇羞涩，没说这是新婚被——小战士误会'], FROST),
        ('Q2', '为什么用"平常"修饰战士？', ['他19岁，部队来了就跟走——身份确实平常', '作者要写"普遍的形象"，不是单个个例', '身份平常，反衬舍身救人的不平常精神'], XIANG),
        ('Q3', '人已逝去，为何还要缝补？', ['让他走得体面——破洞露肉不好入殓', '是对牺牲者起码的尊重', '缝的是崇敬，补的是善良高洁'], GOLD),
    ]
    cw = (CW - Inches(1.0)) / 3
    for i, (q, title, ans, col) in enumerate(qs):
        x = M + i * (cw + Inches(0.5)); y = M + Inches(1.85)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(3.9))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col; card.line.width = Pt(2); card.shadow.inherit = False
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.25), Inches(0.6), Inches(0.6))
        badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background(); badge.shadow.inherit = False
        textbox(s, x + Inches(0.25), y + Inches(0.34), Inches(0.6), Inches(0.45),
                [{'text': q, 'size': 14, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, x + Inches(1.0), y + Inches(0.3), cw - Inches(1.2), Inches(0.7),
                [{'text': title, 'size': 15, 'color': INK, 'bold': True, 'font': HEI, 'line': 1.3}])
        al = [{'text': '· ' + a, 'size': 13, 'color': INK, 'font': KAI, 'line': 1.45, 'space_after': 8} for a in ans]
        textbox(s, x + Inches(0.3), y + Inches(1.2), cw - Inches(0.6), Inches(2.5), al)
    page_num(s)

def s_style():
    s = new_slide(); bg(s)
    place_photo(s, PHOTO['moon'], M, M + Inches(1.5), Inches(5.4), Inches(4.8))
    caption(s, '月光照在她脸上——宁静、圣洁', M, M + Inches(6.35), Inches(5.4))
    kicker(s, '玖 · 风格——清新俊逸')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '不写惨烈战火，专写温柔人性', 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    rx = M + Inches(6.0); rw = Inches(6.3)
    # 茅盾评
    quote_block(s, rx, M + Inches(1.6), rw,
                '"《百合花》的细节描写不但描出了人物的风貌，也描出了人物的精神世界。半瓣花上说人情。"',
                '茅盾', GOLD)
    # 风格关键词
    kws = [('清新俊逸', FROST), ('以小见大', XIANG), ('诗化小说', GOLD), ('人性之美', INK)]
    for i, (k, col) in enumerate(kws):
        x = rx + i * Inches(1.55); y = M + Inches(3.9)
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.45), Inches(0.6))
        tag.fill.solid(); tag.fill.fore_color.rgb = col; tag.line.fill.background(); tag.shadow.inherit = False
        textbox(s, x, y + Inches(0.12), Inches(1.45), Inches(0.4),
                [{'text': k, 'size': 14, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
    textbox(s, rx, M + Inches(4.8), rw, Inches(1.5),
            [{'text': '表现庄严的主题，除了慷慨激昂的笔调，还可以有别的风格——', 'size': 15, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 6},
             {'text': '青春的羞涩里有战争的庄严，温暖的回忆里有战争的现实。这，就是百合花的清新俊逸。', 'size': 15, 'color': FROST, 'bold': True, 'font': KAI, 'line': 1.5}])
    page_num(s)

def s_blackboard():
    s = new_slide(); bg(s, INK)
    kicker(s, '板书 · 结构', color=GOLD)
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '百合花 · 茹志鹃', 'size': 28, 'color': WHITE, 'bold': True, 'font': KAI}])
    board = (
        '┌── 百合花 ＝ 纯洁 ＋ 牺牲 ──┐\n'
        '│                              │\n'
        '│ 【四幕脉络】 借 → 帮 → 牺 → 盖 │\n'
        '│                              │\n'
        '│ 【被子意象演变】              │\n'
        '│   借:不愿(私有)               │\n'
        '│   铺:共用(公共)               │\n'
        '│   盖:献出(牺牲)               │\n'
        '│   → 人的情感在变              │\n'
        '│                              │\n'
        '│ 【三个反复】 以小见大          │\n'
        '│   枪筒树枝 | 3次 | 纯真        │\n'
        '│   衣服破洞 | 3次 | 腼腆→伤痕    │\n'
        '│   两馒头    | 2次 | 以物写人    │\n'
        '│                              │\n'
        '│ 【限知视角】 只知"我"所见       │\n'
        '│   → 真实 ＋ 悬念 ＋ 抒情       │\n'
        '└──────────────────────────────┘'
    )
    tb = s.shapes.add_textbox(M + Inches(0.5), M + Inches(1.6), Inches(7.5), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(board.split('\n')):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = para.add_run(); r.text = line
        r.font.size = Pt(13); r.font.name = 'Consolas'; r.font.color.rgb = SOFT
        set_ea(r, 'Consolas')
    # 右侧金句
    rx = M + Inches(8.5)
    textbox(s, rx, M + Inches(2.0), Inches(3.8), Inches(4.0),
            [{'text': '半瓣花上', 'size': 34, 'color': GOLD, 'bold': True, 'font': KAI, 'space_after': 4},
             {'text': '说人情', 'size': 34, 'color': GOLD, 'bold': True, 'font': KAI, 'space_after': 16},
             {'text': '—— 茅盾', 'size': 16, 'color': SOFT, 'font': HEI, 'space_after': 20},
             {'text': '一刹那间', 'size': 30, 'color': FROST, 'bold': True, 'font': KAI, 'space_after': 4},
             {'text': '肝胆相照', 'size': 30, 'color': FROST, 'bold': True, 'font': KAI, 'space_after': 4},
             {'text': '生死与共', 'size': 30, 'color': FROST, 'bold': True, 'font': KAI, 'space_after': 10},
             {'text': '—— 茹志鹃', 'size': 16, 'color': SOFT, 'font': HEI}])
    page_num(s, dark=True)

def s_homework():
    s = new_slide(); bg(s)
    kicker(s, '巩固 · 作业')
    textbox(s, M, M + Inches(0.7), Inches(10.5), Inches(0.7),
            [{'text': '从细节走向迁移', 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    place_photo(s, PHOTO['pen'], W - M - Inches(1.5), M + Inches(0.5), Inches(1.5), Inches(1.5))
    cols = [
        ('基础作业', FROST, ['写 300 字短文，分析通讯员或新媳妇的形象', '至少引用 3 处细节描写并说明作用']),
        ('提高作业', XIANG, ['以"一床被子"为题，写一篇小小说或散文', '要求：注重细节描写，尝试"以小见大"']),
        ('参考示例', GOLD, ['通讯员：腼腆(破洞/大汗)→纯真(野菊花)→暖心(馒头)→英勇(救担架员)', '"有人情味的英雄"=自然的人+生活实感+崇高牺牲']),
    ]
    cw = (CW - Inches(1.0)) / 3
    for i, (title, col, items) in enumerate(cols):
        x = M + i * (cw + Inches(0.5)); y = M + Inches(1.8)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(4.0))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col; card.line.width = Pt(2); card.shadow.inherit = False
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.65))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col; hdr.line.fill.background(); hdr.shadow.inherit = False
        textbox(s, x, y + Inches(0.12), cw, Inches(0.45),
                [{'text': title, 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        il = [{'text': '· ' + it, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 10} for it in items]
        textbox(s, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.6), Inches(3.0), il)
    page_num(s)

def s_summary():
    s = new_slide(); bg(s)
    kicker(s, '总结 · 一朵百合花的光')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.7),
            [{'text': '战争里的青春，为何让我们心头一热？', 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])
    blocks = [
        ('形象之美', FROST, ['小通讯员：有人情味的英雄', '腼腆(破洞)·纯真(野菊)·暖心(馒头)·英勇(救战友)', '新媳妇：那朵纯洁的百合花', '羞涩真实，又深明大义、重情重义']),
        ('手法之美', XIANG, ['以小见大：一床被子写尽战争主题', '细节反复：枪筒·破洞·馒头，层层递进', '限知视角："我"只所见——真实又留悬念']),
        ('主题之美', GOLD, ['茅盾评："半瓣花上说人情"', '一刹那间，肝胆相照，生死与共', '青春的羞涩里，藏着战争的庄严', '温暖的回忆里，是现实的重量']),
    ]
    cw = (CW - Inches(1.0)) / 3
    for i, (title, col, items) in enumerate(blocks):
        x = M + i * (cw + Inches(0.5)); y = M + Inches(1.8)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, Inches(4.2))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = col; card.line.width = Pt(2); card.shadow.inherit = False
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.65))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col; hdr.line.fill.background(); hdr.shadow.inherit = False
        textbox(s, x, y + Inches(0.12), cw, Inches(0.45),
                [{'text': title, 'size': 18, 'color': WHITE, 'bold': True, 'font': HEI, 'align': PP_ALIGN.CENTER}])
        il = [{'text': '· ' + it, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 10} for it in items]
        textbox(s, x + Inches(0.3), y + Inches(0.9), cw - Inches(0.6), Inches(3.2), il)
    page_num(s)

# ===================================================================
# BUILD
# ===================================================================
s_cover(); s_contents(); s_intro(); s_background(); s_plot()
s_quilt_imagery(); s_detail_repeats(); s_characters(); s_perspective()
s_questions(); s_style(); s_blackboard(); s_homework(); s_summary()

prs.save(OUT_PPTX)
print(f'✓ 百合花融合版已生成：{OUT_PPTX}')
print(f'  共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页')
