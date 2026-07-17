# -*- coding: utf-8 -*-
"""几何+内容自检：
1) 每张幻灯片用了哪张图（按媒体文件映射回 classroom_images_v2 文件名）
2) 照片(非满铺)与文本框是否重叠（满铺封面/章节/四宫格/封底除外）
3) 文本框之间是否重叠（同页 body 文本）
"""
import os, sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

DESKTOP = os.path.expanduser('~/Desktop')
IMG_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'classroom_images_v2')
PATH = sys.argv[1] if len(sys.argv) > 1 else os.path.join(DESKTOP, '沁园春长沙-课堂版-杂志v7.pptx')

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height

def emu_in(v): return Emu(v).inches
def area(r): 
    return (r[2]-r[0])*(r[3]-r[1])
def overlap(a, b):
    ix = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    iy = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    return ix*iy

# map media rId -> filename
def media_name(slide, rel_id):
    rel = slide.part.rels[rel_id]
    return os.path.basename(rel.target_ref)

issues = []
for i, slide in enumerate(prs.slides, 1):
    pics, txts = [], []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            r = [sh.left, sh.top, sh.left+sh.width, sh.top+sh.height]
            # which media?
            fname = None
            for rel in slide.part.rels.values():
                if 'media' in rel.reltype and sh._element.find('.//'+'{http://schemas.openxmlformats.org/drawingml/2006/main}blip') is not None:
                    pass
            # determine media via r:embed on blip
            blip = sh._element.find('.//'+'{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if rid:
                    fname = media_name(slide, rid)
            full = (r[0]<=0 and r[1]<=0 and r[2]>=SW and r[3]>=SH)
            pics.append({'rect':r, 'file':fname, 'full':full})
        elif sh.has_text_frame and sh.text_frame.text.strip():
            txts.append([sh.left, sh.top, sh.left+sh.width, sh.top+sh.height, sh.text_frame.text[:18]])

    # pic vs txt overlap (skip full-bleed pics)
    for p in pics:
        if p['full']:
            continue
        for t in txts:
            ov = overlap(p['rect'], t[:4])
            if ov > 0.02*area(p['rect']):  # >2% overlap
                issues.append(f'  幻灯片{i:02d}: 照片[{p["file"]}] 与文本[{t[4]}] 重叠 {(ov/area(p["rect"])*100):.0f}%')

    # txt vs txt overlap (both must be non-trivial)
    for a in range(len(txts)):
        for b in range(a+1, len(txts)):
            ov = overlap(txts[a][:4], txts[b][:4])
            if ov > 0.15*min(area(txts[a][:4]), area(txts[b][:4])):
                issues.append(f'  幻灯片{i:02d}: 文本[{txts[a][4]}] 与 文本[{txts[b][4]}] 重叠')

    pics_used = [p['file'] for p in pics]
    print(f'幻灯片{i:02d}: 图={pics_used}')

print('\n--- 内容映射校验 ---')
# expected mapping per slide index (1-based)
expected = {
 1:'orange_isle', 3:'mao', 4:'mao', 5:'autumn_forest', 6:'autumn_forest',
 9:'red_mountains', 10:'red_mountains', 11:'green_river', 12:'boats',
 13:'eagle', 14:'fish', 15:['red_mountains','green_river','eagle','fish'],
 16:'eagle', 17:'eagle', 18:'fish', 20:'rapids', 21:'mao', 22:'rapids',
 23:'orange_isle', 24:'autumn_forest', 27:'orange_isle',
}
for i, exp in expected.items():
    slide = list(prs.slides)[i-1]
    files = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blip = sh._element.find('.//'+'{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if rid:
                    files.append(media_name(slide, rid))
    exp_list = exp if isinstance(exp, list) else [exp]
    miss = [e for e in exp_list if e not in files]
    extra = [f for f in files if f not in exp_list]
    if miss or extra:
        issues.append(f'  幻灯片{i:02d}: 期望[{exp_list}] 实际[{files}] 缺={miss} 多={extra}')
    else:
        print(f'  ✓ 幻灯片{i:02d} 图片映射正确: {files}')

if issues:
    print('\n发现问题:')
    for x in issues:
        print(x)
else:
    print('\n几何+内容自检全部通过。')
