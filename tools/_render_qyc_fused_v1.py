# -*- coding: utf-8 -*-
"""沁园春·长沙 课堂版 PPT — 融合版 v1 (fused from top high-praise lessons)
融合 top 2-3 高赞课的招牌设计：
  · cv43571662「自由」立骨课：以"自由"串全词 + "竞"字《齐物论》郭象注 + 课内外互证"击/翔" + 9.3阅兵当代打通 + 怅→人民不自由
  · 课堂实录4：视角变换 + 具身表演(手插腰/慢转换) + 闭眼"大脑荧屏"想象四步
  · jiaoxueshilu/861：四幅"青春画卷"脚手架 + "最美奋斗者"导入/打通
设计系统沿用 mag_v3（真实照片 + 严格网格）。python-pptx 生成可编辑 .pptx。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as ET

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_DIR = os.path.join(BASE, 'classroom_images_v2')
OUT_PPTX = os.path.join(BASE, 'preview_v7', 'qinyuanchun_fused.pptx')

# ---------- design system ----------
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
INK   = RGBColor(0x1C, 0x2A, 0x33)
FROST = RGBColor(0xB2, 0x3A, 0x2A)
XIANG = RGBColor(0x2E, 0x7D, 0x6B)
MUTED = RGBColor(0x6B, 0x62, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xC8, 0xA8, 0x6B)
SOFT  = RGBColor(0xE7, 0xDF, 0xD2)

KAI = '楷体'
HEI = '微软雅黑'

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.7)          # page margin
CW = Inches(13.333 - 1.4)
CH = Inches(7.5 - 1.4)

IMG = {
    'orange_isle':   os.path.join(IMG_DIR, 'orange_isle.jpg'),
    'red_mountains': os.path.join(IMG_DIR, 'red_mountains.jpg'),
    'green_river':   os.path.join(IMG_DIR, 'green_river.jpg'),
    'eagle':         os.path.join(IMG_DIR, 'eagle.jpg'),
    'fish':          os.path.join(IMG_DIR, 'fish.jpg'),
    'boats':         os.path.join(IMG_DIR, 'boats.jpg'),
    'mao_art':       os.path.join(IMG_DIR, 'mao_art.jpg'),
    'mao_mural':     os.path.join(IMG_DIR, 'mao_mural_v2.jpg'),
    'rapids':        os.path.join(IMG_DIR, 'rapids.jpg'),
    'autumn_china':  os.path.join(IMG_DIR, 'autumn_china.jpg'),
}

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def set_ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def bg(slide, color=PAPER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def place_photo(slide, path, x, y, w, h):
    """铺满指定框并居中裁剪；包围盒严格等于 (x,y,w,h)，绝不溢出。"""
    pic = slide.shapes.add_picture(path, x, y)
    iw, ih = pic.width, pic.height
    img_ratio = iw / ih
    box_ratio = w / h
    pic.width = w
    pic.height = h
    if img_ratio > box_ratio:
        vis = box_ratio / img_ratio
        crop = (1 - vis) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        vis = (h * img_ratio) / w
        crop = (1 - vis) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    pic.left = x; pic.top = y
    return pic

def scrim(slide, x, y, w, h, color, alpha_pct):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    srgb = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    alpha = ET.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(int(alpha_pct * 100000)))
    return shp

def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'space_after' in p: para.space_after = Pt(p['space_after'])
        if 'space_before' in p: para.space_before = Pt(p['space_before'])
        if 'line' in p: para.line_spacing = p['line']
        runs = p['runs'] if 'runs' in p else [p]
        for r in runs:
            run = para.add_run()
            run.text = r['text']
            run.font.size = Pt(r['size'])
            run.font.bold = r.get('bold', False)
            run.font.name = r.get('font', HEI)
            run.font.color.rgb = r['color']
            set_ea(run, r.get('font', HEI))
    return tb

def rule(slide, x, y, w, color=FROST, thick=2.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thick / 72.0))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def kicker(slide, text, x=M, y=M, color=FROST):
    textbox(slide, x, y, Inches(6), Inches(0.4),
            [{'text': text, 'size': 13, 'color': color, 'bold': True, 'font': HEI,
              'space_after': 0}])
    rule(slide, x, y + Inches(0.5), Inches(0.9), color, 2.4)

def caption(slide, text, x, y, w, color=MUTED):
    textbox(slide, x, y, w, Inches(0.5),
            [{'text': text, 'size': 11, 'color': color, 'font': HEI, 'italic': False}])

def new_slide():
    return prs.slides.add_slide(BLANK)

PAGE = [0]
def page_num(slide, dark=False):
    PAGE[0] += 1
    c = WHITE if dark else MUTED
    textbox(slide, W - Inches(1.3), H - Inches(0.55), Inches(0.8), Inches(0.35),
            [{'text': f'{PAGE[0]:02d}', 'size': 11, 'color': c, 'font': HEI,
              'align': PP_ALIGN.RIGHT}])

# ===================================================================
# SLIDE BUILDERS
# ===================================================================

def s_cover():
    s = new_slide()
    place_photo(s, IMG['orange_isle'], 0, 0, W, H)
    scrim(s, 0, H - Inches(3.4), W, Inches(3.4), INK, 0.62)
    scrim(s, 0, 0, W, Inches(1.2), INK, 0.35)
    kicker(s, '高中语文 · 必修上册 · 第一课', M, M, GOLD)
    textbox(s, M, H - Inches(2.9), CW, Inches(1.7),
            [{'text': '沁园春 · 长沙', 'size': 60, 'color': WHITE, 'bold': True,
              'font': KAI, 'space_after': 0}])
    textbox(s, M, H - Inches(1.15), CW, Inches(0.9),
            [{'runs': [
                {'text': '毛泽东', 'size': 22, 'color': WHITE, 'font': HEI, 'bold': True},
                {'text': '    写于 1925 年 · 橘子洲头', 'size': 16, 'color': SOFT, 'font': HEI},
            ]}])
    textbox(s, M, H - Inches(0.62), CW, Inches(0.4),
            [{'text': '贯穿线索：以「自由」一词立骨', 'size': 13, 'color': GOLD, 'font': HEI, 'bold': True}])
    page_num(s, dark=True)

def s_contents():
    s = new_slide(); bg(s)
    kicker(s, 'CONTENTS · 本课导览')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.8),
            [{'text': '六个篇章，读懂一首词', 'size': 30, 'color': INK, 'bold': True, 'font': HEI}])
    items = [
        ('壹', '知人论世', '作者 · 时代 · 创作背景'),
        ('贰', '初读感知', '诵读全词 · 字音节拍'),
        ('叁', '上阕写景', '万山 · 漫江 · 百舸 · 鹰鱼'),
        ('肆', '炼字品析', '击 / 翔 / 竞 / 独 的妙处'),
        ('伍', '下阕抒情', '同学少年 · 中流击水'),
        ('陆', '主旨升华', '悲秋与颂秋 · 谁主沉浮'),
    ]
    y = M + Inches(1.7)
    rowh = Inches(0.72)
    for num, t, sub in items:
        textbox(s, M, y, Inches(1.0), rowh,
                [{'text': num, 'size': 30, 'color': FROST, 'bold': True, 'font': KAI}])
        textbox(s, M + Inches(1.1), y, Inches(4.5), rowh,
                [{'text': t, 'size': 21, 'color': INK, 'bold': True, 'font': HEI,
                  'space_after': 0}])
        textbox(s, M + Inches(5.6), y + Inches(0.06), Inches(6.0), rowh,
                [{'text': sub, 'size': 13, 'color': MUTED, 'font': HEI}])
        rule(s, M, y + rowh - Inches(0.04), CW, SOFT, 1.0)
        y = y + rowh + Inches(0.02)
    textbox(s, M, y + Inches(0.1), CW, Inches(0.4),
            [{'text': '一条主线：以「自由」串起上阕之景、下阕之情、时代之问', 'size': 13, 'color': XIANG, 'bold': True, 'font': HEI}])
    page_num(s)

def s_divider(num, title, sub, img_key):
    s = new_slide()
    place_photo(s, IMG[img_key], 0, 0, W, H)
    scrim(s, 0, 0, W, H, INK, 0.55)
    textbox(s, M, Inches(2.3), Inches(3), Inches(1.6),
            [{'text': num, 'size': 110, 'color': FROST, 'bold': True, 'font': KAI}])
    rule(s, M + Inches(0.1), Inches(4.0), Inches(1.4), GOLD, 3)
    textbox(s, M, Inches(4.2), CW, Inches(0.7),
            [{'text': title, 'size': 40, 'color': WHITE, 'bold': True, 'font': HEI}])
    textbox(s, M, Inches(5.15), CW, Inches(0.6),
            [{'text': sub, 'size': 16, 'color': SOFT, 'font': HEI}])
    page_num(s, dark=True)

def s_feature(img_key, kicker_txt, headline, paras, cap):
    s = new_slide(); bg(s)
    px = M
    pw = Inches(6.0)
    ph = CH
    place_photo(s, IMG[img_key], px, M, pw, ph)
    caption(s, cap, px, M + ph + Inches(0.12), pw, MUTED)
    tx = M + pw + Inches(0.6)
    tw = W - tx - M
    kicker(s, kicker_txt, tx, M)
    textbox(s, tx, M + Inches(0.7), tw, Inches(1.2),
            [{'text': headline, 'size': 26, 'color': INK, 'bold': True, 'font': HEI, 'line': 1.15}])
    body = []
    for t in paras:
        body.append({'text': t, 'size': 15, 'color': INK, 'font': HEI, 'line': 1.5, 'space_after': 10})
    textbox(s, tx, M + Inches(2.0), tw, Inches(4.6), body, anchor=MSO_ANCHOR.TOP)
    page_num(s)

def s_editorial(kicker_txt, headline, paras, img_key=None, pull=None):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)
    textbox(s, M, M + Inches(0.7), CW - Inches(5.5), Inches(1.4),
            [{'text': headline, 'size': 28, 'color': INK, 'bold': True, 'font': HEI, 'line': 1.15}])
    colw = Inches(6.4)
    body = []
    for t in paras:
        body.append({'text': t, 'size': 15, 'color': INK, 'font': HEI, 'line': 1.55, 'space_after': 12})
    textbox(s, M, M + Inches(2.0), colw, Inches(4.8), body, anchor=MSO_ANCHOR.TOP)
    if img_key:
        ix = M + colw + Inches(0.5)
        iw = W - ix - M
        place_photo(s, IMG[img_key], ix, M + Inches(1.9), iw, Inches(4.3))
        if pull:
            caption(s, pull, ix, M + Inches(1.9) + Inches(4.3) + Inches(0.1), iw)
    page_num(s)

def s_poem():
    s = new_slide(); bg(s)
    kicker(s, 'FULL TEXT · 全词诵读')
    top = M + Inches(0.8)
    textbox(s, M, top, CW, Inches(0.6),
            [{'text': '沁园春 · 长沙', 'size': 24, 'color': FROST, 'bold': True, 'font': KAI,
              'align': PP_ALIGN.CENTER}])
    upper = ('独立寒秋，湘江北去，橘子洲头。看万山红遍，层林尽染；漫江碧透，百舸争流。'
             '鹰击长空，鱼翔浅底，万类霜天竞自由。怅寥廓，问苍茫大地，谁主沉浮？')
    lower = ('携来百侣曾游。忆往昔峥嵘岁月稠。恰同学少年，风华正茂；书生意气，挥斥方遒。'
             '指点江山，激扬文字，粪土当年万户侯。曾记否，到中流击水，浪遏飞舟？')
    textbox(s, M + Inches(1.0), top + Inches(0.85), CW - Inches(2.0), Inches(2.4),
            [{'text': upper, 'size': 22, 'color': INK, 'font': KAI, 'line': 1.7,
              'align': PP_ALIGN.CENTER, 'space_after': 0}])
    textbox(s, M + Inches(1.0), top + Inches(3.25), CW - Inches(2.0), Inches(2.4),
            [{'text': lower, 'size': 22, 'color': INK, 'font': KAI, 'line': 1.7,
              'align': PP_ALIGN.CENTER}])
    page_num(s)

def s_grid4():
    s = new_slide(); bg(s)
    kicker(s, '上阕 · 四幅画面')
    textbox(s, M, M + Inches(0.7), CW, Inches(0.6),
            [{'text': '远 · 近 · 动 · 静，构成秋日湘江图卷', 'size': 24, 'color': INK, 'bold': True, 'font': HEI}])
    cells = [
        (IMG['red_mountains'], '远观 · 万山红遍', '层林尽染，漫山红透'),
        (IMG['green_river'],   '近察 · 漫江碧透', '江水清澈，百舸争流'),
        (IMG['eagle'],         '仰视 · 鹰击长空', '搏击云端，充满力量'),
        (IMG['fish'],          '俯看 · 鱼翔浅底', '轻快游弋，自在生机'),
    ]
    gx = M; gy = M + Inches(1.6)
    cw = Inches(5.75); ch = Inches(2.35)
    gap = Inches(0.3)
    for i, (img, t, sub) in enumerate(cells):
        col = i % 2; row = i // 2
        x = gx + col * (cw + gap)
        y = gy + row * (ch + gap)
        place_photo(s, img, x, y, cw, ch)
        scrim(s, x, y + ch - Inches(0.85), cw, Inches(0.85), INK, 0.6)
        textbox(s, x + Inches(0.2), y + ch - Inches(0.78), cw - Inches(0.4), Inches(0.4),
                [{'text': t, 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI}])
        textbox(s, x + Inches(0.2), y + ch - Inches(0.42), cw - Inches(0.4), Inches(0.35),
                [{'text': sub, 'size': 11, 'color': SOFT, 'font': HEI}])
    page_num(s)

def s_compare(left_img, left_title, left_lines, right_title, right_lines, conclusion, left_color=FROST):
    """Comparison layout: large image left (full-height) | two stacked cards right (full-width) | conclusion bar"""
    s = new_slide(); bg(s)
    kicker(s, '炼字 · 对比品读')

    # Left image — tall crop, full height of content area
    img_w = Inches(4.8)
    img_h = Inches(5.5)
    place_photo(s, left_img, M, M + Inches(0.65), img_w, img_h)

    textbox(s, M, M + Inches(6.2), img_w, Inches(0.32),
            [{'text': left_title, 'size': 13, 'color': left_color,
              'bold': True, 'font': KAI, 'align': PP_ALIGN.CENTER}])

    # Right panel — two stacked cards filling the right half
    rx = M + img_w + Inches(0.35)
    rw = W - rx - M
    card_h = Inches(2.5)
    gap = Inches(0.18)
    card_top = M + Inches(0.65)

    # --- upper card: left ---
    lc_y = card_top
    lc_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, lc_y, rw, card_h)
    lc_card.fill.solid(); lc_card.fill.fore_color.rgb = WHITE
    lc_card.line.color.rgb = FROST; lc_card.line.width = Pt(2)

    lc_hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, lc_y, rw, Inches(0.55))
    lc_hdr.fill.solid(); lc_hdr.fill.fore_color.rgb = FROST
    lc_hdr.line.fill.background()
    textbox(s, rx + Inches(0.2), lc_y + Inches(0.1), rw - Inches(0.4), Inches(0.4),
            [{'text': left_title, 'size': 18, 'color': WHITE, 'bold': True, 'font': HEI}])

    ll = [{'text': t, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 6} for t in left_lines]
    textbox(s, rx + Inches(0.2), lc_y + Inches(0.65), rw - Inches(0.4), card_h - Inches(0.8), ll)

    # VS badge between the two cards
    vs_x = rx + rw / 2 - Inches(0.28)
    vs_y = lc_y + card_h + gap / 2 - Inches(0.28)
    vs_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, vs_x, vs_y, Inches(0.56), Inches(0.56))
    vs_circle.fill.solid(); vs_circle.fill.fore_color.rgb = GOLD
    vs_circle.line.fill.background()
    textbox(s, vs_x, vs_y + Inches(0.12), Inches(0.56), Inches(0.44),
            [{'text': 'VS', 'size': 15, 'color': WHITE, 'bold': True,
              'font': HEI, 'align': PP_ALIGN.CENTER}])

    # --- lower card: right ---
    rc_y = lc_y + card_h + gap
    rc_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, rc_y, rw, card_h)
    rc_card.fill.solid(); rc_card.fill.fore_color.rgb = WHITE
    rc_card.line.color.rgb = XIANG; rc_card.line.width = Pt(2)

    rc_hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, rc_y, rw, Inches(0.55))
    rc_hdr.fill.solid(); rc_hdr.fill.fore_color.rgb = XIANG
    rc_hdr.line.fill.background()
    textbox(s, rx + Inches(0.2), rc_y + Inches(0.1), rw - Inches(0.4), Inches(0.4),
            [{'text': right_title, 'size': 18, 'color': WHITE, 'bold': True, 'font': HEI}])

    rl = [{'text': t, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.5, 'space_after': 6} for t in right_lines]
    textbox(s, rx + Inches(0.2), rc_y + Inches(0.65), rw - Inches(0.4), card_h - Inches(0.8), rl)

    cby = H - Inches(0.7)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, cby, CW, Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    textbox(s, M + Inches(0.3), cby + Inches(0.1), CW - Inches(0.6), Inches(0.35),
            [{'text': conclusion, 'size': 12, 'color': WHITE, 'bold': True, 'font': HEI, 'line': 1.2}])

    page_num(s)


def s_table(title, headers, rows, kicker_txt='板书 · 结构'):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)
    textbox(s, M, M + Inches(0.65), CW, Inches(0.55),
            [{'text': title, 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])

    col_labels = ['层面', '上阕（景）', '下阕（情）']
    label_y = M + Inches(1.35)
    col_x = [M, M + Inches(1.7), M + Inches(7.0)]
    col_w = [Inches(1.5), Inches(4.9), W - (M + Inches(7.0)) - M]
    for i, (lx, lw, lbl) in enumerate(zip(col_x, col_w, col_labels)):
        textbox(s, lx, label_y, lw, Inches(0.38),
                [{'text': lbl, 'size': 12, 'color': MUTED, 'font': HEI,
                  'align': PP_ALIGN.CENTER}])

    row_colors = [FROST, XIANG, GOLD, INK]
    row_badges = ['起', '展', '合', '法']
    card_top = label_y + Inches(0.5)
    card_h = Inches(0.92)
    gap = Inches(0.14)

    for r_idx, (row_data, badge, bcolor) in enumerate(zip(rows, row_badges, row_colors)):
        y = card_top + r_idx * (card_h + gap)
        acc_w = Inches(0.1)
        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, y, acc_w, card_h)
        acc.fill.solid(); acc.fill.fore_color.rgb = bcolor
        acc.line.fill.background()
        card_body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        M + acc_w, y, CW - acc_w, card_h)
        card_body.fill.solid(); card_body.fill.fore_color.rgb = WHITE
        card_body.line.color.rgb = RGBColor(0xE8, 0xE2, 0xD9)

        cir_sz = Inches(0.52)
        cir_cx = M + acc_w + (col_w[0] - cir_sz) / 2
        cir = s.shapes.add_shape(MSO_SHAPE.OVAL, cir_cx, y + (card_h - cir_sz) / 2,
                                  cir_sz, cir_sz)
        cir.fill.solid(); cir.fill.fore_color.rgb = bcolor
        cir.line.fill.background()
        textbox(s, cir_cx, y + (card_h - cir_sz) / 2 + Inches(0.12),
                cir_sz, cir_sz,
                [{'text': badge, 'size': 18, 'color': WHITE, 'bold': True, 'font': KAI,
                  'align': PP_ALIGN.CENTER}])

        for c_idx in range(1, 3):
            cx = col_x[c_idx] + Inches(0.08)
            cw = col_w[c_idx] - Inches(0.16)
            textbox(s, cx, y + Inches(0.22), cw, card_h - Inches(0.36),
                    [{'text': row_data[c_idx], 'size': 14, 'color': INK,
                      'font': KAI if c_idx == 1 else HEI, 'line': 1.4,
                      'align': PP_ALIGN.CENTER}])

    page_num(s)

def s_closing():
    s = new_slide()
    place_photo(s, IMG['orange_isle'], 0, 0, W, H)
    scrim(s, 0, 0, W, H, INK, 0.55)
    cy = Inches(2.9)
    textbox(s, M, cy, CW, Inches(2.0),
            [{'text': '问苍茫大地，谁主沉浮？', 'size': 46, 'color': WHITE, 'bold': True,
              'font': KAI, 'align': PP_ALIGN.CENTER, 'line': 1.3}])
    rule(s, Inches(5.6), cy + Inches(2.0), Inches(2.1), GOLD, 2.5)
    textbox(s, M, cy + Inches(2.35), CW, Inches(0.5),
            [{'text': '— 沁园春 · 长沙 · 课堂版 —', 'size': 14, 'color': SOFT, 'font': HEI,
              'align': PP_ALIGN.CENTER}])
    page_num(s, dark=True)

# ===================================================================
# BUILD
# ===================================================================
s_cover()
s_contents()
s_divider('壹', '知人论世', '读懂作者，才读懂这首词', 'orange_isle')

# 作者页——上图下文（横幅美术作品 + 下方文字）
def s_author_art(img_key, kicker_txt, headline, paras, cap):
    s = new_slide(); bg(s)
    banner_h = Inches(2.5)
    place_photo(s, IMG[img_key], M, M, CW, banner_h)
    caption(s, cap, M, M + banner_h + Inches(0.08), CW)
    kicker(s, kicker_txt, M, M + banner_h + Inches(0.5))
    textbox(s, M, M + banner_h + Inches(1.15), CW, Inches(0.8),
            [{'text': headline, 'size': 26, 'color': INK, 'bold': True, 'font': HEI, 'line': 1.15}])
    body = []
    for t in paras:
        body.append({'text': t, 'size': 15, 'color': INK, 'font': HEI, 'line': 1.5, 'space_after': 10})
    textbox(s, M, M + banner_h + Inches(2.1), Inches(10.8), Inches(1.9), body)
    page_num(s)

s_author_art('mao_art', '作者 · 毛泽东', '1925：三十二岁的革命者',
          ['1925 年，毛泽东三十二岁，正值大革命前夜。',
           '他回湖南领导农民运动，又奉命赴广州，途经长沙，重游橘子洲。',
           '时值工农运动风起云涌，而前路艰难，他心潮起伏。',
           '独立寒秋，抚今追昔，豪情与忧思一并涌上笔端。'],
          '《潇湘八景图》（南宋·牧溪）— 湘江秋意图')

# ---- 背景页：左图右文时间线 ----
def s_timeline(kicker_txt, headline, events, img_key, pull):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)
    textbox(s, M, M + Inches(0.65), CW, Inches(0.55),
            [{'text': headline, 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])

    pw = Inches(4.8)
    place_photo(s, IMG[img_key], M, M + Inches(1.4), pw, Inches(3.8))
    caption(s, pull, M, M + Inches(5.25), pw)

    rx = M + pw + Inches(0.45)
    rw = W - rx - M
    card_h = Inches(0.85)
    gap = Inches(0.12)
    t_colors = [FROST, XIANG, GOLD, INK]
    for idx, ev in enumerate(events):
        y = M + Inches(1.4) + idx * (card_h + gap)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  rx, y + card_h / 2 - Inches(0.15),
                                  Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = t_colors[idx % len(t_colors)]
        dot.line.fill.background()
        if idx < len(events) - 1:
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     rx + Inches(0.13), y + card_h / 2 + Inches(0.15),
                                     Inches(0.04), card_h + gap - Inches(0.3))
            ln.fill.solid(); ln.fill.fore_color.rgb = SOFT
            ln.line.fill.background()
        crd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  rx + Inches(0.45), y, rw - Inches(0.45), card_h)
        crd.fill.solid(); crd.fill.fore_color.rgb = WHITE
        crd.line.color.rgb = RGBColor(0xE8, 0xE2, 0xD9)
        textbox(s, rx + Inches(0.58), y + Inches(0.18), rw - Inches(0.68), card_h - Inches(0.28),
                [{'text': ev, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.42}])
    page_num(s)

s_timeline('背景 · 1925', '一首词里的时代',
           ['1925 年，国共第一次合作，北伐前夕，革命形势高涨。',
            '毛泽东回湖南考察农民运动，亲见民众觉醒的力量。',
            '同年秋，他重游长沙橘子洲，面对秋景忆战斗岁月。',
            '彼时人民尚不自由——这正是"怅寥廓"的历史来由。'],
           img_key='autumn_china', pull='秋日的湘江畔，正是写词之地')

# ---- 自由立骨主题页（融合 cv43571662 招牌：以"自由"串全词） ----
def s_keyword(kicker_txt, char_display, source_quote, analysis_lines):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)

    char_w = Inches(2.8)
    char_h = Inches(4.0)
    char_box_y = M + Inches(1.0)
    char_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   M, char_box_y, char_w, char_h)
    char_box.fill.solid(); char_box.fill.fore_color.rgb = FROST
    char_box.line.fill.background()
    # 自适应字号：单字大、双字小
    _csz = 96 if len(char_display) <= 1 else 74
    # 垂直居中：CJK 字符视觉中心在文字区域偏上，用红框高度 42% 处作为文本框起点（实测校准值）
    _ty = char_box_y + char_h * 0.44
    textbox(s, M, _ty, char_w, Inches(1.8),
            [{'text': char_display, 'size': _csz, 'color': WHITE, 'bold': True,
              'font': KAI, 'align': PP_ALIGN.CENTER}])

    rx = M + char_w + Inches(0.5)
    rw = W - rx - M

    qb_h = Inches(0.85)
    qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, M + Inches(1.0), rw, qb_h)
    qbox.fill.solid(); qbox.fill.fore_color.rgb = SOFT
    qbox.line.fill.background()
    qa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, M + Inches(1.0), Inches(0.08), qb_h)
    qa.fill.solid(); qa.fill.fore_color.rgb = GOLD
    qa.line.fill.background()
    textbox(s, rx + Inches(0.25), M + Inches(1.15), rw - Inches(0.35), qb_h - Inches(0.2),
            [{'text': source_quote.replace('"', ''), 'size': 17, 'color': INK,
              'font': KAI, 'bold': True, 'line': 1.35}])

    card_top = M + Inches(2.05)
    card_h = Inches(0.88)
    gap = Inches(0.12)
    for idx, line in enumerate(analysis_lines):
        y = card_top + idx * (card_h + gap)
        nb = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 rx + Inches(0.05), y + Inches(0.18),
                                 Inches(0.42), Inches(0.42))
        nb.fill.solid()
        nb.fill.fore_color.rgb = XIANG if idx < len(analysis_lines) - 1 else INK
        nb.line.fill.background()
        textbox(s, rx + Inches(0.05), y + Inches(0.24), Inches(0.42), Inches(0.42),
                [{'text': str(idx + 1), 'size': 14, 'color': WHITE, 'bold': True,
                  'font': HEI, 'align': PP_ALIGN.CENTER}])
        textbox(s, rx + Inches(0.58), y + Inches(0.2), rw - Inches(0.68), card_h - Inches(0.28),
                [{'text': line, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45}])

    page_num(s)

s_keyword('主题 · 自由立骨', '自由', '"万类霜天竞自由"',
          ['以"自由"一词串起全词——这是全课的骨。',
           '上阕：万物在霜天里"竞自由"，秋景因自由而昂扬。',
           '"怅寥廓"因何而怅？因时代里人民尚不自由——由景入思。',
           '"谁主沉浮"的回答：正是我辈青年，去争那自由。'])

s_divider('贰', '初读感知', '先读准字音，再读懂气象', 'autumn_china')
s_poem()

# ---- 朗读指导页：竖向宽卡片 ----
def s_reading_guide(kicker_txt, headline, guides):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)
    textbox(s, M, M + Inches(0.6), CW, Inches(0.5),
            [{'text': headline, 'size': 26, 'color': INK, 'bold': True, 'font': HEI}])

    labels = ['字音', '节拍', '情感', '建议']
    colors = [FROST, XIANG, GOLD, INK]
    icons = ['01', '02', '03', '04']

    card_y = M + Inches(1.25)
    card_h = Inches(1.05)
    gap = Inches(0.16)
    for idx, (lbl, g) in enumerate(zip(labels, guides)):
        bar_w = Inches(0.1)
        bar_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, card_y, bar_w, card_h)
        bar_shape.fill.solid(); bar_shape.fill.fore_color.rgb = colors[idx]
        bar_shape.line.fill.background()
        card_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         M + bar_w, card_y, CW - bar_w, card_h)
        card_shape.fill.solid(); card_shape.fill.fore_color.rgb = WHITE
        card_shape.line.color.rgb = RGBColor(0xE8, 0xE2, 0xD9)
        cir_r = Inches(0.38)
        cir = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  M + bar_w + Inches(0.22), card_y + Inches(0.18),
                                  cir_r, cir_r)
        cir.fill.solid(); cir.fill.fore_color.rgb = colors[idx]
        cir.line.fill.background()
        textbox(s, M + bar_w + Inches(0.22), card_y + Inches(0.25),
                cir_r, cir_r,
                [{'text': icons[idx], 'size': 15, 'color': WHITE, 'bold': True,
                  'font': HEI, 'align': PP_ALIGN.CENTER}])
        label_x = M + bar_w + cir_r + Inches(0.4)
        textbox(s, label_x, card_y + Inches(0.18),
                CW - bar_w - cir_r - Inches(0.6), Inches(0.38),
                [{'text': lbl, 'size': 17, 'color': colors[idx], 'bold': True, 'font': HEI}])
        clean_text = g.replace(f'【{lbl}】', '').strip()
        textbox(s, label_x, card_y + Inches(0.56),
                CW - bar_w - cir_r - Inches(0.6), card_h - Inches(0.64),
                [{'text': clean_text, 'size': 14, 'color': INK,
                  'font': KAI, 'line': 1.45}])
        card_y += card_h + gap

    page_num(s)

s_reading_guide('朗读 · 指导', '上阕写景 · 下阕抒情',
                ['【字音】舸 gě　寥 liáo　廓 kuò　怅 chàng　稠 chóu　遒 qiú　遏 è',
                 '【节拍】上阕"看"字领起七句，一气读尽；下阕"忆"字领起，渐入昂扬。',
                 '【情感】整体感知：上阕写景、下阕抒情；一幅秋景，生机勃勃。',
                 '【建议】先齐读把握节奏，再分角色读出"问"与"答"的气势。'])
s_divider('叁', '上阕 · 写景', '七句景物，写尽湘江秋色', 'red_mountains')
s_feature('red_mountains', '意象 · 万山红遍', '层林尽染，漫山红透',
          ['"看万山红遍，层林尽染"——放眼望去，群山红透。',
           '红色，是秋日枫林的色彩，也是革命热情的象征。',
           '一个"遍"字，写尽范围之广、气势之盛。',
           '诗人以暖色开篇，奠定全词昂扬的基调。'],
          '岳麓山秋色（长沙）')
s_feature('green_river', '意象 · 漫江碧透', '江水澄碧，百舟竞发',
          ['"漫江碧透，百舸争流"——满江碧水，千帆竞渡。',
           '"碧透"写水之清，"争流"写船之疾、人之奋。',
           '一动一静，江面顿生活力。',
           '远景的红山与近景的碧水，构成冷暖对照。'],
          '湘江碧水（长沙）')
s_feature('boats', '意象 · 百舸争流', '千帆竞发，奋斗不息',
          ['"舸"即大船；"百舸争流"写船多而疾。',
           '众船争流，暗喻时代中人们奋楫向前。',
           '与上句"碧透"相连，江面图景完整呈现。',
           '少年毛泽东从中看到的，是民族的生机。'],
          '岳麓书院·湘江畔（长沙）')
s_feature('eagle', '意象 · 鹰击长空', '搏击云端，充满力量',
          ['"鹰击长空，鱼翔浅底"——一上一下，天地开阔。',
           '"击"字极妙：不是"飞"，是奋力搏击，见出力量。',
           '雄鹰破空，正是诗人胸襟的写照。',
           '与下句"鱼翔"相对，构成立体的生机画卷。'],
          '翱翔的苍鹰（意象参照）')
s_feature('fish', '意象 · 鱼翔浅底', '轻快游弋，自在生机',
          ['"鱼翔浅底"——鱼在清浅水底轻快游动。',
           '"翔"字将鱼写得如鸟飞，化实为灵、化静为动。',
           '课内外互证：以"鸢飞戾天，鱼跃于渊"证"击/翔"之力；',
           '以"潭中鱼可百许头，皆若空游无所依"证水清而自由。'],
          '清澈水底的游鱼（意象参照）')
s_grid4()
# 看字领七句 —— 升级：视角变换 + 具身表演 + 大脑荧屏想象（融合 课堂实录4）
s_editorial('上阕 · 章法', '一个"看"字，领起七句',
          ['"看"字总领以下七句，写尽湘江秋色。',
           '视角依次变换：远眺万山，近察漫江，仰视雄鹰，俯看游鱼。',
           '课堂小招：让学生"演"诗人——手插腰、慢转换，体会伟人气度；',
           '再闭眼把文字"放映"在大脑荧屏：先勾勒形态，再染上色彩，',
           '后添入声响，终融入情感——动静相宜，红碧交错，生机盎然。'],
          img_key='red_mountains', pull='岳麓山秋色（长沙）')
s_divider('肆', '炼字品析', '一字之差，境界全出', 'eagle')
s_compare(IMG['eagle'], '击', ['鹰"击"长空', '奋力搏击，见力度与气势', '写尽雄健昂扬之态', '一个"击"字，鹰便有了性格'],
          '飞', ['普通"飞"行', '仅表动作，平淡无奇', '少了那份搏击苍穹的劲', '一字之差，境界天壤'],
          '结论："击"以力度代动作，让雄鹰有了性格——这正是炼字的力量。')
s_compare(IMG['fish'], '翔', ['鱼"翔"浅底', '如鸟轻飞，化实为灵', '写出自由自在的生机', '拟物于鸟，画面灵动'],
          '游', ['普通"游"动', '仅表位移，呆板', '不见那份轻盈灵动', '化实为常，诗意全无'],
          '结论："翔"拟物于鸟，鱼便有了飞的姿态——拟人化让画面活起来。')
# 竞字 —— 升级：字形 + 《齐物论》郭象注（融合 cv43571662 招牌）
s_keyword('炼字 · 竞', '竞', '"万类霜天竞自由"',
          ['"竞"者，竞相、争着也；其形从"競"，二人对辩相逐。',
           '郭象注《庄子·齐物论》"有竞有争"：并逐曰竞，对辩曰争。',
           '"竞"是万物并相追逐——开阔而不拥挤，自由而不争抢。',
           '由"万类霜天竞自由"自然过渡到"问苍茫大地，谁主沉浮"。'])
# 怅寥廓 —— 升级：落到"人民不自由"（融合 cv43571662）
s_keyword('炼字 · 怅寥廓', '怅', '"怅寥廓，问苍茫大地，谁主沉浮"',
          ['"怅"非消极惆怅，而是面对辽阔天地引发的深沉思索。',
           '因何而怅？1925 年的中国，人民尚不自由——这是历史的怅惘。',
           '"问苍茫大地，谁主沉浮"：惊天一问，是英雄气概，更是担当。',
           '由"人民不自由"反推出"我们主沉浮"——上阕由景入思，埋下伏笔。'])
# 独字 —— 新增：倒装凸显 + 《黄河颂》互文（融合 jiaoxueshilu/858）
s_keyword('炼字 · 独', '独', '"独立寒秋"',
          ['"独立寒秋，湘江北去，橘子洲头"本是倒装。',
           '常序为"寒秋，独立橘子洲头，湘江北去"；词序一倒，"独立"凸显。',
           '凸显的，是先立于天地间的诗人形象——人格的独立与孤高。',
           '可与《黄河颂》"我站在高山之巅"互文：大意象常把"我"置于开篇。'])
s_divider('伍', '下阕 · 抒情', '由景入情，由己及世', 'rapids')
# ---- 同学少年页：左书法图 + 右分析卡片 ----
def s_youth_portrait(img_key, kicker_txt, headline, paras, cap):
    s = new_slide(); bg(s)
    kicker(s, kicker_txt)

    img_w = Inches(4.5)
    img_h = Inches(5.2)
    place_photo(s, IMG[img_key], M, M + Inches(0.65), img_w, img_h)

    textbox(s, M, M + Inches(5.9), img_w, Inches(0.32),
            [{'text': cap, 'size': 10, 'color': MUTED, 'font': KAI, 'align': PP_ALIGN.CENTER}])

    rx = M + img_w + Inches(0.45)
    rw = W - rx - M

    hl_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 rx, M + Inches(0.65), rw, Inches(0.85))
    hl_box.fill.solid(); hl_box.fill.fore_color.rgb = XIANG
    hl_box.line.fill.background()
    textbox(s, rx + Inches(0.25), M + Inches(0.78), rw - Inches(0.5), Inches(0.65),
            [{'text': headline, 'size': 21, 'color': WHITE, 'bold': True,
              'font': HEI, 'align': PP_ALIGN.CENTER}])

    card_top = M + Inches(1.68)
    card_h = Inches(1.02)
    gap = Inches(0.10)
    y_colors = [FROST, XIANG, GOLD, INK]
    for idx, line in enumerate(paras):
        y = card_top + idx * (card_h + gap)
        color = y_colors[idx % len(y_colors)]

        crd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 rx, y, rw, card_h)
        crd.fill.solid(); crd.fill.fore_color.rgb = SOFT
        crd.line.color.rgb = color
        crd.line.width = Pt(1.5)

        nb = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 rx + Inches(0.14), y + Inches(0.24),
                                 Inches(0.48), Inches(0.48))
        nb.fill.solid(); nb.fill.fore_color.rgb = color
        nb.line.fill.background()
        textbox(s, rx + Inches(0.14), y + Inches(0.30),
                 Inches(0.48), Inches(0.48),
                [{'text': str(idx + 1), 'size': 16, 'color': WHITE, 'bold': True,
                  'font': HEI, 'align': PP_ALIGN.CENTER}])

        textbox(s, rx + Inches(0.72), y + Inches(0.18),
                rw - Inches(0.88), card_h - Inches(0.30),
                [{'text': line, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45}])

    page_num(s)

s_youth_portrait('mao_mural', '形象 · 同学少年', '风华正茂，挥斥方遒',
          ['"恰同学少年，风华正茂"——写早年战友与自己。',
           '"书生意气，挥斥方遒"——热情奔放，才华横溢。',
           '"指点江山，激扬文字"——评说国事；"激扬"即激浊扬清。',
           '"粪土当年万户侯"——"粪土"意动，视权贵如粪土，何等气概。'],
          '书法作品（毛泽东诗词《沁园春·长沙》）')
s_feature('rapids', '典故 · 中流击水', '浪遏飞舟，壮志不改',
          ['"曾记否，到中流击水，浪遏飞舟"——江心奋游，浪头挡住船行。',
           '化用"中流击楫"：祖逖渡江击楫，誓复中原。',
           '毛泽东反用其意：不是忧伤，而是豪情壮志。',
           '这一问一答，巧妙回应上阕"谁主沉浮"——主沉浮者，正是我辈青年。'],
          '湘江中流（意象参照）')
# ---- 当代打通页（融合 cv43571662 阅兵 + jiaoxueshilu/861 奋斗者） ----
s_feature('orange_isle', '当代 · 青春回响', '从湘江到今天',
          ['九十八年后再看"万类霜天竞自由"：今日中国，自由舒展。',
           '课堂延伸：看 9.3 阅兵，感"竞自由"之盛；以"最美奋斗者"',
           '（钟南山、叶连平）对照"同学少年"——奋斗，是青春不变的底色。',
           '一首词，连起两代青年的壮志：彼争自由，吾承其志。'],
          '今日长沙湘江（橘子洲航拍）')
s_divider('陆', '主旨升华', '同样的秋，不同的胸怀', 'orange_isle')
s_compare(IMG['autumn_china'], '古人悲秋', ['"万里悲秋常作客"（杜甫）', '"怎一个愁字了得"（李清照）', '秋=凋零=伤感=身世之悲', '萧瑟凄凉，是生命衰颓的哀叹'],
          '毛泽东颂秋', ['万山红遍，生机盎然', '鹰击鱼翔，力量勃发', '秋=昂扬=壮志=家国之怀', '万类霜天竞自由'],
          '结论：同写秋景，古人见衰飒，毛泽东见蓬勃——胸襟不同，秋天便不同。这正是情景交融的写法。',
          left_color=MUTED)
# ---- 意象理论页（融合 cv43571662《人间词话》意象说） ----
s_editorial('理论 · 意象', '一切景语皆情语',
          ['意象 = 主观意态 + 客观物象（《人间词话》：以我观物，故物皆着我之色彩）。',
           '词中红山、碧水、雄鹰、游鱼，本是无情物；',
           '一经诗人"我"的注视，便都染上昂扬自由的情志。',
           '这正解答了"悲秋 vs 颂秋"：景同，情异；情景交融，方成此词。'],
          img_key='autumn_china', pull='岳麓书院秋色（长沙）')
s_table('板书结构 · 立→看→怅→问→忆→记（以"自由"立骨）',
        ['层面', '上阕（景）', '下阕（情）'],
        [['起', '独立寒秋·橘子洲头', '携来百侣·曾游'],
         ['展', '万山·漫江·百舸·鹰鱼', '同学少年·风华正茂'],
         ['合', '怅寥廓·谁主沉浮', '到中流击水·浪遏飞舟'],
         ['法', '借景抒情·以自由立骨', '叙议结合·用典']])
# ---- 作业页：杂志风三栏任务卡 ----
def s_exercise():
    s = new_slide(); bg(s)
    kicker(s, '作业 · 分层')
    textbox(s, M, M + Inches(0.65), CW, Inches(0.55),
            [{'text': '课后任务', 'size': 28, 'color': INK, 'bold': True, 'font': HEI}])

    cards = [
        ('01', '基础', FROST, '背诵全词，默写上阕；解释"击、翔、竞、独"的妙处。'),
        ('02', '提高', XIANG, '以"自由"为眼，写 200 字短评：毛泽东为何"颂秋"。'),
        ('03', '拓展', INK,  '联系"最美奋斗者"，谈今日青年如何"主沉浮"。'),
    ]
    card_y = M + Inches(1.4)
    card_h = Inches(1.15)
    gap = Inches(0.2)
    for num, label, color, content in cards:
        bar_w = Inches(0.12)
        bar_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, card_y, bar_w, card_h)
        bar_shape.fill.solid(); bar_shape.fill.fore_color.rgb = color
        bar_shape.line.fill.background()
        card_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         M + bar_w, card_y, CW - bar_w, card_h)
        card_shape.fill.solid(); card_shape.fill.fore_color.rgb = WHITE
        card_shape.line.color.rgb = RGBColor(0xE8, 0xE2, 0xD9)
        cir_r = Inches(0.36)
        cir = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  M + bar_w + Inches(0.28), card_y + Inches(0.18),
                                  cir_r, cir_r)
        cir.fill.solid(); cir.fill.fore_color.rgb = color
        cir.line.fill.background()
        textbox(s, M + bar_w + Inches(0.28), card_y + Inches(0.24), cir_r, cir_r,
                [{'text': num, 'size': 16, 'color': WHITE, 'bold': True, 'font': HEI,
                  'align': PP_ALIGN.CENTER}])
        textbox(s, M + bar_w + cir_r + Inches(0.45), card_y + Inches(0.18),
                CW - bar_w - cir_r - Inches(0.65), Inches(0.42),
                [{'text': label, 'size': 17, 'color': color, 'bold': True, 'font': HEI,
                  'align': PP_ALIGN.CENTER}])
        textbox(s, M + bar_w + cir_r + Inches(0.45), card_y + Inches(0.58),
                CW - bar_w - cir_r - Inches(0.65), Inches(0.58),
                [{'text': content, 'size': 14, 'color': INK, 'font': KAI, 'line': 1.45,
                  'align': PP_ALIGN.CENTER}])
        card_y += card_h + gap

    tip_y = H - Inches(1.2)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, tip_y, CW, Inches(0.55))
    tip.fill.solid(); tip.fill.fore_color.rgb = SOFT
    tip.line.fill.background()
    textbox(s, M + Inches(0.3), tip_y + Inches(0.13), CW - Inches(0.6), Inches(0.38),
            [{'text': '( 参考答案随教师版详案发放 )',
              'size': 13, 'color': MUTED, 'font': KAI, 'align': PP_ALIGN.CENTER}])
    page_num(s)

s_exercise()
s_closing()

# ===================================================================
# SAVE
# ===================================================================
os.makedirs(os.path.dirname(OUT_PPTX), exist_ok=True)
prs.save(OUT_PPTX)
print(f'SAVED {OUT_PPTX} slides={len(prs.slides)}')
