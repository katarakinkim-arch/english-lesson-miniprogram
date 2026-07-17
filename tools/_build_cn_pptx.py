# -*- coding: utf-8 -*-
"""用 python-pptx 把 17 张语文 v2 设计稿 PNG 生成 l-cn-bs-u1-1.pptx（整页背景）。
与英语 pilot 同构：16 设计页 + 1 素材附录 = 17 页。"""
import os
from pptx import Presentation
from pptx.util import Inches

DESIGN = 'cn_preview'
OUT_PPTX = 'l-cn-bs-u1-1.pptx'

slides = [
    'slide_01_cover.png', 'slide_02_contents.png', 'slide_03_overview.png',
    'slide_04_objectives.png', 'slide_05_keypoints.png',
    'slide_06_step1.png', 'slide_07_step2.png', 'slide_08_step3.png',
    'slide_09_step4.png', 'slide_10_step5.png', 'slide_11_step6.png',
    'slide_12_blackboard.png', 'slide_13_homework_basic.png',
    'slide_14_homework_adv.png', 'slide_15_reflection.png',
    'slide_16_summary.png', 'slide_17_appendix.png',
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for i, name in enumerate(slides, 1):
    slide = prs.slides.add_slide(blank)
    img_path = os.path.join(DESIGN, name)
    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    print(f'  {i:02d}  {name}')

prs.save(OUT_PPTX)
print(f'\nPPTX -> {OUT_PPTX} ({os.path.getsize(OUT_PPTX)//1024} KB, {len(prs.slides)} 页)')
